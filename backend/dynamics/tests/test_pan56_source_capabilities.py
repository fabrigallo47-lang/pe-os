"""PAN-56 — honest, addressable multi-format parsing under SourceEnvelope."""
from __future__ import annotations

import json
import builtins
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch


ROOT = Path(__file__).resolve().parents[3]
sys.path.insert(0, str(ROOT))

import app.v20_router as router  # noqa: E402
from tools.excel_formula_graph import compile_workbook  # noqa: E402
from tools.extract_v2 import (  # noqa: E402
    UnsupportedSourceError,
    _capture_workbook_graphs,
    parse_source,
)
from tools.source_capabilities import capability_manifest  # noqa: E402
from tools.source_envelope import build_source_envelope, extractor_source_record  # noqa: E402


EXPECTATIONS = ROOT / "tools" / "fixtures" / "pan56_capability_expectations.json"
FORMULA_FIXTURE = ROOT / "tools" / "fixtures" / "pan51_formula_model.xlsx"


def _record(path: Path, document_type: str = "Diligence source") -> dict:
    return extractor_source_record(build_source_envelope(
        path,
        "pan56-case",
        "2026-09-01T10:00:00Z",
        declared_metadata={
            "document_type": document_type,
            "issuer": "Example management",
            "effective_date": "2026-06-30",
        },
    ))


def _write_docx(path: Path) -> None:
    # PAN-103: parse_docx reads via docx2python (native table-grid access),
    # which requires a genuinely valid OOXML package -- a bare
    # word/document.xml with no [Content_Types].xml or relationships (the
    # old raw-XML reader's minimum) no longer opens. Build a real, minimal
    # document through python-docx itself instead of hand-rolling package
    # structure (same fix already applied to _write_pptx for PAN-101).
    from docx import Document

    document = Document()
    document.add_paragraph("Revenue was EUR 20m in FY2025A.")
    document.add_paragraph("Management expects growth in FY2026E.")
    document.save(str(path))


def _write_pptx(path: Path) -> None:
    # PAN-102/PAN-101: parse_pptx reads via python-pptx (native chart/table/
    # notes access), which requires a genuinely valid OOXML package -- a
    # bare ppt/slides/slide1.xml with no [Content_Types].xml, presentation
    # part, or relationships (the old raw-XML reader's minimum) no longer
    # opens. Build a real, minimal deck through python-pptx itself instead
    # of hand-rolling package structure.
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(0, 0, 9144000, 685800)
    text_box.text_frame.text = "FY2026E base case revenue"
    presentation.save(str(path))


class _FakePDF:
    def __init__(self, texts: list[str]):
        self.pages = [SimpleNamespace(extract_text=lambda text=text: text) for text in texts]

    def __enter__(self):
        return self

    def __exit__(self, *_args):
        return False


class PAN56CapabilityContractTests(unittest.TestCase):
    def test_versioned_matrix_covers_every_required_source_family(self):
        expected = json.loads(EXPECTATIONS.read_text(encoding="utf-8"))
        manifest = capability_manifest()
        self.assertEqual(manifest["schema"], expected["schema"])
        self.assertTrue(manifest["invariants"]["no_fake_extraction"])
        by_id = {item["capability_id"]: item for item in manifest["capabilities"]}
        self.assertTrue(set(expected["required_capabilities"]).issubset(by_id))
        for capability_id in expected["required_capabilities"]:
            with self.subTest(capability_id=capability_id):
                item = by_id[capability_id]
                self.assertTrue(item["provenance_fields"])
                self.assertTrue(item["locator_semantics"])
                self.assertIn("known_at", item["period_semantics"])
        self.assertEqual(by_id["scanned_pdf_ocr"]["support"], "UNSUPPORTED")
        self.assertTrue(by_id["scanned_pdf_ocr"]["action_if_unavailable"])

    def test_csv_routing_carries_envelope_provenance_locator_and_period_context(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "operating.csv"
            path.write_text("period,revenue\nFY2025A,20\nFY2026E,25\n", encoding="utf-8")
            record = _record(path)
            envelope = record["source_envelope"]
            chunks = parse_source(path, max_words=20, source_record=record)
        self.assertEqual(envelope["capability_contract"], "panta.source-capabilities/1.0")
        self.assertEqual(envelope["parser_capability"], "csv")
        self.assertTrue(all(chunk.source_type == "csv" for chunk in chunks))
        self.assertRegex(chunks[0].locator, r"operating\.csv::rows:1-3")
        self.assertEqual(chunks[0].provenance["source_version_id"], envelope["source_version_id"])
        self.assertEqual(chunks[0].provenance["locator"], chunks[0].locator)
        self.assertEqual(chunks[0].period_context["effective_date"], "2026-06-30")
        self.assertEqual(chunks[0].period_context["known_at"], "2026-09-01T10:00:00Z")
        self.assertEqual(chunks[0].period_context["semantics"], "DECLARED_ONLY")

    def test_docx_and_pptx_use_native_openxml_locators(self):
        with tempfile.TemporaryDirectory() as tmp:
            docx = Path(tmp) / "report.docx"
            pptx = Path(tmp) / "deck.pptx"
            _write_docx(docx)
            _write_pptx(pptx)
            document_chunks = parse_source(docx, source_record=_record(docx))
            slide_chunks = parse_source(pptx, source_record=_record(pptx, "Management deck"))
        self.assertEqual(document_chunks[0].source_type, "docx")
        # PAN-103: locator label is "blocks" now, not "paragraphs" --
        # docx2python-derived chunks may carry paragraphs or real tables,
        # and the old label stopped being accurate.
        self.assertRegex(document_chunks[0].locator, r"report\.docx::blocks:\d+-\d+")
        self.assertIn("FY2025A", document_chunks[0].body)
        self.assertEqual(slide_chunks[0].source_type, "pptx")
        self.assertIn("deck.pptx::slide:1", slide_chunks[0].locator)
        self.assertEqual(slide_chunks[0].page_or_slide_number, 1)

    def test_transcript_and_email_exports_keep_native_message_or_time_locators(self):
        with tempfile.TemporaryDirectory() as tmp:
            transcript = Path(tmp) / "call.vtt"
            transcript.write_text(
                "WEBVTT\n\n00:00:01.000 --> 00:00:04.000\nRevenue grew in FY2025A.\n",
                encoding="utf-8",
            )
            email = Path(tmp) / "update.eml"
            email.write_text(
                "From: ceo@example.com\nTo: deal@example.com\n"
                "Date: Tue, 1 Sep 2026 09:30:00 +0000\n"
                "Subject: Trading update\nMIME-Version: 1.0\n"
                "Content-Type: multipart/mixed; boundary=pan56\n\n"
                "--pan56\nContent-Type: text/plain; charset=utf-8\n\n"
                "August revenue was EUR 2m.\n"
                "--pan56\nContent-Type: application/octet-stream\n"
                "Content-Disposition: attachment; filename=model.xlsx\n\n"
                "not parsed as fake spreadsheet evidence\n--pan56--\n",
                encoding="utf-8",
            )
            transcript_chunks = parse_source(transcript, source_record=_record(transcript, "Call transcript"))
            email_chunks = parse_source(email, source_record=_record(email, "Email export"))
        self.assertIn("call.vtt::cue:1:00:00:01.000-->00:00:04.000", transcript_chunks[0].locator)
        self.assertEqual(transcript_chunks[0].provenance["parser_capability"], "transcript_export")
        self.assertIn("update.eml::message:1:body", email_chunks[0].locator)
        self.assertEqual(email_chunks[0].provenance["excluded_attachments"], 1)
        self.assertEqual(
            email_chunks[0].provenance["attachment_policy"],
            "SEPARATE_SOURCE_ENVELOPE_REQUIRED",
        )
        self.assertEqual(email_chunks[0].period_context["document_date_source"], "email-header:Date")
        self.assertTrue(email_chunks[0].period_context["document_date"].startswith("2026-09-01"))

    def test_native_pdf_is_page_addressable_and_scans_fail_actionably(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "memo.pdf"
            path.write_bytes(b"fixture bytes are handled by the fake local reader")
            record = _record(path)
            native_reader = SimpleNamespace(open=lambda _path: _FakePDF(["Revenue EUR 20m FY2025A."]))
            with patch.dict(sys.modules, {"pdfplumber": native_reader}):
                native = parse_source(path, source_record=record)
            scanned_reader = SimpleNamespace(open=lambda _path: _FakePDF(["", "   "]))
            with patch.dict(sys.modules, {"pdfplumber": scanned_reader}):
                with self.assertRaises(UnsupportedSourceError) as caught:
                    parse_source(path, source_record=record)
        self.assertEqual(native[0].locator, "p1")
        self.assertEqual(native[0].page_or_slide_number, 1)
        response = caught.exception.to_dict()
        self.assertEqual(response["schema"], "panta.source-capabilities/1.0")
        self.assertEqual(response["code"], "OCR_REQUIRED")
        self.assertEqual(response["capability_id"], "scanned_pdf_ocr")
        self.assertIn("OCR", response["action"])

    def test_unsupported_and_malformed_sources_never_produce_fake_chunks(self):
        with tempfile.TemporaryDirectory() as tmp:
            for filename, expected_code in (
                ("legacy.xls", "LEGACY_EXCEL_UNSUPPORTED"),
                ("mail.msg", "OUTLOOK_MSG_UNSUPPORTED"),
                ("archive.bin", "UNSUPPORTED_SOURCE"),
                ("broken.docx", "OPENXML_INVALID"),
            ):
                path = Path(tmp) / filename
                path.write_bytes(b"not the claimed format")
                with self.subTest(filename=filename):
                    with self.assertRaises(UnsupportedSourceError) as caught:
                        parse_source(path)
                    self.assertEqual(caught.exception.to_dict()["code"], expected_code)
                    self.assertTrue(caught.exception.to_dict()["action"])


class PAN56FormulaContinuityTests(unittest.TestCase):
    def test_v2_sidecar_and_v20_admission_preserve_formula_graph_without_formulas_package(self):
        expected = json.loads(EXPECTATIONS.read_text(encoding="utf-8"))["formula_fixture"]
        record = _record(FORMULA_FIXTURE, "LBO Model")
        real_import = builtins.__import__

        def import_without_formulas(name, globals=None, locals=None, fromlist=(), level=0):
            if name == "formulas" and level == 0:
                raise ImportError("PAN-56 verifies the optional evaluator is absent")
            return real_import(name, globals, locals, fromlist, level)

        with patch("builtins.__import__", side_effect=import_without_formulas):
            sidecar = _capture_workbook_graphs([FORMULA_FIXTURE], {FORMULA_FIXTURE: record})
        workbook = sidecar["workbooks"][0]
        formula_cells = [
            cell for cell in workbook["graph"]["cells"].values()
            if cell["kind"] == "formula"
        ]
        self.assertEqual(len(formula_cells), expected["formula_count"])
        self.assertEqual(workbook["source_id"], record["source_id"])
        self.assertEqual(
            workbook["graph"]["evaluation_warnings"][0]["code"],
            "OPTIONAL_FORMULA_EVALUATOR_UNAVAILABLE",
        )
        self.assertTrue(all(str(cell["value"]).startswith("=") for cell in formula_cells))
        self.assertTrue(any(cell["precedents"] for cell in formula_cells))

        compiled = compile_workbook(FORMULA_FIXTURE)
        previous_pipeline = router.PIPELINE_OUT
        with tempfile.TemporaryDirectory() as tmp:
            try:
                router.PIPELINE_OUT = Path(tmp)
                admitted = router._admit_excel_formula_graph(
                    "pan56-case",
                    {"proposal_id": "pan56-proposal", "excel_formula_graph": compiled},
                    "reviewer-pan56",
                )
                stored = router._load_excel_model_graphs("pan56-case")
            finally:
                router.PIPELINE_OUT = previous_pipeline
        self.assertIsNotNone(admitted)
        self.assertEqual(len(stored), 1)
        self.assertEqual(len(stored[0]["formulas"]), expected["formula_count"])
        self.assertTrue(any(
            edge["rel"] == expected["required_relation"]
            for edge in stored[0]["edges"]
        ))
        self.assertEqual(stored[0]["admission"]["reviewed_by"], "reviewer-pan56")


if __name__ == "__main__":
    unittest.main(verbosity=2)
