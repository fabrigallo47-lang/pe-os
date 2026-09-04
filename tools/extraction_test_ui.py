#!/usr/bin/env python3
"""Side-by-side extraction test UI -- the original document on the left,
exactly what tools/extract_v2_physical.py::parse_source produced from it on the
right, linked so clicking a chunk jumps the original to the place it came
from.

The point is comparison. Every chunk carries a locator (`p7`,
`model.xlsx::Assumptions!12:20`); on its own that is a string in a JSON
dump. Rendered next to the page or the rows it names, it is either
obviously right or obviously wrong, which is the only way to tell whether
a parser change helped.

This calls the real parse_source with no shortcuts, no page caps and no
reimplemented parsing, so what you see is what the pipeline sees. PDFs go
through Granite-Docling-258M when torch/transformers/docling_core are
importable; the banner says which mode is live, because a plain-text
fallback silently producing worse chunks is the failure this tool exists
to catch.

Dev/QA tool only -- separate module and port from the product server
(app/server.py), so it can never touch vault state or production routes.
Nothing leaves this machine: uploads are held in a temp dir for the life
of the process and deleted on exit.

Run:  .venv/bin/python tools/extraction_test_ui.py [port]   (default 4192)
"""
from __future__ import annotations

import atexit
import hmac
import html
import json
import os
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import traceback
import uuid
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from urllib.parse import parse_qs, urlparse

from tools.extract_v2_physical import (
    UnsupportedSourceError,
    _granite_docling_available,
    annotate_chunk,
    assemble,
    parse_source,
    validate,
    _to_e3_manifest,
)
from tools.claim_overlap import analyze_claims
from tools.llm_provider import anthropic_client_kwargs, configured_api_key

SUPPORTED_EXTENSIONS = [
    ".pdf", ".docx", ".pptx", ".xlsx", ".xlsm",
    ".md", ".markdown", ".txt", ".html", ".htm", ".csv",
    ".srt", ".vtt", ".eml", ".mbox",
]
DECLARED_UNSUPPORTED = [".xls (legacy Excel)", ".msg (Outlook binary)"]

# Rendering the original is a preview, not the parse. These caps keep a
# 50k-row workbook from producing a browser-killing HTML table; the
# extraction itself is never capped.
MAX_RENDER_ROWS = 400
MAX_RENDER_COLS = 40

_UPLOAD_ROOT = Path(tempfile.mkdtemp(prefix="panta-extract-ui-"))
atexit.register(lambda: shutil.rmtree(_UPLOAD_ROOT, ignore_errors=True))

# sid -> {"path": Path, "kind": str}
_SESSIONS: dict[str, dict[str, Any]] = {}


# --------------------------------------------------------------------------
# Rendering the ORIGINAL document (left pane)
#
# PDFs go to the browser's own viewer, which is the real thing. Everything
# else is rendered to same-origin HTML so the parent page can script it --
# that is what makes "click a chunk, scroll the original to it" possible.
# --------------------------------------------------------------------------

RENDER_STYLE = """
<style>
  :root { color-scheme: light; }
  body { font: 13px/1.55 -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
         margin: 0; padding: 18px 20px; color: #16181d; background: #fff; }
  table.grid { border-collapse: collapse; font: 11.5px/1.4 ui-monospace, SFMono-Regular, Menlo, monospace; }
  table.grid th, table.grid td { border: 1px solid #dfe3e8; padding: 2px 6px;
         max-width: 260px; overflow: hidden; text-overflow: ellipsis; white-space: nowrap; }
  table.grid th { background: #f2f4f7; color: #5b6472; font-weight: 600; position: sticky; top: 0; z-index: 2; }
  table.grid th.rownum { left: 0; z-index: 3; }
  table.grid td.rownum { background: #f2f4f7; color: #8a929e; text-align: right;
         position: sticky; left: 0; font-variant-numeric: tabular-nums; }
  table.grid td.formula { color: #0b6b3a; }
  tr.hit td { background: #fff3c4 !important; }
  tr.hit td.rownum { background: #ffe27a !important; color: #6a4b00; }
  .sheet { display: none; }
  .sheet.active { display: block; }
  .tabs { display: flex; flex-wrap: wrap; gap: 4px; margin-bottom: 12px;
          position: sticky; top: 0; background: #fff; padding-bottom: 6px; z-index: 5; }
  .tabs button { font: 12px ui-monospace, Menlo, monospace; padding: 4px 10px; cursor: pointer;
          border: 1px solid #d4d9e0; background: #f7f8fa; border-radius: 5px; color: #3d4450; }
  .tabs button.active { background: #16181d; color: #fff; border-color: #16181d; }
  .truncated { color: #9aa2ae; font-size: 11.5px; margin-top: 8px; font-style: italic; }
  pre.raw { font: 12px/1.6 ui-monospace, SFMono-Regular, Menlo, monospace;
            white-space: pre-wrap; word-break: break-word; margin: 0; }
  pre.raw .ln { color: #b6bcc6; user-select: none; display: inline-block;
            width: 3.5em; text-align: right; padding-right: 1em; }
  mark.hit { background: #fff3c4; box-shadow: 0 0 0 3px #fff3c4; }
  .slide { border: 1px solid #d4d9e0; border-radius: 6px; margin-bottom: 16px; background: #fff;
           position: relative; overflow: hidden; }
  .slide .num { position: absolute; top: 6px; right: 8px; font: 11px ui-monospace, Menlo, monospace;
           color: #9aa2ae; z-index: 2; }
  .shape { position: absolute; font-size: 11px; line-height: 1.3; overflow: hidden;
           border: 1px dashed #d9dde3; padding: 2px 3px; box-sizing: border-box; }
  .note { color: #6b7280; font-size: 12px; background: #f7f8fa; border: 1px solid #e5e8ec;
          border-radius: 6px; padding: 8px 10px; margin-bottom: 14px; }
</style>
"""


def _render_page(inner: str, script: str = "") -> str:
    return (f"<!doctype html><html><head><meta charset='utf-8'>{RENDER_STYLE}</head>"
            f"<body>{inner}<script>{script}</script></body></html>")


def _render_text(path: Path) -> str:
    """Raw bytes as the parser sees them, with line numbers.

    Deliberately not prettified: for .md/.csv/.eml/.mbox/.srt the literal
    text IS the original, and any prettifying would hide exactly the
    whitespace and delimiter bugs this pane is meant to expose.
    """
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines() or [""]
    body = "".join(
        f"<span class='ln'>{i}</span>{html.escape(line)}\n"
        for i, line in enumerate(lines, 1)
    )
    return _render_page(f"<pre class='raw' id='raw'>{body}</pre>")


def _render_html_source(path: Path) -> str:
    # An .html source is shown as its own rendered self -- that is what a
    # reader of the artifact sees, and so what the chunks should match.
    return path.read_text(encoding="utf-8", errors="replace")


def _render_docx(path: Path) -> str:
    """Visual DOCX render via macOS `textutil` (local, no network, no new dep).

    textutil keeps headings, bold and real table grids, which is what makes
    a DOCX table chunk checkable at a glance. If it is unavailable or
    fails, fall back to declaring that rather than silently showing a
    worse view that could be mistaken for the document.
    """
    try:
        out = subprocess.run(
            ["textutil", "-convert", "html", "-stdout", str(path)],
            capture_output=True, timeout=30, check=True,
        )
        return out.stdout.decode("utf-8", errors="replace")
    except (OSError, subprocess.SubprocessError) as exc:
        return _render_page(
            f"<div class='note'>Visual render unavailable: <code>textutil</code> "
            f"failed ({html.escape(str(exc))}). The extraction on the right is "
            f"unaffected -- only this preview is missing.</div>"
        )


def _render_xlsx(path: Path) -> str:
    """The workbook as a real grid: sheet tabs, column letters, row numbers.

    Both workbooks are opened because a cell's formula and its cached
    result are different facts -- parse_xlsx keeps both, so the preview
    has to show both or a reviewer cannot tell which one a chunk quoted.
    The formula is shown when present, with the cached value beside it.
    """
    import openpyxl
    from openpyxl.utils import get_column_letter

    formulas = openpyxl.load_workbook(path, data_only=False, read_only=True)
    values = openpyxl.load_workbook(path, data_only=True, read_only=True)

    tabs, panes = [], []
    for idx, sheet_name in enumerate(formulas.sheetnames):
        ws, vws = formulas[sheet_name], values[sheet_name]
        active = " active" if idx == 0 else ""
        tabs.append(
            f"<button class='tab{active}' data-sheet='{html.escape(sheet_name)}' "
            f"onclick=\"showSheet('{html.escape(sheet_name)}')\">{html.escape(sheet_name)}</button>"
        )

        rows_html, n_cols, truncated_rows = [], 0, False
        for r, (row, vrow) in enumerate(zip(ws.iter_rows(), vws.iter_rows()), start=1):
            if r > MAX_RENDER_ROWS:
                truncated_rows = True
                break
            cells = []
            for c, (cell, vcell) in enumerate(zip(row, vrow), start=1):
                if c > MAX_RENDER_COLS:
                    break
                raw, cached = cell.value, vcell.value
                is_formula = isinstance(raw, str) and raw.startswith("=")
                if is_formula:
                    shown = f"{raw}  → {cached!r}" if cached is not None else str(raw)
                else:
                    shown = "" if raw is None else str(raw)
                cls = " class='formula'" if is_formula else ""
                title = html.escape(shown, quote=True)
                cells.append(
                    f"<td{cls} title='{title}'>{html.escape(shown)}</td>"
                )
            n_cols = max(n_cols, len(cells))
            body = "".join(cells)
            rows_html.append(
                f"<tr data-row='{r}'><td class='rownum'>{r}</td>{body}</tr>"
            )

        header = "".join(
            f"<th>{get_column_letter(c)}</th>" for c in range(1, max(n_cols, 1) + 1)
        )
        note = (f"<div class='truncated'>preview stops at row {MAX_RENDER_ROWS} / "
                f"column {MAX_RENDER_COLS} -- extraction read the whole sheet.</div>"
                if truncated_rows else "")
        panes.append(
            f"<div class='sheet{active}' data-sheet='{html.escape(sheet_name)}'>"
            f"<table class='grid'><thead><tr><th class='rownum'>#</th>{header}</tr></thead>"
            f"<tbody>{''.join(rows_html)}</tbody></table>{note}</div>"
        )

    script = """
      function showSheet(name) {
        document.querySelectorAll('.sheet').forEach(function (s) {
          s.classList.toggle('active', s.dataset.sheet === name);
        });
        document.querySelectorAll('.tab').forEach(function (t) {
          t.classList.toggle('active', t.dataset.sheet === name);
        });
      }
      // Called from the parent frame when a chunk is clicked.
      window.panta_focus = function (sheet, fromRow, toRow) {
        if (sheet) showSheet(sheet);
        document.querySelectorAll('tr.hit').forEach(function (t) { t.classList.remove('hit'); });
        var pane = sheet
          ? document.querySelector('.sheet[data-sheet="' + CSS.escape(sheet) + '"]')
          : document.querySelector('.sheet.active');
        if (!pane) return false;
        var first = null;
        for (var r = fromRow; r <= toRow; r++) {
          var tr = pane.querySelector('tr[data-row="' + r + '"]');
          if (tr) { tr.classList.add('hit'); if (!first) first = tr; }
        }
        if (first) first.scrollIntoView({ block: 'center', behavior: 'smooth' });
        return !!first;
      };
    """
    return _render_page(
        f"<div class='tabs'>{''.join(tabs)}</div>{''.join(panes)}", script
    )


def _render_pptx(path: Path) -> str:
    """Slides laid out at true shape coordinates via python-pptx.

    Not a pixel render -- no local renderer exists without LibreOffice --
    but positions, sizes and reading order are real, which is enough to
    see whether a slide chunk picked up the title, the body and the table
    it should have, and in the right order.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return _render_page(
            "<div class='note'>python-pptx is not installed, so no slide preview. "
            "The extraction on the right is unaffected.</div>"
        )

    prs = Presentation(str(path))
    scale = 900 / float(prs.slide_width or Emu(9144000))
    sw = float(prs.slide_width) * scale
    sh = float(prs.slide_height) * scale

    out = ["<div class='note'>Structural preview: real shape positions and sizes, "
           "not a pixel render (no local renderer without LibreOffice).</div>"]
    for i, slide in enumerate(prs.slides, 1):
        shapes = []
        for shape in slide.shapes:
            try:
                left, top = float(shape.left) * scale, float(shape.top) * scale
                width, height = float(shape.width) * scale, float(shape.height) * scale
            except (TypeError, ValueError):
                continue  # a shape with no explicit geometry cannot be placed
            if shape.has_table:
                rows = [
                    "<tr>" + "".join(f"<td>{html.escape(c.text)}</td>" for c in row.cells) + "</tr>"
                    for row in shape.table.rows
                ]
                inner = f"<table class='grid'>{''.join(rows)}</table>"
            elif shape.has_text_frame:
                inner = html.escape(shape.text_frame.text).replace("\n", "<br>")
            else:
                inner = f"<i style='color:#9aa2ae'>&lt;{html.escape(shape.shape_type and str(shape.shape_type) or 'shape')}&gt;</i>"
            shapes.append(
                f"<div class='shape' style='left:{left:.0f}px;top:{top:.0f}px;"
                f"width:{width:.0f}px;height:{height:.0f}px'>{inner}</div>"
            )
        out.append(
            f"<div class='slide' data-slide='{i}' style='width:{sw:.0f}px;height:{sh:.0f}px'>"
            f"<span class='num'>slide {i}</span>{''.join(shapes)}</div>"
        )

    script = """
      window.panta_focus = function (sheet, fromRow, toRow) {
        var el = document.querySelector('.slide[data-slide="' + fromRow + '"]');
        if (!el) return false;
        document.querySelectorAll('.slide').forEach(function (s) { s.style.outline = ''; });
        el.style.outline = '3px solid #f0b429';
        el.scrollIntoView({ block: 'start', behavior: 'smooth' });
        return true;
      };
    """
    return _render_page("".join(out), script)


def _render_csv(path: Path) -> str:
    import csv as csvmod

    with path.open(newline="", encoding="utf-8", errors="replace") as fh:
        rows = list(csvmod.reader(fh))
    truncated = len(rows) > MAX_RENDER_ROWS
    body = "".join(
        f"<tr data-row='{i}'><td class='rownum'>{i}</td>"
        + "".join(f"<td>{html.escape(c)}</td>" for c in row[:MAX_RENDER_COLS])
        + "</tr>"
        for i, row in enumerate(rows[:MAX_RENDER_ROWS], 1)
    )
    note = (f"<div class='truncated'>preview stops at row {MAX_RENDER_ROWS} -- "
            f"extraction read all {len(rows)} rows.</div>" if truncated else "")
    return _render_page(f"<table class='grid'><tbody>{body}</tbody></table>{note}")


# suffix -> (kind for the left pane, renderer or None for "serve the file itself")
_RENDERERS = {
    ".xlsx": _render_xlsx, ".xlsm": _render_xlsx,
    ".docx": _render_docx,
    ".pptx": _render_pptx,
    ".csv": _render_csv,
    ".html": _render_html_source, ".htm": _render_html_source,
}


def _kind_for(path: Path) -> str:
    return "pdf" if path.suffix.lower() == ".pdf" else "html"

# --------------------------------------------------------------------------
# Server
#
# stdlib http.server, not a web framework. Partly because CLAUDE.md keeps
# v1 on stdlib + PyYAML, and partly for a measured reason: importing
# fastapi in this venv takes minutes, because pydantic's plugin loader
# stats the metadata of all ~200 installed packages and this working tree
# sits on a slow filesystem. Four routes do not justify paying that on
# every restart of a tool whose whole point is a fast edit/retry loop.
# --------------------------------------------------------------------------

# Resolved once, off the request path, so the port is live immediately and
# a slow first import never looks like the parser hanging.
#
# This mattered more than it sounds: with .venv inside the iCloud-synced
# Desktop, `import transformers` blocked for ~20 minutes and then failed
# with OSError(EINVAL-ish) "Operation canceled" while iCloud tried to
# materialise evicted .dist-info files. Moving the venv to ~/venvs fixed
# it (4s), but the async resolve stays: an unavailable model should
# degrade to a banner, never to a dead server.
_MODE: dict[str, str] | None = None


def _resolve_pdf_mode() -> dict[str, str]:
    global _MODE
    try:
        return _resolve_pdf_mode_inner()
    except Exception as exc:  # noqa: BLE001
        # _granite_docling_available() only catches ImportError, so a
        # filesystem error out of the venv (iCloud returning ECANCELED for
        # an evicted .dist-info file, seen on this machine) escapes and
        # would otherwise kill this thread silently, leaving the banner
        # saying "checking" forever. Say what actually broke instead.
        _MODE = {"label": "unavailable", "level": "warn", "detail":
                 f"could not resolve the PDF model: {type(exc).__name__}: {exc}"}
        return _MODE


def _resolve_pdf_mode_inner() -> dict[str, str]:
    global _MODE
    if _granite_docling_available():
        _MODE = {"label": "Granite-Docling-258M", "detail":
                 "real model: tables extracted as structure, pictures declared as typed markers",
                 "level": "ok"}
    else:
        _MODE = {"label": "plain-text fallback", "detail":
                 "torch / transformers / docling_core not importable -- no tables, needs a text "
                 "layer (see requirements-pdf-ml.txt)", "level": "warn"}
    return _MODE


def _pdf_mode() -> dict[str, str]:
    if _MODE is None:
        return {"label": "checking\u2026", "level": "pending", "detail":
                "importing torch / transformers / docling_core (a few seconds); "
                "non-PDF formats already work"}
    return _MODE


# Shared-secret gate. Deliberately minimal -- this is a dev/QA tool, not a
# product surface -- but it is the difference between "reachable only by
# the person at this machine" and "reachable by the internet", which is
# the difference that matters once real deal documents are uploaded.
AUTH_TOKEN = os.environ.get("PE_OS_UI_TOKEN", "")


def _authorised(handler: BaseHTTPRequestHandler) -> bool:
    if not AUTH_TOKEN:
        return True     # loopback-only mode; main() enforces that pairing
    supplied = ""
    query = parse_qs(urlparse(handler.path).query).get("token")
    if query:
        supplied = query[0]
    else:
        cookie = handler.headers.get("Cookie") or ""
        for part in cookie.split(";"):
            name, _, value = part.strip().partition("=")
            if name == "pe_os_ui_token":
                supplied = value
                break
    # constant-time: a timing oracle on a shared secret is cheap to avoid
    return hmac.compare_digest(supplied, AUTH_TOKEN)


CONTENT_TYPES = {
    ".pdf": "application/pdf", ".html": "text/html", ".htm": "text/html",
    ".csv": "text/csv", ".txt": "text/plain",
}


# The classic-Docling engine lives in its own virtualenv (see
# tools/docling_engine.py for why) and is reached by subprocess.
DOCLING_PYTHON = Path(
    os.environ.get("PE_OS_DOCLING_PYTHON", Path.home() / "venvs/pe-os-docling/bin/python")
)
PADDLE_PYTHON = Path(
    os.environ.get("PE_OS_PADDLE_PYTHON", Path.home() / "venvs/pe-os-paddle13/bin/python")
)
ENGINES = {
    "granite": "Granite-Docling-258M (VLM, reads the page image)",
    "docling": "Docling classic (layout model + TableFormer)",
    "paddle": "PaddleOCR-VL (PP-DocLayoutV3 + 0.9B VLM)",
}
# engine -> (interpreter, engine script)
SUBPROCESS_ENGINES = {
    "docling": (lambda: DOCLING_PYTHON, "docling_engine.py"),
    "paddle": (lambda: PADDLE_PYTHON, "paddle_engine.py"),
}


def _docling_available() -> bool:
    return DOCLING_PYTHON.exists()


def _paddle_available() -> bool:
    return PADDLE_PYTHON.exists()


def _subprocess_convert_page(pdf: Path, engine: str):
    """Run the classic pipeline once for the whole PDF, then hand parse_pdf a
    per-page lookup shaped like the Granite converter.

    One subprocess per document, not per page: the classic pipeline is a
    document-level converter, and re-running it per page would both be
    slower and give it less context than it normally has.
    """
    # The result comes back through a file: TableFormer writes warnings to
    # stdout, so stdout cannot carry JSON (see tools/docling_engine.py).
    interpreter, script = SUBPROCESS_ENGINES[engine]
    out_json = pdf.with_suffix(f".{engine}.json")
    proc = subprocess.run(
        [str(interpreter()), str(ROOT / "tools" / script),
         str(pdf), str(out_json)],
        capture_output=True, timeout=3600,
    )
    if not out_json.exists():
        tail = proc.stderr.decode("utf-8", errors="replace")[-500:]
        raise RuntimeError(f"{engine} engine wrote no result (exit {proc.returncode}); stderr: {tail}")
    payload = json.loads(out_json.read_text())
    if "error" in payload:
        raise RuntimeError(f"{engine} engine failed: {payload['error']}")

    pages = payload.get("pages", {})

    def convert(image: Any, page_num: int) -> tuple[str, list[str]]:
        page = pages.get(str(page_num)) or {}
        return page.get("markdown", ""), page.get("pictures", [])

    convert.warnings = payload.get("warnings", [])   # surfaced in the UI
    return convert


def _parse_page_spec(spec: str, total: int) -> list[int]:
    """`12-14`, `7`, `1-4,9` -> [12,13,14] / [7] / [1,2,3,4,9] (1-based)."""
    pages: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part.lstrip("-"):
            first, _, last = part.partition("-")
            pages.extend(range(int(first), int(last) + 1))
        else:
            pages.append(int(part))
    seen, out = set(), []
    for n in pages:
        if 1 <= n <= total and n not in seen:
            seen.add(n)
            out.append(n)
    return out


def _slice_pdf(path: Path, pages: list[int]) -> Path:
    """Write the requested pages out as their own PDF.

    Subsetting the file rather than teaching parse_pdf a page argument
    keeps the thing under test untouched -- the parser still runs its
    normal whole-document path, just over a shorter document. At ~49s
    per page on CPU this is what makes iterating on a 43-page deck
    possible at all.
    """
    import pypdfium2 as pdfium

    src = pdfium.PdfDocument(str(path))
    dst = pdfium.PdfDocument.new()
    dst.import_pages(src, [n - 1 for n in pages])   # pdfium is 0-based
    out = path.with_name(f"{path.stem}__p{pages[0]}-{pages[-1]}{path.suffix}")
    dst.save(str(out))
    return out


def _extract_payload(name: str, data: bytes, page_spec: str = "",
                     engine: str = "granite") -> dict[str, Any]:
    """Store the upload, run the real parser over it, describe the result."""
    sid = uuid.uuid4().hex
    folder = _UPLOAD_ROOT / sid
    folder.mkdir(parents=True, exist_ok=True)
    # Dispatch is by suffix, so the stored copy must keep the real filename.
    path = folder / (Path(name).name or "upload")
    path.write_bytes(data)
    _SESSIONS[sid] = {"path": path, "kind": _kind_for(path)}

    payload: dict[str, Any] = {
        "sid": sid,
        "filename": path.name,
        "kind": _kind_for(path),
        "size_bytes": len(data),
        "pdf_mode": _pdf_mode(),
        "page_map": None,
        "engine": engine,
        "engine_label": ENGINES.get(engine, engine),
    }

    # The left pane always shows the whole uploaded file, so _SESSIONS keeps
    # the original; only the parser sees the subset.
    target = path
    if page_spec.strip() and path.suffix.lower() == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(path) as doc:
                total = len(doc.pages)
            wanted = _parse_page_spec(page_spec, total)
            if wanted:
                target = _slice_pdf(path, wanted)
                # chunk page N of the slice is really page page_map[N-1]
                payload["page_map"] = wanted
                payload["total_pages"] = total
        except Exception as exc:  # noqa: BLE001 -- fall back to the whole file
            payload["page_warning"] = f"page range ignored ({type(exc).__name__}: {exc})"

    # Timer starts BEFORE engine setup: the docling model runs in its
    # subprocess here, and timing only the parse_source call afterwards
    # reported 0.09s for a 20-second conversion -- flattering the engine
    # whose work happens outside the loop, which is exactly the number a
    # model comparison must not get wrong.
    started = time.perf_counter()

    convert_page = None
    if engine in SUBPROCESS_ENGINES and path.suffix.lower() == ".pdf":
        interpreter = SUBPROCESS_ENGINES[engine][0]()
        if not interpreter.exists():
            payload["status"] = "error"
            payload["error"] = (
                f"{engine} engine not installed: no interpreter at {interpreter}. "
                "See deploy/README.md for the environment it needs."
            )
            payload["elapsed_s"] = round(time.perf_counter() - started, 2)
            return payload
        try:
            convert_page = _subprocess_convert_page(target, engine)
            payload["engine_warnings"] = getattr(convert_page, "warnings", [])
        except Exception as exc:  # noqa: BLE001
            payload["status"] = "error"
            payload["error"] = f"{type(exc).__name__}: {exc}"
            payload["elapsed_s"] = round(time.perf_counter() - started, 2)
            return payload
        payload["engine_setup_s"] = round(time.perf_counter() - started, 2)

    try:
        chunks = parse_source(target, convert_page=convert_page)
    except UnsupportedSourceError as exc:
        payload["status"] = "rejected"
        payload["rejection"] = exc.to_dict()
    except Exception as exc:  # noqa: BLE001 -- a dev tool shows the failure
        payload["status"] = "error"
        payload["error"] = f"{type(exc).__name__}: {exc}"
        payload["traceback"] = traceback.format_exc()
    else:
        payload["status"] = "ok"
        # Kept in-process (not re-parsed) so step 2 runs against the exact
        # same Chunk objects step 1 showed, not a fresh parse that could
        # legitimately differ for a page-range slice or a re-uploaded file.
        _SESSIONS[sid]["chunks"] = chunks
        payload["chunks"] = [{
            "chunk_id": c.chunk_id,
            "locator": c.locator,
            "body": c.body,
            "word_count": c.word_count,
            "section_heading": c.section_heading,
            "page_or_slide_number": c.page_or_slide_number,
            "provenance": c.provenance,
            "period_context": c.period_context,
        } for c in chunks]
    payload["elapsed_s"] = round(time.perf_counter() - started, 2)
    return payload


# --------------------------------------------------------------------------
# Step 2: semantic extraction (L2) + the claim-identity graph
#
# Reuses annotate_chunk/validate/assemble/_to_e3_manifest exactly as
# tools/extract_v2_physical.py's own CLI does, and claim_overlap.
# analyze_claims() exactly as it runs in production -- this tool shows
# what the pipeline does, it does not run a second, UI-only approximation
# of it.
# --------------------------------------------------------------------------

_ANTHROPIC_CLIENT: Any = None


def _anthropic_client() -> Any:
    global _ANTHROPIC_CLIENT
    if _ANTHROPIC_CLIENT is None:
        import anthropic
        api_key = configured_api_key()
        if not api_key:
            raise RuntimeError(
                "no Anthropic API key configured (checked the usual env/.env "
                "locations) -- semantic extraction needs one, physical "
                "extraction does not"
            )
        _ANTHROPIC_CLIENT = anthropic.Anthropic(**anthropic_client_kwargs(api_key))
    return _ANTHROPIC_CLIENT


def _semantic_payload(sid: str) -> dict[str, Any]:
    """Run real L2 + L3 + L4 over the chunks step 1 already parsed."""
    session = _SESSIONS.get(sid)
    if not session or "chunks" not in session:
        return {"status": "error", "error": "run extraction (step 1) first"}
    chunks = session["chunks"]

    started = time.perf_counter()
    client = _anthropic_client()
    raw_claims = []
    chunk_errors = []
    for chunk in chunks:
        try:
            raw_claims.extend(annotate_chunk(chunk, client, deal="ui-test"))
        except Exception as exc:  # noqa: BLE001 -- one bad chunk must not void the rest
            chunk_errors.append({"chunk_id": chunk.chunk_id, "error": f"{type(exc).__name__}: {exc}"})

    canonicals = [validate(r) for r in raw_claims]
    graph = assemble(canonicals)
    manifest = _to_e3_manifest(graph, deal="ui-test", manifest="UI",
                               sources_used=[{"source_id": "SRC-UI-UPLOAD"}])
    overlap = analyze_claims(manifest)
    session["manifest"] = manifest  # step 3 (archetype) reads this back

    return {
        "status": "ok",
        "elapsed_s": round(time.perf_counter() - started, 2),
        "chunk_errors": chunk_errors,
        "manifest": manifest,
        "overlap": overlap,
    }


# --------------------------------------------------------------------------
# Step 3: deal-archetype signal
#
# Honest label: THERE IS NO PRODUCTION CLASSIFIER. METRIC_ENUM is built for
# LBO/buyout underwriting (see docs/dizionario_estrazione.md); a claim
# outside that vocabulary lands on "Other" (PAN-117) rather than being
# force-fit onto the nearest-sounding entry. This step turns that same
# signal -- which labels landed on Other, and what they contain -- into a
# rough, visible-as-a-guess-not-a-fact hint of which known vocabulary (if
# any) the document actually matches. It never blocks or relabels a claim.
# --------------------------------------------------------------------------

_FAMILY_KEYWORDS: dict[str, tuple[str, ...]] = {
    "LBO / buyout": (
        "sponsor equity", "leverage", "covenant", "ebitda", "first-lien",
        "revolver", "ddtl", "rollover", "moic", "irr", "entry multiple",
        "exit multiple", "net debt",
    ),
    "Venture / growth": (
        "pre-money", "post-money", "round", "runway", "arr", "mrr",
        "burn rate", "cap table", "dilution", "safe", "seed", "series",
        "ownership stake", "funding",
    ),
}


def _archetype_payload(sid: str) -> dict[str, Any]:
    session = _SESSIONS.get(sid)
    manifest = session.get("manifest") if session else None
    if not manifest:
        return {"status": "error", "error": "run semantic extraction (step 2) first"}

    fields = manifest["extraction_metadata"]["compiler_fields_per_claim"]
    total = len(fields)
    other = [f for f in fields if f.get("metric") == "Other"]
    other_labels = " | ".join((f.get("metric_label") or "").lower() for f in other)
    known_metrics = [f.get("metric") for f in fields if f.get("metric") != "Other"]

    family_hits: dict[str, list[str]] = {name: [] for name in _FAMILY_KEYWORDS}
    for name, keywords in _FAMILY_KEYWORDS.items():
        for kw in keywords:
            if kw in other_labels:
                family_hits[name].append(kw)
        # A known METRIC_ENUM value that is itself LBO-shaped vocabulary
        # counts too -- Sponsor Equity/Leverage/etc are only IN the enum
        # because that is the domain it was built for.
        for metric in known_metrics:
            if metric and metric.lower() in _FAMILY_KEYWORDS.get(name, ()):
                family_hits[name].append(metric)

    return {
        "status": "ok",
        "note": "euristica sperimentale -- nessun classificatore di archetipo "
                "esiste in produzione. Segnale grezzo da parole chiave, non "
                "una decisione. Vedi G5 in docs/dizionario_estrazione.md.",
        "total_claims": total,
        "other_count": len(other),
        "other_ratio": round(len(other) / total, 2) if total else None,
        "family_hits": {name: sorted(set(hits)) for name, hits in family_hits.items()},
    }


class Handler(BaseHTTPRequestHandler):
    protocol_version = "HTTP/1.1"

    def log_message(self, fmt: str, *args: Any) -> None:
        sys.stderr.write("  %s\n" % (fmt % args))

    # -- helpers ----------------------------------------------------------
    def _send(self, body: bytes, ctype: str, code: int = 200,
              extra: dict[str, str] | None = None) -> None:
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        for k, v in (extra or {}).items():
            self.send_header(k, v)
        self.end_headers()
        self.wfile.write(body)

    def _html(self, markup: str, code: int = 200) -> None:
        self._send(markup.encode("utf-8"), "text/html; charset=utf-8", code)

    def _json(self, obj: Any, code: int = 200) -> None:
        self._send(json.dumps(obj).encode("utf-8"), "application/json", code)

    def _session(self, sid: str) -> dict[str, Any] | None:
        return _SESSIONS.get(sid)

    # -- routes -----------------------------------------------------------
    def do_GET(self) -> None:
        if not _authorised(self):
            return self._html("<h1>401</h1><p>append ?token=... to the URL</p>", 401)
        route = urlparse(self.path).path
        if route == "/":
            token = parse_qs(urlparse(self.path).query).get("token")
            if token and AUTH_TOKEN:
                # Store it once so later fetches (which carry no query
                # string) authenticate without the secret in every URL.
                body = index().encode("utf-8")
                self.send_response(200)
                self.send_header("Content-Type", "text/html; charset=utf-8")
                self.send_header("Content-Length", str(len(body)))
                self.send_header(
                    "Set-Cookie",
                    f"pe_os_ui_token={token[0]}; Path=/; HttpOnly; SameSite=Strict",
                )
                self.end_headers()
                return self.wfile.write(body)
            return self._html(index())
        if route == "/mode":
            return self._json({**_pdf_mode(),
                               "docling_available": _docling_available(),
                               "paddle_available": _paddle_available()})
        if route.startswith("/original/"):
            session = self._session(route.split("/original/", 1)[1])
            if not session:
                return self._html("<p>unknown or expired upload</p>", 404)
            path: Path = session["path"]
            ctype = CONTENT_TYPES.get(path.suffix.lower(), "application/octet-stream")
            # inline, so the browser's own PDF viewer takes it over downloading
            return self._send(path.read_bytes(), ctype,
                              extra={"Content-Disposition": "inline"})
        if route.startswith("/render/"):
            session = self._session(route.split("/render/", 1)[1])
            if not session:
                return self._html("<p>unknown or expired upload</p>", 404)
            return self._html(render_original(session["path"]))
        self._html("<p>not found</p>", 404)

    def do_POST(self) -> None:
        if not _authorised(self):
            return self._json({"status": "error", "error": "unauthorised"}, 401)
        parsed = urlparse(self.path)
        query = parse_qs(parsed.query)

        if parsed.path == "/semantic":
            sid = (query.get("sid") or [""])[0]
            try:
                return self._json(_semantic_payload(sid))
            except Exception as exc:  # noqa: BLE001 -- never hang the browser
                return self._json({"status": "error", "error": f"{type(exc).__name__}: {exc}",
                                   "traceback": traceback.format_exc()}, 500)

        if parsed.path == "/archetype":
            sid = (query.get("sid") or [""])[0]
            try:
                return self._json(_archetype_payload(sid))
            except Exception as exc:  # noqa: BLE001
                return self._json({"status": "error", "error": f"{type(exc).__name__}: {exc}",
                                   "traceback": traceback.format_exc()}, 500)

        if parsed.path != "/extract":
            return self._html("<p>not found</p>", 404)
        # The browser posts the raw bytes with the filename in the query
        # string: no multipart parsing, and `cgi` is gone in 3.13+ anyway.
        name = (query.get("name") or ["upload"])[0]
        page_spec = (query.get("pages") or [""])[0]
        engine = (query.get("engine") or ["granite"])[0]
        data = self.rfile.read(int(self.headers.get("Content-Length") or 0))
        try:
            self._json(_extract_payload(name, data, page_spec, engine))
        except Exception as exc:  # noqa: BLE001 -- never hang the browser
            self._json({"status": "error", "error": f"{type(exc).__name__}: {exc}",
                        "traceback": traceback.format_exc(), "elapsed_s": 0,
                        "filename": name, "kind": "html", "sid": "",
                        "size_bytes": len(data), "pdf_mode": _pdf_mode()}, 500)


def render_original(path: Path) -> str:
    renderer = _RENDERERS.get(path.suffix.lower())
    try:
        return renderer(path) if renderer else _render_text(path)
    except Exception as exc:  # noqa: BLE001 -- a preview failure must never
        # look like an extraction failure; say which one broke.
        return _render_page(
            f"<div class='note'>Preview failed: "
            f"<code>{html.escape(type(exc).__name__)}: {html.escape(str(exc))}</code><br>"
            f"This is the <b>preview</b> only -- the extraction result on the right "
            f"is what parse_source actually returned.</div>"
        )


INDEX_HTML = r"""<!doctype html>
<html><head><meta charset="utf-8"><title>Extraction test — PANTA</title>
<style>
  :root { color-scheme: light; --line: #e3e6ea; --muted: #6b7280; --ink: #16181d; }
  * { box-sizing: border-box; }
  body { margin: 0; height: 100vh; display: flex; flex-direction: column; color: var(--ink);
         font: 13px/1.5 -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #f4f5f7; }
  header { display: flex; align-items: center; gap: 14px; padding: 10px 16px; background: #fff;
           border-bottom: 1px solid var(--line); flex: none; flex-wrap: wrap; }
  h1 { font-size: 14px; margin: 0; font-weight: 650; letter-spacing: -0.01em; }
  .badge { font: 11px ui-monospace, SFMono-Regular, Menlo, monospace; padding: 3px 9px;
           border-radius: 999px; border: 1px solid; }
  .badge.ok   { background: #e8f6ee; color: #106b3c; border-color: #b7e3ca; }
  .badge.warn { background: #fdf1dc; color: #8a5a00; border-color: #f2ddb0; }
  .badge.pending { background: #eef1f5; color: #55607a; border-color: #d7dce4; }
  .pages { width: 128px; padding: 5px 8px; border: 1px solid #d4d9e0; border-radius: 6px;
           font: 12px ui-monospace, Menlo, monospace; }
  .spacer { flex: 1; }
  .stats { font: 11.5px ui-monospace, Menlo, monospace; color: var(--muted); }
  button.go { background: var(--ink); color: #fff; border: 0; padding: 6px 14px; border-radius: 6px;
              cursor: pointer; font-size: 13px; font-weight: 500; }
  button.go:disabled { opacity: .45; cursor: default; }
  main { flex: 1; display: flex; min-height: 0; }
  .pane { display: flex; flex-direction: column; min-width: 0; background: #fff; }
  #left  { flex: 1 1 0; border-right: 0; }
  #right { flex: 1 1 0; overflow-y: auto; border-left: 1px solid var(--line); }
  #divider { flex: none; width: 6px; cursor: col-resize; background: var(--line); }
  #divider:hover { background: #b9c0c9; }
  .pane-head { flex: none; padding: 7px 14px; border-bottom: 1px solid var(--line); background: #fafbfc;
               font: 11.5px ui-monospace, Menlo, monospace; color: var(--muted);
               display: flex; justify-content: space-between; gap: 10px; }
  #frame { flex: 1; width: 100%; border: 0; min-height: 0; }
  #chunks { padding: 10px 12px 40px; }
  .drop { margin: 40px auto; max-width: 460px; text-align: center; color: var(--muted);
          border: 1.5px dashed #c9cfd7; border-radius: 10px; padding: 34px 20px; background: #fafbfc; }
  .chunk { border: 1px solid var(--line); border-radius: 7px; margin-bottom: 9px; overflow: hidden;
           cursor: pointer; background: #fff; }
  .chunk:hover { border-color: #aeb6c0; }
  .chunk.sel { border-color: #f0b429; box-shadow: 0 0 0 2px #fdf1dc; }
  .chunk-head { background: #f7f8fa; padding: 6px 10px; display: flex; justify-content: space-between;
                gap: 10px; font: 11.5px ui-monospace, Menlo, monospace; }
  .chunk-head .loc { color: #0b4f9e; font-weight: 600; word-break: break-all; }
  .chunk-head .meta { color: var(--muted); white-space: nowrap; }
  .chunk-body { padding: 9px 10px; white-space: pre-wrap; word-break: break-word;
                font: 11.5px/1.55 ui-monospace, SFMono-Regular, Menlo, monospace; max-height: 300px;
                overflow-y: auto; }
  details.prov { border-top: 1px dashed var(--line); }
  details.prov summary { padding: 5px 10px; font-size: 11px; color: var(--muted); cursor: pointer; }
  details.prov pre { margin: 0; padding: 0 10px 9px; font: 11px/1.5 ui-monospace, Menlo, monospace;
                     color: #3d4450; white-space: pre-wrap; word-break: break-word; }
  .panel { margin: 14px 12px; border-radius: 8px; padding: 12px 14px; font-size: 12.5px; }
  .panel.reject { background: #fff7ed; border: 1px solid #f2d5ae; }
  .panel.error  { background: #fef2f2; border: 1px solid #f3c6c6; }
  .panel h2 { margin: 0 0 8px; font-size: 13px; }
  .panel.reject h2 { color: #92400e; } .panel.error h2 { color: #b42318; }
  .panel dt { font-weight: 650; margin-top: 6px; font-size: 11.5px; }
  .panel dd { margin: 1px 0 0; font: 11.5px ui-monospace, Menlo, monospace; word-break: break-word; }
  .hint { color: var(--muted); font-size: 11.5px; padding: 0 2px 8px; }
  button.go.step { background: #fff; color: var(--ink); border: 1px solid #d4d9e0; }
  button.go.step:disabled { opacity: .4; }
  #tabs { display: flex; gap: 2px; padding: 0 16px; background: #fff; border-bottom: 1px solid var(--line); flex: none; }
  .tab { border: 0; background: none; padding: 8px 12px; font: 12px ui-monospace, Menlo, monospace;
         color: var(--muted); cursor: pointer; border-bottom: 2px solid transparent; }
  .tab:disabled { opacity: .35; cursor: default; }
  .tab.sel { color: var(--ink); border-bottom-color: #f0b429; font-weight: 600; }
  .tabpane { padding: 10px 12px 40px; }
  table.claims { width: 100%; border-collapse: collapse; font: 11.5px ui-monospace, Menlo, monospace; }
  table.claims th { text-align: left; font-size: 10.5px; text-transform: uppercase; letter-spacing: .04em;
                    color: var(--muted); padding: 6px 8px; border-bottom: 1px solid var(--line); position: sticky; top: 0; background: #fff; }
  table.claims td { padding: 6px 8px; border-bottom: 1px solid #f0f1f3; vertical-align: top; }
  tr.row-other td.metric { color: #7c3aed; }
  tr.row-flagged { background: #fff7ed; }
  .pill { display: inline-block; padding: 1px 7px; border-radius: 20px; font-size: 10.5px; white-space: nowrap; }
  .pill.asserted { background: #eef1f5; color: #55607a; }
  .pill.observed { background: #e4edea; color: #2b5c56; }
  .pill.attested { background: #e3f3ea; color: #2e8b57; }
  .pill.derived  { background: #f6ebd8; color: #a9762e; }
  .flagmark { color: #b4483c; font-weight: 700; cursor: help; }
  .cluster { display: inline-block; margin: 8px 10px 8px 0; vertical-align: top; border: 1px solid var(--line); border-radius: 8px; padding: 8px; background: #fafbfc; }
  .cluster .id { font-size: 10.5px; color: var(--muted); margin-bottom: 4px; max-width: 260px; }
  .cluster.contra { border-color: #f3c6c6; background: #fef7f7; }
  .cluster.corro { border-color: #b7e3ca; background: #f3fbf6; }
  .archrow { display: flex; justify-content: space-between; padding: 5px 0; border-bottom: 1px dashed var(--line); font: 12px ui-monospace, Menlo, monospace; }
  .archnote { background: #fdf1dc; border: 1px solid #f2ddb0; color: #8a5a00; border-radius: 8px; padding: 10px 12px; margin-bottom: 12px; font-size: 12px; }
</style></head>
<body>
<header>
  <h1>PANTA pipeline test</h1>
  <span class="badge __MODE_LEVEL__" title="__MODE_DETAIL__">PDF: __MODE_LABEL__</span>
  <input type="file" id="file" accept="__ACCEPT__">
  <select id="engine" class="pages" title="Which PDF model converts each page to markdown. Everything after that step — chunking, locators, fallback — is identical, so this compares models, not pipelines.">
    <option value="granite">Granite-Docling-258M</option>
    <option value="docling">Docling classic (TableFormer)</option>
    <option value="paddle">PaddleOCR-VL</option>
  </select>
  <input type="text" id="pages" class="pages" placeholder="pages e.g. 12-14"
         title="PDF only. Blank = whole document. At ~49s/page on CPU, a range is how you iterate on a long deck.">
  <button class="go" id="go">1 · Estrazione fisica</button>
  <button class="go step" id="go2" disabled title="Chiama il vero L2 (annotate_chunk) su ogni chunk — costa una chiamata API per chunk.">2 · Estrazione semantica + grafo</button>
  <button class="go step" id="go3" disabled title="Euristica sperimentale, non un classificatore di produzione.">3 · Archetipo di deal</button>
  <span class="spacer"></span>
  <span class="stats" id="stats"></span>
</header>
<nav id="tabs">
  <button class="tab sel" data-tab="chunks">1 · Chunk fisici</button>
  <button class="tab" data-tab="claims" id="tabClaims" disabled>2 · Claim semantici</button>
  <button class="tab" data-tab="graph" id="tabGraph" disabled>2 · Grafo</button>
  <button class="tab" data-tab="archetype" id="tabArchetype" disabled>3 · Archetipo</button>
</nav>
<main>
  <section class="pane" id="left">
    <div class="pane-head"><span>original</span><span id="leftname"></span></div>
    <iframe id="frame" sandbox="allow-same-origin allow-scripts"></iframe>
  </section>
  <div id="divider"></div>
  <section class="pane" id="right">
    <div class="pane-head"><span id="paneheadlabel">extraction — parse_source()</span><span id="rightcount"></span></div>
    <div id="chunks" class="tabpane"><div class="drop">Pick a file and hit "1 · Estrazione fisica".<br>
      <small>__SUPPORTED__<br><br>Declared unsupported (rejected on purpose): __UNSUPPORTED__</small></div></div>
    <div id="claims" class="tabpane" hidden><div class="drop">Run step 1, then step 2.</div></div>
    <div id="graph" class="tabpane" hidden><div class="drop">Run step 2 first.</div></div>
    <div id="archetype" class="tabpane" hidden><div class="drop">Run step 2, then step 3.</div></div>
  </section>
</main>
<script>
const $ = (id) => document.getElementById(id);
let SID = null, KIND = null, PAGEMAP = null;

$('go').onclick = run;
$('go2').onclick = runSemantic;
$('go3').onclick = runArchetype;
$('file').onchange = () => { if ($('file').files.length) run(); };

document.querySelectorAll('.tab').forEach(tab => {
  tab.onclick = () => {
    if (tab.disabled) return;
    document.querySelectorAll('.tab').forEach(t => t.classList.remove('sel'));
    tab.classList.add('sel');
    const label = { chunks: 'extraction — parse_source()',
                     claims: 'L2 — annotate_chunk() → validate()',
                     graph: 'L4 — assemble() + claim_overlap.analyze_claims()',
                     archetype: 'deal-archetype heuristic (experimental)' }[tab.dataset.tab];
    $('paneheadlabel').textContent = label;
    document.querySelectorAll('.tabpane').forEach(p => { p.hidden = (p.id !== tab.dataset.tab); });
  };
});

async function run() {
  const f = $('file').files[0];
  if (!f) return;
  $('go').disabled = true;
  const t0 = performance.now();
  const tick = setInterval(() => {
    $('stats').textContent = 'extracting… ' + ((performance.now() - t0) / 1000).toFixed(1) + 's';
  }, 100);
  $('chunks').innerHTML = '<div class="drop">Running the real pipeline.<br><small>' +
    'PDFs with the Granite model run the VLM once per page — a long document takes minutes.</small></div>';
  try {
    // Raw body + filename in the query string: the server is stdlib
    // http.server, and multipart parsing would buy nothing here.
    const res = await fetch('/extract?name=' + encodeURIComponent(f.name) +
                            '&pages=' + encodeURIComponent($('pages').value.trim()) +
                            '&engine=' + encodeURIComponent($('engine').value),
                            { method: 'POST', body: f });
    render(await res.json());
  } catch (e) {
    $('chunks').innerHTML = '<div class="panel error"><h2>Request failed</h2><pre>' +
      esc(String(e)) + '</pre></div>';
  } finally {
    clearInterval(tick);
    $('go').disabled = false;
  }
}

function esc(s) {
  return String(s).replace(/[&<>"]/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;'}[c]));
}

// A sliced PDF renumbers from 1, but the left pane shows the whole file --
// so every page shown or jumped to is mapped back to the real document.
function truePage(n) {
  return (PAGEMAP && PAGEMAP[n - 1]) ? PAGEMAP[n - 1] : n;
}

function render(d) {
  SID = d.sid; KIND = d.kind; PAGEMAP = d.page_map;
  $('leftname').textContent = d.filename + '  ·  ' + (d.size_bytes / 1024).toFixed(0) + ' KB';
  // The browser's own PDF viewer is unreliable inside a sandboxed frame, so
  // sandbox only the rendered-HTML pane (which may be an uploaded .html).
  const frame = $('frame');
  if (d.kind === 'pdf') frame.removeAttribute('sandbox');
  else frame.setAttribute('sandbox', 'allow-same-origin allow-scripts');
  frame.src = d.kind === 'pdf' ? '/original/' + SID + '#page=1' : '/render/' + SID;

  const badge = document.querySelector('.badge');
  badge.className = 'badge ' + d.pdf_mode.level;
  badge.textContent = 'PDF: ' + d.pdf_mode.label;
  badge.title = d.pdf_mode.detail;

  if (d.status === 'rejected') {
    $('stats').textContent = 'rejected in ' + d.elapsed_s + 's';
    $('rightcount').textContent = 'REJECTED';
    $('chunks').innerHTML = '<div class="panel reject"><h2>Declared rejection — not a crash</h2>' +
      '<dl>' + Object.entries(d.rejection).map(([k, v]) =>
        '<dt>' + esc(k) + '</dt><dd>' + esc(typeof v === 'string' ? v : JSON.stringify(v)) + '</dd>'
      ).join('') + '</dl></div>';
    return;
  }
  if (d.status === 'error') {
    $('stats').textContent = 'failed in ' + d.elapsed_s + 's';
    $('rightcount').textContent = (d.engine_label || '') + ' · ERROR';
    $('chunks').innerHTML = '<div class="panel error"><h2>Unexpected error</h2><pre>' +
      esc(d.error) + '</pre>' +
      (d.traceback ? '<details><summary>traceback</summary><pre>' +
        esc(d.traceback) + '</pre></details>' : '') + '</div>';
    return;
  }

  const words = d.chunks.reduce((a, c) => a + c.word_count, 0);
  const scope = d.page_map
    ? ' · pages ' + d.page_map.join(',') + ' of ' + d.total_pages
    : '';
  $('stats').textContent = (d.engine || 'granite') + ' · ' +
                           d.chunks.length + ' chunks · ' + words + ' words · ' +
                           d.elapsed_s + 's' + scope;
  $('rightcount').textContent = (d.engine_label || '') + ' · ' + d.chunks.length + ' chunks';
  $('go2').disabled = false;
  $('tabClaims').disabled = false;
  $('tabGraph').disabled = false;
  $('chunks').innerHTML =
    '<div class="hint">Click a chunk to jump the original to where it came from.' +
      (d.page_warning ? ' <b>' + esc(d.page_warning) + '</b>' : '') + '</div>' +
    (d.engine_warnings && d.engine_warnings.length
      ? '<details class="panel reject"><summary><b>' + d.engine_warnings.length +
        ' model warning(s)</b> — the engine reported low confidence on parts of this document' +
        '</summary><pre>' + esc(d.engine_warnings.join('\n')) + '</pre></details>'
      : '') +
    d.chunks.map((c, i) => {
      const page = c.page_or_slide_number ? ' · p' + truePage(c.page_or_slide_number) : '';
      const prov = (c.provenance && Object.keys(c.provenance).length) ||
                   (c.period_context && Object.keys(c.period_context).length)
        ? '<details class="prov"><summary>provenance · period_context</summary><pre>' +
          esc(JSON.stringify({ provenance: c.provenance, period_context: c.period_context }, null, 2)) +
          '</pre></details>' : '';
      return '<div class="chunk" data-i="' + i + '">' +
        '<div class="chunk-head"><span class="loc">' + esc(c.locator) + '</span>' +
        '<span class="meta">' + c.word_count + 'w' + page + '</span></div>' +
        '<div class="chunk-body">' + esc(c.body) + '</div>' + prov + '</div>';
    }).join('');

  document.querySelectorAll('.chunk').forEach(el => {
    el.onclick = () => {
      document.querySelectorAll('.chunk.sel').forEach(x => x.classList.remove('sel'));
      el.classList.add('sel');
      focusOriginal(d.chunks[+el.dataset.i]);
    };
  });
}

async function runSemantic() {
  if (!SID) return;
  $('go2').disabled = true;
  const t0 = performance.now();
  const tick = setInterval(() => {
    $('stats').textContent = 'semantic extraction… ' + ((performance.now() - t0) / 1000).toFixed(1) + 's';
  }, 100);
  $('claims').innerHTML = '<div class="drop">Running annotate_chunk() → validate() → assemble() on every chunk.' +
    '<br><small>One real API call per chunk — this can take a while on a long document.</small></div>';
  try {
    const res = await fetch('/semantic?sid=' + encodeURIComponent(SID), { method: 'POST' });
    renderSemantic(await res.json());
  } catch (e) {
    $('claims').innerHTML = '<div class="panel error"><h2>Request failed</h2><pre>' + esc(String(e)) + '</pre></div>';
  } finally {
    clearInterval(tick);
    $('go2').disabled = false;
  }
}

function renderSemantic(d) {
  if (d.status !== 'ok') {
    $('claims').innerHTML = '<div class="panel error"><h2>Semantic extraction failed</h2><pre>' +
      esc(d.error || JSON.stringify(d, null, 2)) +
      (d.traceback ? '\n\n' + d.traceback : '') + '</pre></div>';
    return;
  }
  const manifest = d.manifest, overlap = d.overlap;
  const fieldsById = {};
  manifest.extraction_metadata.compiler_fields_per_claim.forEach(f => { fieldsById[f.claim_id] = f; });

  $('stats').textContent = manifest.claims.length + ' claims admitted · ' +
    manifest.extraction_metadata.rejected_count + ' rejected · ' +
    (d.chunk_errors.length ? d.chunk_errors.length + ' chunk errors · ' : '') + d.elapsed_s + 's';

  const errBlock = d.chunk_errors.length
    ? '<details class="panel reject"><summary><b>' + d.chunk_errors.length +
      ' chunk(s) failed to annotate</b> — the rest were still processed</summary><pre>' +
      esc(d.chunk_errors.map(e => e.chunk_id + ': ' + e.error).join('\n')) + '</pre></details>'
    : '';

  const rows = manifest.claims.map(c => {
    const f = fieldsById[c.claim_id] || {};
    const isOther = f.metric === 'Other';
    const flags = f.nonblocking_validation_errors || [];
    const rowClasses = ['row-' + (isOther ? 'other' : 'ok'), flags.length ? 'row-flagged' : ''].join(' ').trim();
    const metricLabel = isOther ? ('Other: ' + esc(f.metric_label || '?')) : esc(f.metric || '');
    return '<tr class="' + rowClasses + '">' +
      '<td class="metric">' + metricLabel + '</td>' +
      '<td>' + esc(c.value ?? '') + (c.unit ? ' ' + esc(c.unit) : '') + '</td>' +
      '<td><span class="pill ' + esc(c.epistemic_class || '') + '">' + esc(c.epistemic_class || '') + '</span></td>' +
      '<td>' + esc(f.basis || '') + '</td>' +
      '<td>' + esc(f.measurement || '') + '</td>' +
      '<td>' + esc(c.period || '') + '</td>' +
      '<td title="' + esc(c.statement || '') + '">' + esc((c.statement || '').slice(0, 90)) +
        ((c.statement || '').length > 90 ? '…' : '') + '</td>' +
      '<td>' + (flags.length
        ? '<span class="flagmark" title="' + esc(flags.join(' | ')) + '">⚑</span>' : '') + '</td>' +
      '</tr>';
  }).join('');

  $('claims').innerHTML = errBlock +
    '<div class="hint">Metric in purple = landed on the "Other" escape hatch (PAN-117), not force-fit onto a ' +
    'lookalike enum value. ⚑ = derivation arithmetic disagrees with the stored value (PAN-125) — flagged, still admitted.</div>' +
    '<table class="claims"><thead><tr><th>Metric</th><th>Value</th><th>Epistemic</th><th>Basis</th>' +
    '<th>Measurement</th><th>Period</th><th>Statement</th><th></th></tr></thead><tbody>' + rows + '</tbody></table>';

  $('graph').innerHTML =
    '<div class="hint">' + overlap.overlap_groups + ' identity group(s) with 2+ claims — ' +
    overlap.corroboration_groups + ' corroborating, ' + overlap.contradiction_candidate_groups +
    ' contradiction candidate(s). ' + overlap.claims_skipped_unresolvable +
    ' claim(s) skipped as unresolvable (missing entity/metric/period). This module never adjudicates — ' +
    'it only reports what agrees and what does not.</div>' +
    (overlap.findings.length ? overlap.findings.map(f => {
      const cls = f.status === 'CORROBORATION' ? 'corro' : 'contra';
      const idLine = Object.entries(f.identity).filter(([, v]) => v && v !== 'unspecified' && v !== 'none')
        .map(([k, v]) => k + '=' + v).join(', ');
      const claimLines = f.claims.map(c =>
        '<div>' + (c.shows_working ? '<b>shows working</b> · ' : '') +
        '<span class="pill ' + esc(c.epistemic_class || '') + '">' + esc(c.epistemic_class || '') + '</span> ' +
        esc(c.value ?? '') + ' — <span title="' + esc(c.statement || '') + '">' +
        esc((c.statement || '').slice(0, 60)) + '…</span> <small>(' + esc(c.source_id || '') + ')</small></div>'
      ).join('');
      return '<div class="cluster ' + cls + '"><div class="id">' + esc(f.status) + ' — ' + esc(idLine) +
        '</div>' + claimLines + '</div>';
    }).join('') : '<div class="drop">No overlapping identities found — every claim in this document stands alone.</div>');

  $('go3').disabled = false;
  $('tabArchetype').disabled = false;
}

async function runArchetype() {
  if (!SID) return;
  $('go3').disabled = true;
  $('archetype').innerHTML = '<div class="drop">Scanning admitted claims for family keywords…</div>';
  try {
    const res = await fetch('/archetype?sid=' + encodeURIComponent(SID), { method: 'POST' });
    renderArchetype(await res.json());
  } catch (e) {
    $('archetype').innerHTML = '<div class="panel error"><h2>Request failed</h2><pre>' + esc(String(e)) + '</pre></div>';
  } finally {
    $('go3').disabled = false;
  }
}

function renderArchetype(d) {
  if (d.status !== 'ok') {
    $('archetype').innerHTML = '<div class="panel error"><h2>Archetype heuristic failed</h2><pre>' +
      esc(d.error || JSON.stringify(d, null, 2)) + '</pre></div>';
    return;
  }
  const families = Object.entries(d.family_hits).map(([name, hits]) =>
    '<div class="archrow"><span>' + esc(name) + '</span><span>' +
    (hits.length ? hits.length + ' hit(s): ' + hits.map(esc).join(', ') : 'no hits') + '</span></div>'
  ).join('');
  $('archetype').innerHTML =
    '<div class="archnote">' + esc(d.note) + '</div>' +
    '<div class="archrow"><span>Total admitted claims</span><span>' + d.total_claims + '</span></div>' +
    '<div class="archrow"><span>On the "Other" escape hatch</span><span>' + d.other_count +
      ' (' + (d.other_ratio === null ? '—' : Math.round(d.other_ratio * 100) + '%') + ')</span></div>' +
    families;
}

// Map a chunk back onto the rendered original. Locator first (it is the
// parser's own claim about where the text is, so this also tests the
// locator); text search only as a fallback for formats whose locator is a
// section label rather than a coordinate.
function focusOriginal(c) {
  const frame = $('frame');
  if (KIND === 'pdf') {
    const page = c.page_or_slide_number || parseInt((c.locator.match(/p(\d+)/) || [])[1], 10);
    if (page) frame.src = '/original/' + SID + '#page=' + truePage(page);
    return;
  }
  const win = frame.contentWindow;
  if (!win) return;
  const cell = c.locator.match(/::(.+)!(\d+):(\d+)$/);   // book.xlsx::Sheet!12:20
  if (cell && win.panta_focus) {
    if (win.panta_focus(cell[1], +cell[2], +cell[3])) return;
  }
  if (c.page_or_slide_number && win.panta_focus) {
    if (win.panta_focus(null, c.page_or_slide_number, c.page_or_slide_number)) return;
  }
  // Fallback: find the chunk's first real line of text in the rendered doc.
  const line = c.body.split('\n')
    .map(s => s.trim())
    .filter(s => s.length > 12 && !/^(Workbook|Sheet|Source):/.test(s))[0];
  if (line && win.find) {
    win.getSelection().removeAllRanges();
    win.find(line.slice(0, 60), false, false, true);
  }
}

// The PDF-mode badge starts as "checking" because torch is slow to import
// here; poll until it settles rather than making the user reload.
(function pollMode() {
  const badge = document.querySelector('.badge');
  if (!badge.className.includes('pending')) return;
  fetch('/mode').then(r => r.json()).then(m => {
    badge.className = 'badge ' + m.level;
    badge.textContent = 'PDF: ' + m.label;
    badge.title = m.detail;
    [['docling', m.docling_available, 'Docling classic'],
     ['paddle', m.paddle_available, 'PaddleOCR-VL']].forEach(([k, ok, label]) => {
      if (ok === false) {
        const opt = document.querySelector('#engine option[value="' + k + '"]');
        if (opt) { opt.disabled = true; opt.textContent = label + ' (not installed)'; }
      }
    });
    if (m.level === 'pending') setTimeout(pollMode, 2000);
  }).catch(() => setTimeout(pollMode, 4000));
})();

// Draggable split.
(function () {
  let dragging = false;
  $('divider').addEventListener('mousedown', e => { dragging = true; e.preventDefault(); });
  window.addEventListener('mouseup', () => { dragging = false; $('frame').style.pointerEvents = ''; });
  window.addEventListener('mousemove', e => {
    if (!dragging) return;
    $('frame').style.pointerEvents = 'none';   // don't lose the drag inside the iframe
    const pct = Math.min(85, Math.max(15, (e.clientX / window.innerWidth) * 100));
    $('left').style.flex = '0 0 ' + pct + '%';
    $('right').style.flex = '1 1 0';
  });
})();
</script>
</body></html>
"""


def index() -> str:
    mode = _pdf_mode()
    return (INDEX_HTML
            .replace("__MODE_LEVEL__", mode["level"])
            .replace("__MODE_LABEL__", html.escape(mode["label"]))
            .replace("__MODE_DETAIL__", html.escape(mode["detail"], quote=True))
            .replace("__ACCEPT__", ",".join(SUPPORTED_EXTENSIONS))
            .replace("__SUPPORTED__", html.escape(", ".join(SUPPORTED_EXTENSIONS)))
            .replace("__UNSUPPORTED__", html.escape(", ".join(DECLARED_UNSUPPORTED))))


def main(port: int = 4192) -> None:
    host = os.environ.get("PE_OS_UI_HOST", "127.0.0.1")

    # Binding beyond loopback puts an unauthenticated file-upload form on a
    # network. On a laptop this tool is reachable only by its user; on a
    # server it would accept any document from anyone who finds the port,
    # and hand back the extracted text. Refuse rather than default to
    # something the operator did not ask for.
    if host != "127.0.0.1" and not AUTH_TOKEN:
        raise SystemExit(
            f"refusing to bind {host} without auth.\n"
            "Set PE_OS_UI_TOKEN to a long random secret, e.g.\n"
            "  export PE_OS_UI_TOKEN=$(python3 -c 'import secrets;print(secrets.token_urlsafe(32))')\n"
            "and put TLS in front of it (see deploy/README.md) -- uploads and\n"
            "extracted text travel in cleartext otherwise."
        )

    server = ThreadingHTTPServer((host, port), Handler)
    threading.Thread(target=_resolve_pdf_mode, daemon=True).start()
    print("resolving PDF model availability in the background...")
    shown = "127.0.0.1" if host in ("127.0.0.1", "0.0.0.0") else host
    print(f"extraction test UI → http://{shown}:{port}   (ctrl-c to stop)")
    if AUTH_TOKEN:
        print("auth: ON (send ?token=... once; it is stored in a cookie)")
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nstopped")


if __name__ == "__main__":
    main(int(sys.argv[1]) if len(sys.argv) > 1 else 4192)
