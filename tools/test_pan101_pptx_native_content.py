"""PAN-101 -- PPTX native chart-data, table, and speaker-notes extraction.

parse_pptx previously read only DrawingML text nodes (a:t) via hand-rolled
XML -- zero table structure, zero chart data, zero speaker notes. Real PE
decks carry native OOXML chart parts (not flat pictures) with exact
category/series/value numbers cached in the chart XML; this is a
deterministic, zero-OCR way to get real numbers out of a chart. Verifies
against a synthetic deck built the same way python-pptx itself would
produce a real one, not a hand-rolled minimal package.
"""

import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

from tools.extract_v2 import parse_pptx  # noqa: E402


class PAN101PptxNativeContentTests(unittest.TestCase):
    def setUp(self):
        self.temporary = tempfile.TemporaryDirectory()
        self.path = Path(self.temporary.name) / "deck.pptx"

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Revenue Bridge"

        chart_data = CategoryChartData()
        chart_data.categories = ["FY2023A", "FY2024A", "FY2025A"]
        chart_data.add_series("Revenue", (56.8, 65.4, 74.0))
        presentation_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(6), Inches(3), chart_data
        )
        self.chart_shape_name = presentation_chart.name

        table_shape = slide.shapes.add_table(3, 2, Inches(1), Inches(5), Inches(4), Inches(1.5))
        table = table_shape.table
        table.cell(0, 0).text, table.cell(0, 1).text = "Metric", "Value"
        table.cell(1, 0).text, table.cell(1, 1).text = "Revenue", "74.0"
        table.cell(2, 0).text, table.cell(2, 1).text = "EBITDA", "10.2"

        slide.notes_slide.notes_text_frame.text = "Confirm EBITDA with QoE before IC."

        presentation.save(str(self.path))

    def tearDown(self):
        self.temporary.cleanup()

    def _slide1_body(self) -> str:
        chunks = parse_pptx(self.path)
        self.assertEqual(len(chunks), 1)
        self.assertEqual(chunks[0].page_or_slide_number, 1)
        return chunks[0].body

    def test_native_chart_data_extracted_with_real_numbers(self):
        body = self._slide1_body()
        self.assertIn("FY2023A=56.8", body)
        self.assertIn("FY2024A=65.4", body)
        self.assertIn("FY2025A=74.0", body)
        self.assertIn("Revenue", body)

    def test_native_table_extracted_as_structured_grid(self):
        body = self._slide1_body()
        self.assertIn("Revenue | 74.0", body)
        self.assertIn("EBITDA | 10.2", body)

    def test_speaker_notes_are_captured(self):
        body = self._slide1_body()
        self.assertIn("Confirm EBITDA with QoE before IC.", body)

    def test_slide_title_text_still_captured(self):
        body = self._slide1_body()
        self.assertIn("Revenue Bridge", body)

    def test_flat_image_chart_is_a_declared_gap_not_a_silent_drop(self):
        import io

        from PIL import Image

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        png_bytes = io.BytesIO()
        Image.new("RGB", (4, 4), color="white").save(png_bytes, format="PNG")
        png_bytes.seek(0)
        picture = slide.shapes.add_picture(png_bytes, Inches(1), Inches(1), Inches(2), Inches(2))

        with_picture_path = Path(self.temporary.name) / "deck_with_picture.pptx"
        presentation.save(str(with_picture_path))

        chunks = parse_pptx(with_picture_path)
        self.assertEqual(len(chunks), 1)
        self.assertIn(picture.name, chunks[0].body)
        self.assertIn("IMAGE_NOT_EXTRACTED", chunks[0].body)

    def test_chartex_waterfall_chart_is_a_declared_gap_not_a_silent_drop(self):
        """PowerPoint's native waterfall/funnel/histogram chart family (the
        standard chart type for an EBITDA bridge) uses a 'chart-ex' graphicData
        namespace python-pptx's has_chart never matches -- confirmed by direct
        source-code reading of pptx/shapes/graphfrm.py (PAN-101 research).
        python-pptx cannot create a real chart-ex chart via its API (no
        creation support), so this simulates one the honest way: build a real
        native chart, then swap only its graphicData uri attribute to the
        real, confirmed chart-ex namespace -- everything else (rels, cached
        series data, package structure) stays exactly as a real chart
        produces, so this exercises the real detection path, not a fabricated
        shortcut.
        """
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        chart_data = CategoryChartData()
        chart_data.categories = ["FY2023A"]
        chart_data.add_series("EBITDA bridge", (1.0,))
        graphic_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(3), Inches(3), chart_data
        )
        graphic_data_el = graphic_frame._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}graphicData"
        )
        graphic_data_el.set("uri", "http://schemas.microsoft.com/office/drawing/2014/chartex")
        self.assertFalse(graphic_frame.has_chart)  # confirms the simulated condition is real

        path = Path(self.temporary.name) / "waterfall.pptx"
        presentation.save(str(path))

        chunks = parse_pptx(path)
        self.assertEqual(len(chunks), 1)
        self.assertIn(graphic_frame.name, chunks[0].body)
        self.assertIn("UNSUPPORTED_GRAPHIC_FRAME", chunks[0].body)
        self.assertIn("waterfall", chunks[0].body.lower())

    def test_invalid_pptx_is_a_declared_rejection_not_a_crash(self):
        from tools.extract_v2 import UnsupportedSourceError

        bad = Path(self.temporary.name) / "not_a_deck.pptx"
        bad.write_bytes(b"not a real zip")
        with self.assertRaises(UnsupportedSourceError):
            parse_pptx(bad)


if __name__ == "__main__":
    unittest.main()
