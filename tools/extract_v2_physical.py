#!/usr/bin/env python3
"""
extract_v2_physical.py — E3 extraction pipeline for PANTA.

Architecture:
  L1  Document parser    → deterministic chunks with stable locators
  L2  LLM annotator      → one call per chunk, tool_use schema-constrained
  L3  Claim validator    → deterministic normalization + stable_id
  L4  Graph assembler    → deterministic dedup + conflict detection

Claim schema aligned to PANTA_Keystone_Canonical_Investment_Case_v1.1.json
(CAP-003). Does NOT touch extract.py.

Manifest modes prevent temporal leakage between knowledge snapshots:
  K-PRE   CIM, DR, IA, QL, QoE  — state of knowledge before IC memo
  K-IC    K-PRE + IC memo        — as-of IC gate (2026-03-10)
  K-LIVE  K-IC + board packs     — post-close, event-by-event

Usage:
  python3 tools/extract_v2_physical.py --manifest K-IC --deal keystone --dry-run

  export OPENROUTER_API_KEY=sk-or-...
  export PEOS_LLM_PROVIDER=openrouter  # defaults to z-ai/glm-5.2
  python3 tools/extract_v2_physical.py --manifest K-PRE --deal keystone \\
      --output pipeline_out/e3/k_pre

  python3 tools/extract_v2_physical.py --source vault/inbox/keystone_ic_memo.md \\
      --deal keystone --output pipeline_out/e3/single

  python3 tools/extract_v2_physical.py --manifest K-IC --deal keystone \\
      --compare pipeline_out/keystone_full_story/graph.json
"""
from __future__ import annotations

import argparse
import csv
import hashlib
import json
import mailbox
import math
import os
import re
import sys
import copy
import textwrap
from functools import lru_cache
import time
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from datetime import date, datetime
from email import policy as email_policy
from email.parser import BytesParser
from email.utils import parsedate_to_datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from tools.llm_provider import (  # noqa: E402
    anthropic_client_kwargs,
    configured_api_key,
    configured_model,
    missing_key_message,
    forces_tool_choice,
    openrouter_extra_body,
    thinking_parameter,
)
from tools.archetype_pack import (  # noqa: E402
    DEFAULT_ARCHETYPE, evidence_state_vocabulary, extraction_vocabulary,
    load_pack, missing_required_identity, workstream_ids,
)
from tools.derivation_verifier import verify_derivation  # noqa: E402
from tools.object_identity import claim_id as canonical_claim_id  # noqa: E402
from tools.source_envelope import extractor_source_record  # noqa: E402
from tools.source_capabilities import (  # noqa: E402
    CAPABILITY_SCHEMA,
    capability_failure,
    resolve_source_capability,
)
from tools.source_catalog import (  # noqa: E402
    load_source_catalog,
    source_record_from_catalog,
    unmapped_source_types,
)

VAULT_INBOX = ROOT / "vault" / "inbox"
MODEL = configured_model("claude-haiku-4-5-20251001")
# A merged markdown chunk (parse_markdown folds adjacent short sections
# together, see its docstring) can legitimately need up to CLAIM_TOOL's own
# cap of 20 claims in one tool call. Measured directly: a 202-word, 16-claim
# chunk hit stop_reason=max_tokens at 4096 and returned zero claims -- not
# occasionally, 3/3 attempts -- because a truncated tool_use block parses to
# no claims at all, silently, with no exception to catch. 8192 leaves real
# headroom under that same load.
MAX_TOKENS = int(os.environ.get("PEOS_EXTRACT_V2_MAX_TOKENS", "8192"))
MAX_PROVIDER_RETRIES = int(os.environ.get("PEOS_LLM_CHUNK_RETRIES", "3"))

# ─────────────────────────────────────────────────────────────────────────────
# Source registry — maps vault/inbox filenames to canonical SRC-xxx IDs.
# Aligned with the sources dict in the benchmark CIC v1.1.
# known_at = when the Firm first received this document.
# ─────────────────────────────────────────────────────────────────────────────

SOURCE_REGISTRY: dict[str, dict] = {
    # filename (without extension or with) → canonical source record
    "keystone_seller_cim": {
        "source_id": "SRC-CIM",
        "name": "Project Keystone CIM",
        "party": "Alderstone management and Hawthorne Capital Markets",
        "doc_type": "CIM",
        "effective_date": "2025-10-27",
        "known_at": "2025-10-27",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_data_room_extract": {
        "source_id": "SRC-DR",
        "name": "Keystone data room extract",
        "party": "Alderstone management / VDR export",
        "doc_type": "Data Room",
        "effective_date": "2025-11-15",
        "known_at": "2025-11-15",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_firm_initial_assessment": {
        "source_id": "SRC-IA",
        "name": "Firm initial assessment",
        "party": "The Firm investment team",
        "doc_type": "Internal",
        "effective_date": "2025-12-01",
        "known_at": "2025-12-01",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_question_list": {
        "source_id": "SRC-QL",
        "name": "Keystone question list",
        "party": "The Firm diligence team",
        "doc_type": "Internal",
        "effective_date": "2026-01-10",
        "known_at": "2026-01-10",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_qoe_report": {
        "source_id": "SRC-QOE",
        "name": "Quality of Earnings report",
        "party": "Independent financial diligence provider",
        "doc_type": "QoE Report",
        "effective_date": "2026-02-20",
        "known_at": "2026-02-20",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_firm_model_summary": {
        "source_id": "SRC-MODEL-SUM",
        "name": "Firm model summary",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_lbo_model_working": {
        "source_id": "SRC-MODEL",
        "name": "Keystone LBO model",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone-model-part1": {
        "source_id": "SRC-MODEL",
        "name": "Keystone LBO model (part 1)",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone-model-part2": {
        "source_id": "SRC-MODEL",
        "name": "Keystone LBO model (part 2)",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone-model-part3": {
        "source_id": "SRC-MODEL",
        "name": "Keystone LBO model (part 3)",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone-model-part4": {
        "source_id": "SRC-MODEL",
        "name": "Keystone LBO model (part 4)",
        "party": "The Firm investment team",
        "doc_type": "LBO Model",
        "effective_date": "2026-03-05",
        "known_at": "2026-03-05",
        "manifest": ["K-PRE", "K-IC", "K-LIVE"],
    },
    "keystone_ic_memo": {
        "source_id": "SRC-IC",
        "name": "IC memo 2026-03-10",
        "party": "The Firm investment team / Investment Committee",
        "doc_type": "IC Memo",
        "effective_date": "2026-03-10",
        "known_at": "2026-03-10",
        "manifest": ["K-IC", "K-LIVE"],  # not in K-PRE
    },
    "keystone_monitoring_boardpack1_dec2026": {
        "source_id": "SRC-BP1",
        "name": "Board pack Dec 2026",
        "party": "Alderstone management / board package",
        "doc_type": "Board Pack",
        "effective_date": "2026-12-31",
        "known_at": "2026-12-31",
        "manifest": ["K-LIVE"],
    },
    "keystone_monitoring_boardpack2_mar2027": {
        "source_id": "SRC-BP2",
        "name": "Board pack Mar 2027",
        "party": "Alderstone management / board package",
        "doc_type": "Board Pack",
        "effective_date": "2027-03-31",
        "known_at": "2027-03-31",
        "manifest": ["K-LIVE"],
    },
    "keystone_monitoring_junecompliance2027": {
        "source_id": "SRC-JUN27",
        "name": "June 2027 compliance certificate",
        "party": "Alderstone management / board package",
        "doc_type": "Board Pack",
        "effective_date": "2027-06-30",
        "known_at": "2027-06-30",
        "manifest": ["K-LIVE"],
    },
    "keystone_monitoring_augustamendment2027": {
        "source_id": "SRC-AUG27",
        "name": "August 2027 amendment",
        "party": "Alderstone management / board package",
        "doc_type": "Amendment",
        "effective_date": "2027-08-15",
        "known_at": "2027-08-15",
        "manifest": ["K-LIVE"],
    },
    "keystone_monitoring_recovery_exit2031": {
        "source_id": "SRC-EXIT31",
        "name": "Exit 2031 package",
        "party": "Alderstone management / board package",
        "doc_type": "Board Pack",
        "effective_date": "2031-06-01",
        "known_at": "2031-06-01",
        "manifest": ["K-LIVE"],
    },
}

# Sources per manifest — derived from registry
MANIFEST_SOURCES: dict[str, list[str]] = {
    "K-PRE": [k for k, v in SOURCE_REGISTRY.items() if "K-PRE" in v["manifest"]],
    "K-IC": [k for k, v in SOURCE_REGISTRY.items() if "K-IC" in v["manifest"]],
    "K-LIVE": [k for k, v in SOURCE_REGISTRY.items() if "K-LIVE" in v["manifest"]],
    "ALL": list(SOURCE_REGISTRY.keys()),
}

def _source_record(path: Path, catalog: dict[str, dict] | None = None) -> dict:
    """Return registry metadata first, then deal-catalog metadata or fallback."""
    stem = path.stem
    # Try exact match first, then strip suffixes like _part1, _extract
    record = SOURCE_REGISTRY.get(stem)
    if not record and catalog is not None:
        record = source_record_from_catalog(path, catalog)
    if not record:
        for key in SOURCE_REGISTRY:
            if stem.startswith(key) or key.startswith(stem):
                record = SOURCE_REGISTRY[key]
                break
    if not record:
        record = {
            "source_id": f"SRC-{stem.upper()[:12]}",
            "name": stem,
            "party": "unknown",
            "doc_type": "Other",
            "effective_date": "",
            "known_at": "",
            "manifest": ["ALL"],
        }
    return record


# ─────────────────────────────────────────────────────────────────────────────
# Schema definitions (bounded enums)
# ─────────────────────────────────────────────────────────────────────────────

METRIC_ENUM: list[str] = [
    "Revenue", "Recurring Revenue", "Revenue Growth",
    "Gross Profit", "Gross Margin",
    "EBITDA", "EBITDA Margin", "EBITDA Add-back", "EBITDA Adjustment",
    "EBIT", "Net Income", "Free Cash Flow", "Operating Cash Flow",
    "Capex", "Working Capital", "DSO", "DPO", "Inventory Days",
    "Earnings Quality Risk", "Revenue Quality", "Adjustment Supportability",
    "Customer Concentration", "Customer Count", "Active Billing Accounts",
    "Customer Retention", "Customer Churn", "Contract Terms", "Market Position", "Market Size",
    "Enterprise Value", "Equity Value", "Entry Multiple", "Exit Multiple", "Exit EV",
    "Net Debt", "Gross Debt", "Leverage", "Interest Coverage",
    "Sponsor Equity", "Seller Rollover", "First-Lien Debt",
    "Revolver Capacity", "DDTL Availability", "Covenant EBITDA",
    "Covenant Threshold", "Covenant Headroom",
    # Credit-agreement-defined variants, distinct from the general "Leverage"
    # entry the same way "Covenant EBITDA" is distinct from plain "EBITDA":
    # a defined term with its own add-back/measurement conventions, not the
    # same fact wearing a longer name.
    "Total Net Leverage Ratio", "Minimum Liquidity",
    "MOIC", "IRR", "Exit Horizon", "Supported Price",
    "Net Working Capital", "Net Working Capital Target", "Net Working Capital Adjustment",
    "Headcount", "Team Tenure", "Acquisition Count",
    "Systems Integration Risk", "Integration Risk", "Operational Risk",
    "Key Person Risk", "Regulatory Risk", "Competition Risk",
    "IC Conditions", "IC Vote", "Decision Coherence",
    "Customer Contract Terms",
    # Escape hatch. This enum is LBO/buyout-shaped, and a closed enum with no
    # exit forces a wrong pick on anything outside that shape: measured on a
    # venture term-sheet fragment, a EUR 6.0m primary round became "Sponsor
    # Equity" (a buyout sponsor's own cheque — a different thing), a EUR 24.0m
    # pre-money became "Enterprise Value", and "eighteen to twenty-four months
    # of runway" became "Exit Horizon" = 21, a number stated nowhere in the
    # source. Naming the gap is strictly better than filling it wrongly:
    # metric "Other" carries the real concept in `metric_label`, and
    # normalize_metric() maps it to "" — so the claim is *unresolvable* by
    # design, stays visible in the ledger, and is never silently matched
    # against a claim it only superficially resembles.
    "Other",
]
_seen: set[str] = set()
METRIC_ENUM = [m for m in METRIC_ENUM if not (m in _seen or _seen.add(m))]  # type: ignore

UNIT_ENUM: list[str | None] = [
    None, "", "$m", "£m", "€m", "$m/year", "$m/quarter",
    "%", "x", "bps", "days", "headcount", "turns", "$", "£", "€",
]

PERIOD_MAP: dict[str, str] = {
    "FY2025A": "2025-12-31", "FY2025": "2025-12-31",
    "FY2026E": "2026-12-31", "FY2026": "2026-12-31",
    "FY2027E": "2027-12-31", "FY2027": "2027-12-31",
    "FY2028E": "2028-12-31", "FY2028": "2028-12-31",
    "FY2029E": "2029-12-31", "FY2029": "2029-12-31",
    "FY2030E": "2030-12-31", "FY2030": "2030-12-31",
    "FY2031E": "2031-12-31", "FY2031": "2031-12-31",
    "OPENING": "2026-03-31",
    "LTM": "LTM",
    "Q1 2026": "2026-06-30", "Q2 2026": "2026-09-30",
    "Q3 2026": "2026-12-31", "Q4 2026": "2027-03-31",
    "Q1 2027": "2027-06-30", "Q2 2027": "2027-09-30",
}
# Period is now free-text — no enum constraint.
# PERIOD_MAP is kept for period_iso normalization only (not for constraining the field).

EPISTEMIC_CLASS_ENUM = ["asserted", "observed", "derived", "attested", "institutional_act"]
DIRECTION_ENUM = ["supports", "contradicts", "context"]

# What kind of assertion a claim is. The distinction that matters is the last
# one: a sentence with no number earns a claim only if someone could later
# confirm or refute it. "No material litigation as of 30 June" is checkable;
# "a scaled regional platform" is not, and admitting it puts a seller's adjective
# into the case as though it were evidence.
#
# This criterion used to be implicit, and that is exactly where extraction
# wobbled: across runs of one paragraph the quantified claims were identical
# every time, while a valueless MarketPosition claim from the phrase "scaled
# regional platform" appeared in some runs and not others.
CLAIM_KIND_ENUM = [
    "QUANTITATIVE",     # carries a number or a measured quantity
    "DEFINITION",       # states how something is defined or calculated
    "CONDITION",        # states a requirement, contingency or covenant
    "ATTRIBUTION",      # states who said, did or decided something
    "NEGATIVE",         # explicitly asserts an absence — checkable, so admissible
    # A checkable proposition about the business that carries no number: a
    # capability assessment, a deployment or proof state, a buyer's stance.
    # Added because CHARACTERISATION had become the residual bucket for
    # everything non-numeric, and CHARACTERISATION is deleted: measured over
    # five documents of a venture corpus, 42 of 103 claims died there, and
    # sampling them found proof states, buyer urgency and milestone
    # requirements rather than adjectives. Maps to the dictionary's
    # `claim_kind=qualitative` (section 3.2).
    "QUALITATIVE",
    "CHARACTERISATION", # a descriptor with no checkable content — DO NOT EMIT
]

# How to read the number. "More than 600 accounts" recorded as 600 EXACT makes a
# later, entirely consistent "640" look like a contradiction. The bound is not
# part of identity — it describes the value, not the quantity being measured —
# but comparison must respect it.
BOUND_ENUM = ["EXACT", "AT_LEAST", "AT_MOST", "APPROXIMATE", "RANGE", "NONE"]
ARCHETYPE_PACK = load_pack()
# "Use workstream/concept families to select schemas. Category never creates identity or contradiction."
# Topic is extraction routing metadata only; identity and conflict logic deliberately ignore it.
TOPIC_ENUM = workstream_ids(ARCHETYPE_PACK) + ["OTHER"]


def _topic_description(pack: dict[str, Any]) -> str:
    """Give the model the pack's governing questions beside each route ID."""
    workstreams = pack["workstreams"]
    routes = "\n".join(
        f"- {workstream_id}: {workstreams[workstream_id]['governing_question']}"
        for workstream_id in workstream_ids(pack)
    )
    return (
        "Choose the workstream whose governing question best routes this claim; "
        "do not infer identity or contradiction from topic. Use OTHER only when no "
        "workstream applies.\nCanonical workstreams:\n"
        f"{routes}"
    )


TOPIC_DESCRIPTION = _topic_description(ARCHETYPE_PACK)

# ── Identity dimensions (bounded) ─────────────────────────────────────────────
# `perimeter` stays prose because a human reads it. These carry the machine
# identity instead: an exact-match table cannot normalize unbounded language, so
# the dimensions that decide whether two claims are comparable are asked for
# separately and constrained here. See tools/object_identity.py.
BASIS_ENUM = [
    "SellerView", "QoEView", "FirmView", "CovenantView", "ReportedView", "unspecified",
]
SCOPE_ENUM = ["consolidated", "standalone", "customer", "segment", "unspecified"]
SCENARIO_ENUM = ["base", "management", "seller", "upside", "downside", "unspecified"]

# Canonical period grammar. Anything outside it is "none" — an honest gap beats
# an extraction timestamp masquerading as a period (the failure that put
# "as of <ingest date>" on 227 vault claims).
PERIOD_CANONICAL_PATTERN = (
    r"^(FY\d{4}[AE]?|LTM(-\d{4}-\d{2}-\d{2})?|ExitLTM|Opening|EntryToExit"
    r"|CrossPeriod|Recurring|\d{4}-Q[1-4]|\d{4}-\d{2}-\d{2}|none)$"
)

# Known definition IDs — if LLM flags a claim as definition-linked, map to DEF-xxx.
# These align with the definitions[] array in the canonical benchmark.
DEFINITION_ENUM: list[str | None] = [
    None,
    "DEF-FIRM-EBITDA", "DEF-QOE-EBITDA", "DEF-COV-EBITDA",
    "DEF-SELLER-EBITDA", "DEF-NWC", "DEF-NWC-TARGET",
    "DEF-RECURRING-REVENUE", "DEF-CONCENTRATION",
    "DEF-MOIC", "DEF-IRR", "DEF-EV",
]

# ─────────────────────────────────────────────────────────────────────────────
# Tool schema (tool_use — enforced at LLM decoding level)
# ─────────────────────────────────────────────────────────────────────────────

CLAIM_TOOL = {
    "name": "emit_claims",
    "description": (
        "Emit 0-20 financial claims from this document fragment. "
        "Extract only what is explicitly stated. Empty list if no financial claims. "
        "epistemic_class: asserted=seller/mgmt claim, observed=third-party real-time, "
        "attested=formally certified (QoE conclusion, IC decision), derived=YOU computed it."
    ),
    "input_schema": {
        "type": "object",
        "required": ["claims"],
        "properties": {
            "claims": {
                "type": "array",
                # A 250-word Excel chunk can contain more than four distinct
                # financial rows. The old cap deterministically dropped later
                # rows such as PAN-37's DSO assumption.
                "maxItems": 20,
                "items": {
                    "type": "object",
                    "required": [
                        "metric", "value", "unit", "period", "perimeter",
                        "entity", "period_canonical", "scope", "measurement", "basis", "scenario",
                        "claim_kind", "bound",
                        "epistemic_class", "direction", "topic",
                        "statement", "locator_hint",
                    ],
                    "additionalProperties": False,
                    "properties": {
                        "metric": {
                            "type": "string",
                            "enum": METRIC_ENUM,
                            "description": (
                                "Pick the entry that genuinely names this quantity. "
                                "If none does, pick 'Other' and name the real concept "
                                "in metric_label — never the closest-sounding entry. "
                                "A near-miss label is worse than an honest gap: it "
                                "silently merges this claim with unrelated ones."
                            ),
                        },
                        "metric_label": {
                            "type": ["string", "null"],
                            "maxLength": 60,
                            "description": (
                                "REQUIRED when metric is 'Other': the concept's real "
                                "name, in the source's own words (e.g. 'Pre-money "
                                "valuation', 'Post-money valuation', 'Runway months', "
                                "'Paid site count'). Leave null for every other metric."
                            ),
                        },
                        "value": {
                            "type": ["number", "string", "null"],
                        },
                        "unit": {
                            "type": ["string", "null"],
                            "enum": UNIT_ENUM,
                        },
                        "period": {
                            "type": "string",
                            "minLength": 1,
                            "pattern": r".*\S.*",
                            "description": (
                                "The time period or vintage this claim refers to, in the document's own language. "
                                "Use the richest description available. Examples: "
                                "'FY2025A', 'FY2025A / FY2025E seller presentation', 'As of 2025-10-27', "
                                "'LTM Sep-25', '2020 to 2025-10-27', 'FY2026E–FY2030E forecast'. "
                                "Use 'cross-period' only if the claim is genuinely timeless."
                            ),
                        },
                        "perimeter": {
                            "type": "string",
                            "minLength": 1,
                            "pattern": r".*\S.*",
                            "description": (
                                "The precise economic scope — entity + metric definition + any adjustments. "
                                "Write the full descriptive string, not a shorthand. Examples: "
                                "'Alderstone consolidated revenue', "
                                "'Alderstone consolidated EBITDA under seller adjustment perimeter', "
                                "'Alderstone consolidated EBITDA under independent QoE adjustment perimeter', "
                                "'Alderstone consolidated EBITDA under Firm valuation, leverage and returns perimeter', "
                                "'Alderstone customer revenue measured by individual billing account', "
                                "'Alderstone accounts-receivable ledger'. "
                                "Use 'unknown' only if scope is truly unspecified."
                            ),
                        },
                        "entity": {
                            "type": "string",
                            "minLength": 1,
                            "description": (
                                "The bare proper name of the company or counterparty this claim "
                                "measures — nothing else. NOT a sentence, NOT the metric, NOT the "
                                "perimeter. 'Bare' means no surrounding sentence, not a shortened "
                                "name: write it exactly as the source spells it, full legal-form "
                                "words included ('Orion Group', 'Acme Holdings', 'Orion Services "
                                "Ltd') — 'Group', 'Holdings', 'Ltd' are part of the proper name "
                                "here, not boilerplate to strip. The same goes for a leading "
                                "'Project': the deal is called 'Project Cedar', so entity is "
                                "'Project Cedar', not 'Cedar'. Copy the name the document uses, "
                                "whole, in both directions — do not lengthen it either. "
                                "If the text says 'the Company' "
                                "or 'the Target', resolve it to the deal's proper name when the "
                                "document makes it unambiguous, otherwise write 'unspecified'.\n"
                                "A customer/billing-account row is about the CUSTOMER, not the "
                                "reporting company: on a schedule breaking the deal company's "
                                "revenue down by customer, entity is that row's customer or "
                                "ultimate parent (e.g. 'Acme Holdings'), never the deal company's "
                                "own name, even though the deal company is whose revenue is being "
                                "sliced. The deal company's name belongs in the perimeter/"
                                "measurement instead ('Meridian customer revenue — Acme Holdings "
                                "account'). Getting this backwards on a schedule with multiple "
                                "distinct customers collapses every row onto one identity."
                            ),
                        },
                        "period_canonical": {
                            "type": "string",
                            "pattern": PERIOD_CANONICAL_PATTERN,
                            "description": (
                                "The same period as the `period` field, reduced to one canonical "
                                "token: FY2025A, FY2025E, FY2024A, LTM, ExitLTM, Opening, "
                                "EntryToExit, CrossPeriod, Recurring, a single quarter as "
                                "YYYY-Q# (2026-Q2), an ISO date (2026-03-31), or exactly 'none'. "
                                "Write 'none' when the source states no period — never substitute "
                                "today's date, and never guess a fiscal year the document does "
                                "not give.\n"
                                "CrossPeriod and Recurring are NOT the same thing. CrossPeriod is "
                                "one static, time-invariant fact ('incorporated in Delaware'). "
                                "Recurring is a condition re-tested every time its trigger recurs "
                                "— a covenant checked every quarter end, a condition that applies "
                                "'at all times' or 'on each measurement date', or a condition "
                                "re-applied on every occurrence of an event ('immediately after "
                                "each acquisition'). A covenant is not timeless just because it "
                                "has no single stated date: it recurs, whether the trigger is a "
                                "calendar date or an event.\n"
                                "LTM needs its end date when the source states one: write "
                                "LTM-2026-06-30, not bare LTM. Two different LTM figures months "
                                "apart are two different periods, and bare LTM cannot tell them "
                                "apart. Use bare LTM only when the source truly never states which "
                                "date the trailing twelve months ends on."
                            ),
                        },
                        "scope": {
                            "type": "string",
                            "enum": SCOPE_ENUM,
                            "description": (
                                "Economic boundary: consolidated=whole group; standalone=one "
                                "entity; customer=one customer/account; segment=a division.\n"
                                "A figure explicitly framed as one division's or one segment's "
                                "contribution is segment, even when the same sentence also says "
                                "it sits 'within' or 'as part of' the consolidated result — that "
                                "phrase describes the relationship between two DIFFERENT claims "
                                "(the segment figure and the whole-group figure), it does not "
                                "make the segment figure itself a consolidated one. 'The North "
                                "America segment contributed $40.0m within the consolidated "
                                "group result' is scope=segment, value=40.0 — a separate claim "
                                "from the group's own consolidated total, not a restatement of it."
                            ),
                        },
                        "measurement": {
                            "type": "string",
                            "minLength": 1,
                            "description": (
                                "WHICH SLICE of the quantity this figure covers, in the source's "
                                "own words — a service line, department, customer account, "
                                "product, region, or cost category. Write exactly 'total' when "
                                "the figure is the whole, undivided quantity.\n"
                                "This is what separates a breakdown from a disagreement. Three "
                                "service lines reporting 30.3, 20.0 and 14.1 are components of "
                                "one revenue, not three conflicting claims about it — but only "
                                "if each names its slice. Leaving this blank on a component "
                                "makes it collide with the total and with its siblings.\n"
                                "Examples: 'total', 'EHS compliance service line', "
                                "'field inspection', 'Riverton account', 'engineering headcount'.\n"
                                "Name the slice as a NOUN PHRASE saying what kind of thing it "
                                "is, not as a bare identifier: 'RIV-001 billing account', not "
                                "'RIV-001'; 'segment contribution', not just the segment's name; "
                                "'ultimate-parent total' for a parent-level roll-up. A bare code "
                                "tells a reader nothing about what was measured.\n"
                                "When the metric is inherently ABOUT another metric — a "
                                "threshold, headroom, cap or covenant test FOR something — "
                                "this field must name that something, never 'total'. Use "
                                "generic 'Covenant Threshold' as the metric ONLY when no "
                                "more specific, credit-agreement-defined metric entry exists "
                                "(e.g. an FCCR minimum, which has no entry of its own) — "
                                "then this field carries the specificity: 'FCCR'. When a "
                                "specific defined-term entry DOES exist (Total Net Leverage "
                                "Ratio, Minimum Liquidity), use THAT as the metric instead of "
                                "the generic one, even inside a covenant/compliance context — "
                                "the specific entry is more identity-correct, not a fallback. "
                                "A compliance page with a 4.25x leverage cap, a 1.25x FCCR "
                                "minimum and a $3.0m liquidity floor holds a Total Net "
                                "Leverage Ratio claim, a Covenant Threshold/'FCCR' claim, and "
                                "a Minimum Liquidity claim — three different metrics, not "
                                "three Covenant Threshold claims distinguished only by this "
                                "field."
                            ),
                        },
                        "basis": {
                            "type": "string",
                            "enum": BASIS_ENUM,
                            "description": (
                                "Which party's adjustments the figure is stated under. This is "
                                "what makes seller EBITDA and QoE EBITDA two legitimate numbers "
                                "rather than a contradiction, so do not guess: "
                                "SellerView=seller/management adjusted; QoEView=independent "
                                "quality-of-earnings; FirmView=our own underwriting basis; "
                                "CovenantView=credit-agreement definition; ReportedView=statutory "
                                "or unadjusted; unspecified=the source does not say. "
                                "A plain figure pulled straight from the entity's own books or a "
                                "reporting schedule (a billing-account table, an AR ledger, a "
                                "revenue-by-customer breakdown) with no adjustment view mentioned "
                                "is ReportedView, not unspecified — reserve unspecified for when "
                                "even the unadjusted/statutory reading is unclear, not merely "
                                "because no basis word appears in the row. Likewise, a figure "
                                "explicitly attributed to management — 'management forecasts', "
                                "'management targets', a base/downside/upside case — is "
                                "SellerView even when it is a forward-looking forecast rather "
                                "than a historical EBITDA adjustment: WHO produced the figure is "
                                "what basis tracks, not only whether the source used the word "
                                "'adjusted'. unspecified is for when the fragment genuinely does "
                                "not attribute the figure to anyone."
                            ),
                        },
                        "scenario": {
                            "type": "string",
                            "enum": SCENARIO_ENUM,
                            "description": (
                                "Which case the figure belongs to. Base/upside/downside of the "
                                "same metric are NOT contradictions, so label them: "
                                "management=management's own forecast; seller=seller case; "
                                "base=our base case or an actual historical figure; "
                                "unspecified=the source does not distinguish cases."
                            ),
                        },
                        "claim_kind": {
                            "type": "string",
                            "enum": CLAIM_KIND_ENUM,
                            "description": (
                                "What kind of assertion this is. The test for a sentence with no "
                                "number is whether someone could later CONFIRM OR REFUTE it:\n"
                                "  QUANTITATIVE    carries a number or measured quantity\n"
                                "  DEFINITION      states how something is defined or calculated\n"
                                "  CONDITION       states a requirement, contingency or covenant\n"
                                "  ATTRIBUTION     states who said, did or decided something\n"
                                "  NEGATIVE        explicitly asserts an absence — checkable\n"
                                "  QUALITATIVE     a checkable statement about the business "
                                "that carries no number\n"
                                "  CHARACTERISATION a descriptor with no checkable content\n"
                                "QUALITATIVE vs CHARACTERISATION is the line between a claim "
                                "someone could go and check and a word someone chose. Ask what "
                                "would have to be true for it to be WRONG. If you can name the "
                                "evidence that would refute it, it is QUALITATIVE — a capability "
                                "assessment ('classification is not solved for quiet activity'), "
                                "a deployment or proof state ('installed and streaming, "
                                "validation still pending'; 'the prototype drove autonomously on "
                                "a test field but has no customer deployment'), a counterparty's "
                                "stance ('the customer treats the problem as real but not "
                                "urgent and will not redesign around it'), a scaling constraint. "
                                "None of those carry a number and all of them can be refuted by "
                                "the next site visit or reference call. CHARACTERISATION is for "
                                "what survives no such test: 'leading', 'scaled', 'attractive "
                                "momentum', 'strong position' — strip the adjective and there is "
                                "nothing left to check. When a sentence has substance AND "
                                "adjectives, keep the substance and label it QUALITATIVE; do not "
                                "discard a real finding because it was phrased warmly.\n"
                                "Label honestly; what happens to each kind is not your "
                                "concern. 'A scaled regional platform', 'low capital "
                                "expenditure' and 'strong market position' are "
                                "CHARACTERISATION — the seller's adjectives, with nothing to "
                                "check. 'No material litigation as of 30 June' is NEGATIVE: an "
                                "absence someone can verify. Do not reach for NEGATIVE when a "
                                "phrase is merely favourable. 'The dissenting IC member views "
                                "the integration risk as inadequately compensated' is "
                                "ATTRIBUTION even though it carries no number: the checkable "
                                "content is that this specific person holds this view, not "
                                "whether the view itself is correct. A named party's stated "
                                "position, objection or vote is ATTRIBUTION; CHARACTERISATION "
                                "is for un-attributed narrative color with no party to check it "
                                "against.\n"
                                "Attribution alone does not decide this — a seller's adjectives "
                                "are still adjectives when a named party says them. The test is "
                                "whether the VIEW ITSELF contains something confirmable or "
                                "refutable. 'Management describes the company as a leading, "
                                "scaled platform with attractive momentum' is CHARACTERISATION: "
                                "strip the attribution and nothing testable remains. 'The IC "
                                "believes current product truth is narrower than the deck: "
                                "partly detect, not full classification' is ATTRIBUTION: the "
                                "view carries a substantive proposition about what the product "
                                "does, which later evidence can confirm or refute, and the fact "
                                "that this committee holds it is itself checkable. A deal team's "
                                "or committee's reasoning about the target is almost always the "
                                "second kind — it is the most decision-critical content in the "
                                "document, and labelling it CHARACTERISATION removes it from the "
                                "graph entirely.\n"
                                "An instruction telling the reader HOW TO ANALYSE the data in "
                                "this fragment ('aggregate accounts only when the ultimate-"
                                "parent field is identical', 'similar names alone do not "
                                "establish common ownership') is not a claim about the company "
                                "at all, DEFINITION or otherwise — do not emit it as one. It "
                                "carries no fact someone could confirm or refute about the "
                                "business; it is a methodology note for the reader, and its "
                                "content only matters for how YOU apply it to the real claims "
                                "in this fragment. DEFINITION is for the source stating what a "
                                "financial TERM means or how a stated figure was itself "
                                "calculated ('Covenant EBITDA excludes X'), not for aggregation "
                                "or analysis instructions aimed at you."
                            ),
                        },
                        "bound": {
                            "type": "string",
                            "enum": BOUND_ENUM,
                            "description": (
                                "How to read the number.\n"
                                "  EXACT       the figure as stated\n"
                                "  AT_LEAST    'more than 600', 'over 72%', 'at least 4'\n"
                                "  AT_MOST     'up to 15%', 'no more than 3x', 'below 2%'\n"
                                "  APPROXIMATE 'around 11.4', 'circa', 'roughly'\n"
                                "  RANGE       '1%-2%' — put the lower figure in value\n"
                                "  NONE        the claim carries no number\n"
                                "Get this right or a later consistent figure reads as a conflict: "
                                "'more than 600' stored as EXACT 600 makes a subsequent 640 look "
                                "like a contradiction when the two agree."
                            ),
                        },
                        "epistemic_class": {
                            "type": "string",
                            "enum": EPISTEMIC_CLASS_ENUM,
                            "description": (
                                "asserted=seller/management; observed=third-party real-time; "
                                "attested=an independent third party's formal certification "
                                "(QoE conclusion, auditor opinion); institutional_act=a body's own "
                                "formal decision or action, not a third party certifying someone "
                                "else's numbers — an audit committee restating accounts, an IC "
                                "adopting an assumption, a lender's own credit-agreement covenant "
                                "term; derived=computed by you"
                            ),
                        },
                        "direction": {
                            "type": "string",
                            "enum": DIRECTION_ENUM,
                            "description": "supports=positive thesis signal; contradicts=negative signal; context=neutral",
                        },
                        "topic": {
                            "type": "string",
                            "enum": TOPIC_ENUM,
                            "description": TOPIC_DESCRIPTION,
                        },
                        "definition_id": {
                            "type": ["string", "null"],
                            "enum": DEFINITION_ENUM,
                            "description": "If claim uses a specific defined term, reference its DEF-xxx ID.",
                        },
                        "statement": {
                            "type": "string",
                            "maxLength": 280,
                            "description": "One complete sentence stating the claim with full numeric and contextual detail.",
                        },
                        "locator_hint": {
                            "type": "string",
                            "minLength": 1,
                            "pattern": r".*\S.*",
                            "maxLength": 120,
                            "description": "Section heading, table name, slide, or line reference where this appears.",
                        },
                        "derivation": {
                            "type": ["string", "null"],
                            "description": "Required when epistemic_class=derived. State the computation.",
                        },
                        "derivation_operand_indices": {
                            "type": "array",
                            "items": {"type": "integer", "minimum": 0},
                            "description": (
                                "Only when epistemic_class=derived AND an operand is ITSELF "
                                "one of the claims you are emitting in this SAME response — "
                                "the 0-based position of that operand claim in this claims "
                                "array (position 0 = the first claim you emit here, etc; never "
                                "your own position). Omit or leave empty when an operand is a "
                                "plain number stated in the fragment but not extracted as its "
                                "own claim, or when it isn't in this fragment at all. This is "
                                "the only place a claim can reference another — claim_id does "
                                "not exist yet when you are writing this."
                            ),
                        },
                        "alternative_to_index": {
                            "type": ["integer", "null"],
                            "minimum": 0,
                            "description": (
                                "Only when this claim is a downside/upside/alternative "
                                "scenario of a base-case claim ALSO in this response: that "
                                "base claim's 0-based position in this claims array (same "
                                "position convention as derivation_operand_indices). Null when "
                                "this claim IS the base case, when there is no base-case "
                                "counterpart in this fragment, or when scenario=base/"
                                "unspecified."
                            ),
                        },
                        "evidence_state": {
                            "type": ["string", "null"],
                            "enum": [None, *evidence_state_vocabulary(DEFAULT_ARCHETYPE)],
                            "description": (
                                "WHERE ON THE LADDER this evidence sits, when the fragment says. "
                                "Much of what matters in diligence is a state, not a number: a "
                                "relationship at 'paid pilot' and one at 'production deployment' "
                                "differ in the way that decides the case, and revenue that is "
                                "'contracted' is not revenue that is 'collected'. Pick the state "
                                "the source actually evidences, never the one it hopes for — a "
                                "signed pilot is not a deployment. Null when the fragment does "
                                "not place it, which is common and fine; do not infer a state "
                                "from optimism or from the document's tone."
                            ),
                        },
                        "author": {
                            "type": ["string", "null"],
                            "description": "Party making the claim (e.g. management, QoE provider, IC).",
                        },
                    },
                },
            }
        },
    },
}

def metric_vocabulary(archetype_id: str = DEFAULT_ARCHETYPE) -> list[str]:
    """Metric labels extraction may emit for this archetype.

    buyout returns METRIC_ENUM unchanged -- it IS the buyout vocabulary, and
    the benchmark scores against it, so the default path is byte-identical to
    before this existed. venture/growth widen it from their packs.
    """
    try:
        return extraction_vocabulary(archetype_id, METRIC_ENUM)
    except (KeyError, ValueError, FileNotFoundError):
        # An unknown or unreadable archetype must not silently narrow the
        # vocabulary to nothing; fall back to the buyout baseline and let the
        # caller's own archetype validation surface the problem.
        return list(METRIC_ENUM)


@lru_cache(maxsize=8)
def claim_tool_for(archetype_id: str = DEFAULT_ARCHETYPE) -> dict:
    """CLAIM_TOOL with this archetype's metric enum substituted.

    Everything else about the contract -- identity fields, bounds, derivation
    references, the claim_kind rules -- is archetype-independent on purpose.
    Only the vocabulary of WHAT can be measured changes with strategy; how a
    claim is identified does not.
    """
    if archetype_id == DEFAULT_ARCHETYPE:
        return CLAIM_TOOL
    tool = copy.deepcopy(CLAIM_TOOL)
    properties = tool["input_schema"]["properties"]["claims"]["items"]["properties"]
    properties["metric"]["enum"] = metric_vocabulary(archetype_id)
    # States are archetype-specific and NOT portable: "paid pilot" is a venture
    # state and means nothing on a buyout ladder. Swapping the vocabulary is
    # what keeps a state comparable to its own kind.
    states = evidence_state_vocabulary(archetype_id)
    if states:
        properties["evidence_state"]["enum"] = [None, *states]
    else:
        properties["evidence_state"].pop("enum", None)
    return tool


SYSTEM_PROMPT = textwrap.dedent("""
    You are a financial claim extractor for a private equity firm (PANTA system).
    Extract only claims explicitly stated in the fragment. Never infer or interpolate.
    Return an empty list when the fragment contains no financial claims.

    EXHAUSTIVENESS — a computed total does not replace the rows it came from:
    - When a fragment gives you a table or schedule of individual rows AND asks
      you to total or otherwise compute across them, extract BOTH: every
      individual row as its own claim, AND the computed result(s). Do not stop
      at only the final computed figure because it is what the fragment's
      instruction was building toward -- a reader who only sees "10.0%
      concentration" cannot check that number without the rows it came from.
    - A table with 7 rows and a stated or computed total describes 8 facts,
      not 1. Skipping the rows in favor of the summary is not a shorter,
      cleaner extraction; it is a lossy one.

    METRIC CHOICE — an honest gap beats a near-miss:
    - The metric enum is shaped for buyout underwriting. When a quantity has no
      entry that genuinely names it, emit metric "Other" and put its real name
      in metric_label. Do NOT reach for the closest-sounding entry: a venture
      round is not "Sponsor Equity", a pre-money valuation is not "Enterprise
      Value", and runway months are not an "Exit Horizon". A wrong label is
      worse than "Other" — it silently merges this claim with unrelated ones
      and the mistake becomes invisible.
    - A sentence with no quantity and no checkable content is not a claim at
      all. Do not emit a claim with a null value just to have something to
      return: "use of funds is product, field operations and certification"
      is not a Capex claim, a Free Cash Flow claim, or any other claim.
    - A sentence ABOUT another claim — how to read it, how it relates to
      another figure, or what kind of figure it is — is not itself a new
      claim, even when it sits right next to real numbers and reads like a
      rule or a fact. Extract the real claim(s) it refers to; skip the
      sentence that only comments on them. Three shapes this takes:
        * methodology: "Aggregate accounts only when their ultimate-parent
          field is identical" tells you how to read the schedule, not a
          fact about the business.
        * relationship: "the restated revenue supersedes the original",
          "this confirms the prior figure" restates or explains a change
          between two claims you already have; it is not a third fact
          alongside them.
        * status/qualifier: "this target is not a historical result" says
          what KIND of figure the EBITDA target is (forward-looking, not
          actual) — it qualifies that claim, it is not a second, separate
          NEGATIVE claim about the business.
      Getting this wrong doesn't just add noise: it turns one real claim
      into two or three, which silently hurts precision on every fragment
      shaped like this. When genuinely unsure whether a sentence is its own
      fact or commentary on one you're already extracting, ask whether it
      could be TRUE OR FALSE independent of the claim it sits next to — if
      not, it's commentary.

    EPISTEMIC CLASS — apply strictly by document source and claim type:
    - attested: an INDEPENDENT THIRD PARTY formally certifies someone else's
      numbers — QoE conclusion, auditor opinion. The certifier is not the
      party whose figures are being certified, and is not the firm itself.
    - institutional_act: a body's OWN formal decision or action, not a third
      party certifying anything — the IC's own vote or adopted assumption,
      the firm's own underwriting conclusion, a lender's own credit-agreement
      covenant term, an audit committee's own decision to restate. When in
      doubt for an IC memo or the firm's own underwriting fragment → use
      institutional_act, not attested: the firm attesting to its own
      conclusion is not third-party certification.
    - asserted: seller or management stated claims with no third-party verification.
      CIM numbers, management presentations, seller-deck projections → asserted.
      This is the DEFAULT for a company reporting its own figures, including
      its own schedules and breakdowns: "reported accounts", "revenue was
      $74.0m", a billing-account table, a revenue-by-customer schedule. The
      company stating its own numbers is an assertion, however routine or
      factual it looks.
    - observed: a third party measured or witnessed something DIRECTLY and
      recorded it themselves — a QoE analyst's own workpaper note, a call
      transcript quote, a site visit. The test is whether the RECORDER is
      someone other than the party whose figure it is. A schedule of the
      company's own accounts is not "observed" just because it arrived in a
      data room: nobody outside the company witnessed anything, they were
      handed a table. Reach for observed only when the fragment names a
      third party doing the observing.
    - derived: YOU computed this claim from two or more stated values — requires derivation field.
      Only use when you performed arithmetic (ratio, sum, subtraction).

    Source → class mapping (use this every time):
      IC memo (the firm's own vote/assumption) → institutional_act
      Firm underwriting / initial assessment  → institutional_act
      Credit agreement covenant term (lender's own defined condition) → institutional_act
      Audit committee restatement decision    → institutional_act
      QoE report conclusion  → attested  (independent third party certifies)
      Auditor opinion         → attested  (independent third party certifies)
      Signed / audited / statutory accounts → attested  (formally certified,
        not a bare management claim, even when the fragment never says the
        word "auditor" — "the signed FY2024 accounts reported..." IS the
        certification event, distinct from an unsigned management deck).
        This needs the certification word to actually be there: "signed",
        "audited", "statutory", "auditor". A section merely headed
        "Reported accounts", or a sentence saying the company "reported"
        a figure, carries no certification and is asserted.
      Data room management documents → asserted
      Data room transactional/workpaper observations → observed
      Meeting notes / call transcript / DDQ → observed
      Seller CIM / IM → asserted  (seller's marketing claims)
      Management presentation → asserted
      Computed by you → derived

    EXPLICIT DERIVATIONS — when the source hands you the method and the operands:
    - "Never infer or interpolate" governs facts, not arithmetic the source itself
      asks for. If a fragment states both a computation method and every value it
      needs (e.g. "concentration is calculated by dividing that total by revenue,
      rounded to one decimal place" with the total and the revenue both stated),
      PERFORM that arithmetic and emit the numeric result — do not emit only a
      DEFINITION claim describing the method with a null value.
    - A sentence stating HOW something is calculated is a DEFINITION claim by the
      claim_kind rule below — emit that one, AND SEPARATELY emit a second,
      QUANTITATIVE claim carrying the actual computed number. The DEFINITION
      claim documents the method; it does not substitute for doing the division.
    - Mark the quantitative one epistemic_class=derived, put the numeric result in
      value, and state the computation in derivation (e.g. "13.468 / 74.0 * 100,
      rounded to 1dp").
    - Only refuse the quantitative claim when an operand is genuinely missing from
      the fragment, not when the fragment gives you everything and simply expects
      you to do the division.
    - When an operand of a derived claim IS one of the other claims you are
      emitting right now, set derivation_operand_indices to that operand's
      position in this same response (see the tool schema) so the derivation
      chain is machine-checkable, not just prose.
    - Emit the INTERMEDIATE step, do not jump the chain. If getting to the
      answer runs through a subtotal — sum seven account rows into one
      customer's total, then divide that total by revenue — the subtotal is
      its own derived claim, and the final figure derives from THAT subtotal
      plus the divisor, not from the seven rows directly. A chain someone
      can re-check one step at a time is the whole point: "7.0 = 4.0 + 3.0"
      and "10.0% = 7.0 / 70.0" are each verifiable on their own, while
      "10.0% came from these eight numbers somehow" is not. Collapsing the
      chain hides the step most likely to be wrong.
    - Base/downside/upside are NOT independent claims with no relationship to
      each other. When you emit a downside or upside scenario claim and its
      base case is ALSO in this response, set alternative_to_index to the
      base claim's position (see the tool schema) — do not leave scenario
      alone to carry that link.

    PERIOD EXTRACTION — mandatory for every emitted claim:
    - Never leave period blank. Read the workbook column header, table header,
      section title, page heading, or source date when it is not repeated in the row.
    - Preserve the source language: FY2024, FY2025A, LTM Sep-25, Q2 2025,
      Opening Balance Sheet, Budget 2026, or 'as of 2025-10-27'.
    - If the only available time reference is the source effective date, use
      'as of {source effective date}' rather than leaving period blank.
    - Use 'cross-period' only for an explicitly time-invariant fact.
    - Use 'recurring' for a condition re-tested every time its trigger
      recurs — a covenant checked at each quarter end, a requirement stated
      to apply 'at all times' or 'on each measurement date', or a condition
      re-applied on every occurrence of an event ('immediately after each
      acquisition'). This is not cross-period: a recurring test has a
      cadence (calendar or event-triggered), a time-invariant fact does not.

    PERIMETER INFERENCE — mandatory when the source determines scope:
    - Use exactly one canonical evidence-view label from document type:
      Seller View (CIM/management report), QoE View (QoE report),
      Firm View (firm internal memo), or Statutory (audited accounts).
    - Combine that view with the entity and structural scope (standalone,
      consolidated, deal level) rather than returning only the generic label.
    - For concentration, distinguish billing-account from ultimate-parent scope.
    - Use 'unknown' only when neither the fragment nor source metadata supports a scope.

    LOCATOR HINT:
    - The deterministic fragment locator is mandatory, never blank, and always
      retained by the pipeline with section and page/slide metadata when available.
    - Add locator_hint only when the row, cell, table, or subsection inside the
      fragment can be identified. Accepted hints include 'slide 12', 'page 7',
      'Sheet1!D42', 'section "Revenue Analysis"', 'table row 3', and
      'timestamp 00:14:22'; never invent a page or cell address.

    Identity rule — write the FULL perimeter and FULL period, not a shorthand:
    - period: use the document's own language including context, e.g.
        'FY2025A / FY2025E seller presentation'  not just 'FY2025A'
        'As of 2025-10-27'  not just '2025-10-27'
        'FY2026E–FY2030E forecast'  not just 'FY2026E'
    - perimeter: write the complete scope string, e.g.
        'Alderstone consolidated EBITDA under seller adjustment perimeter'
        'Alderstone consolidated EBITDA under independent QoE adjustment perimeter'
        'Alderstone customer revenue measured by individual billing account'
      NOT just 'Alderstone consolidated' or 'Alderstone standalone'.

    Do not collapse different definitions:
    - Firm EBITDA ≠ QoE EBITDA ≠ Covenant EBITDA ≠ Seller EBITDA
    - account-level concentration ≠ parent-level concentration
    - FY2025A ≠ LTM ≠ FY2026E

    A number the source itself disclaims is not a claim about the world, even
    though a number appears in the sentence. "Illustrative only", "for
    drafting purposes", "hypothetical example", "not an operative covenant",
    "not currently in effect" — when the source explicitly says a figure does
    not describe an actual, current state of affairs, do not emit it as a
    claim at all. This is not the same as omitting an ambiguous claim: the
    source is not silent here, it is affirmatively telling you the number is
    not real. Extracting it as a covenant, threshold, or fact anyway is
    inventing the exact thing the sentence just told you isn't true.

    No source statement may be invented. If the fragment is ambiguous, omit the claim.
""").strip()


# ─────────────────────────────────────────────────────────────────────────────
# L1 — Document parser (deterministic)
# ─────────────────────────────────────────────────────────────────────────────

CHUNK_WORDS = 250


@dataclass
class Chunk:
    chunk_id: str
    locator: str
    body: str
    source_path: str
    source_type: str
    source_record: dict
    word_count: int
    section_heading: str | None = None
    page_or_slide_number: int | None = None
    provenance: dict[str, Any] = field(default_factory=dict)
    period_context: dict[str, Any] = field(default_factory=dict)
    # Every individual heading this chunk covers, when parse_markdown merged
    # several short sections into one. section_heading is their combined
    # label (for display); this is the list a claim's own locator_hint is
    # resolved against, so a claim points at ITS section rather than at the
    # whole merged span. See _claim_locator().
    section_headings: list[str] = field(default_factory=list)


class UnsupportedSourceError(ValueError):
    """Raised with a machine-readable response when V2 cannot parse safely."""

    def __init__(self, message: str, response: dict[str, Any] | None = None):
        super().__init__(message)
        self.response = response or {
            "schema": CAPABILITY_SCHEMA,
            "status": "REJECTED",
            "code": "UNSUPPORTED_SOURCE",
            "detail": message,
        }

    def to_dict(self) -> dict[str, Any]:
        return dict(self.response)


def _chunk_hash(body: str) -> str:
    return "ch-" + hashlib.sha256(body.encode()).hexdigest()[:12]


def _split_words(text: str, max_words: int, locator_prefix: str,
                 source_path: str, source_type: str, source_record: dict,
                 section_heading: str | None = None,
                 page_or_slide_number: int | None = None) -> list[Chunk]:
    words = text.split()
    chunks = []
    for i in range(0, len(words), max_words):
        body = " ".join(words[i: i + max_words])
        locator = f"{locator_prefix}:w{i}-{i + len(body.split())}"
        chunks.append(Chunk(
            chunk_id=_chunk_hash(body),
            locator=locator,
            body=body,
            source_path=source_path,
            source_type=source_type,
            source_record=source_record,
            word_count=len(body.split()),
            section_heading=section_heading,
            page_or_slide_number=page_or_slide_number,
        ))
    return chunks


def _reject_source(
    path: Path,
    code: str,
    detail: str,
    *,
    capability_id: str | None = None,
    action: str | None = None,
) -> UnsupportedSourceError:
    response = capability_failure(
        path,
        code,
        capability_id=capability_id,
        action=action,
        detail=detail,
    )
    message = detail
    if response.get("action"):
        message += f" Action: {response['action']}"
    return UnsupportedSourceError(message, response)


def _decorate_chunks(
    chunks: list[Chunk],
    path: Path,
    source_record: dict[str, Any],
    capability_id: str,
) -> list[Chunk]:
    """Attach source identity and honest period semantics to every fragment."""
    envelope = source_record.get("source_envelope") or {}
    for chunk in chunks:
        chunk.provenance = {
            **chunk.provenance,
            "source_id": source_record.get("source_id"),
            "source_version_id": envelope.get("source_version_id"),
            "case_id": envelope.get("case_id"),
            "original_filename": envelope.get("original_filename") or path.name,
            "locator": chunk.locator,
            "parser_capability": capability_id,
            "capability_contract": CAPABILITY_SCHEMA,
        }
        specific = dict(chunk.period_context)
        specific.update({
            "effective_date": envelope.get("effective_date") or source_record.get("effective_date") or None,
            # A chunk-specific known_at (e.g. one message's own Date header
            # in a multi-message .mbox thread) must survive this generic
            # decoration pass, not get overwritten by the source-level
            # default -- that would make every message in a thread carry
            # the same known_at regardless of when it was actually sent.
            "known_at": (
                specific.get("known_at")
                or envelope.get("known_at")
                or source_record.get("known_at")
                or None
            ),
            "semantics": "DECLARED_ONLY",
            "content_period_policy": "preserve source label; never infer at L1",
        })
        chunk.period_context = specific
    return chunks


class _VisibleHTML(HTMLParser):
    """Small deterministic HTML-to-text reader; scripts/styles are excluded."""

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._hidden_depth = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._hidden_depth += 1
        elif not self._hidden_depth and tag.lower() in {"p", "div", "br", "li", "tr", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._hidden_depth:
            self._hidden_depth -= 1
        elif not self._hidden_depth and tag.lower() in {"p", "div", "li", "tr", "h1", "h2", "h3"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._hidden_depth:
            self.parts.append(data)

    def text(self) -> str:
        return "\n".join(line.strip() for line in "".join(self.parts).splitlines() if line.strip())


def _chunks_from_numbered_lines(
    lines: list[str],
    path: Path,
    max_words: int,
    source_type: str,
    source_record: dict[str, Any],
    locator_label: str = "lines",
) -> list[Chunk]:
    chunks: list[Chunk] = []
    pending: list[str] = []
    start = end = 0

    def flush() -> None:
        nonlocal pending, start, end
        if not pending:
            return
        body = "\n".join(pending)
        chunks.append(Chunk(
            chunk_id=_chunk_hash(body),
            locator=f"{path.name}::{locator_label}:{start}-{end}",
            body=body,
            source_path=str(path),
            source_type=source_type,
            source_record=source_record,
            word_count=len(body.split()),
        ))
        pending, start, end = [], 0, 0

    for line_number, line in enumerate(lines, 1):
        if not line.strip():
            continue
        projected = len((" ".join(pending + [line])).split())
        if pending and projected > max_words:
            flush()
        if not pending:
            start = line_number
        end = line_number
        pending.append(line.strip())
    flush()
    return chunks


def parse_markdown(path: Path, max_words: int = CHUNK_WORDS,
                   source_record: dict | None = None) -> list[Chunk]:
    """Split a markdown file into chunks on ``##``/``###`` headings, merging
    adjacent SHORT sections up to ``max_words`` instead of giving every
    heading its own chunk regardless of size.

    A "packet" that bundles several short excerpts under separate headings
    (a reported-accounts line, a QoE line, a customer schedule) used to
    become one chunk per heading even when the whole document was a
    fraction of the word budget. annotate_chunk has no cross-chunk memory,
    so a claim whose operands sit in two of those headings -- "concentration
    is total divided by revenue" where the total is in one section and the
    revenue in another -- was unresolvable by any single call: not because
    the model couldn't do the arithmetic, but because it never saw both
    numbers at once. Merging keeps that context together whenever it fits.
    """
    src = source_record or _source_record(path)
    text = path.read_text(encoding="utf-8")
    parts = re.split(r"(?=^#{2,3} )", text, flags=re.MULTILINE)
    chunks: list[Chunk] = []
    pending: list[str] = []
    pending_words = 0
    # Every REAL heading folded into the pending chunk, not just the first --
    # a headerless intro paragraph ahead of the first "##" must never blank
    # out the actual section name once something real merges in behind it.
    # annotate_chunk puts this straight into the model's prompt as
    # "SECTION HEADING: ...", so losing it back to a generic "section" label
    # is a real information loss to the model, not a cosmetic one.
    pending_headings: list[str] = []

    def _combined_label() -> str:
        if not pending_headings:
            return "section"
        if len(pending_headings) <= 3:
            return " + ".join(pending_headings)
        return " + ".join(pending_headings[:3]) + f" (+{len(pending_headings) - 3} more)"

    def flush() -> None:
        nonlocal pending, pending_words, pending_headings
        if not pending:
            return
        body = "\n".join(pending).strip()
        if body:
            label = _combined_label()
            chunks.append(Chunk(
                chunk_id=_chunk_hash(body),
                locator=f"{path.name}::{label}",
                body=body,
                source_path=str(path),
                source_type="markdown",
                source_record=src,
                word_count=pending_words,
                section_heading=label,
                section_headings=list(pending_headings),
            ))
        pending, pending_words, pending_headings = [], 0, []

    for part in parts:
        if not part.strip():
            continue
        header_match = re.match(r"^(#{2,3} .+)", part)
        section_label = header_match.group(1)[:50].strip() if header_match else "section"
        words = part.split()
        if len(words) > max_words:
            # Too big to merge with anything -- flush whatever is pending,
            # unchanged from before: split this one section on its own.
            flush()
            sub = _split_words(part, max_words, f"{path.name}::{section_label}",
                               str(path), "markdown", src,
                               section_heading=section_label)
            chunks.extend(sub)
            continue
        if pending and pending_words + len(words) > max_words:
            flush()
        if header_match:
            pending_headings.append(section_label)
        pending.append(part.strip())
        pending_words += len(words)
    flush()
    return chunks


def parse_plain_text(path: Path, max_words: int = CHUNK_WORDS,
                     source_record: dict | None = None,
                     source_type: str = "text") -> list[Chunk]:
    src = source_record or _source_record(path)
    try:
        text = path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise _reject_source(
            path,
            "TEXT_ENCODING_UNSUPPORTED",
            f"{path.name} is not valid UTF-8 text ({exc}).",
            action="Export the artifact as UTF-8 without replacing source characters.",
        ) from exc
    return _chunks_from_numbered_lines(text.splitlines(), path, max_words, source_type, src)


def parse_html(path: Path, max_words: int = CHUNK_WORDS,
               source_record: dict | None = None) -> list[Chunk]:
    src = source_record or _source_record(path)
    try:
        markup = path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise _reject_source(
            path,
            "TEXT_ENCODING_UNSUPPORTED",
            f"{path.name} is not valid UTF-8 HTML ({exc}).",
            action="Export the page as UTF-8 HTML or plain text.",
        ) from exc
    reader = _VisibleHTML()
    reader.feed(markup)
    return _chunks_from_numbered_lines(
        reader.text().splitlines(), path, max_words, "html", src, "visible-lines"
    )


def parse_csv(path: Path, max_words: int = CHUNK_WORDS,
              source_record: dict | None = None) -> list[Chunk]:
    src = source_record or _source_record(path)
    try:
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
    except (UnicodeDecodeError, csv.Error) as exc:
        raise _reject_source(
            path,
            "CSV_INVALID",
            f"Cannot parse {path.name} as UTF-8 CSV: {exc}.",
            action="Export a UTF-8 CSV with a single consistent delimiter and header row.",
        ) from exc
    lines = [" | ".join(f"c{column}={value}" for column, value in enumerate(row, 1)) for row in rows]
    return _chunks_from_numbered_lines(lines, path, max_words, "csv", src, "rows")


def _read_openxml_root(path: Path, member: str, capability_id: str) -> ElementTree.Element:
    try:
        with zipfile.ZipFile(path) as archive:
            payload = archive.read(member)
        return ElementTree.fromstring(payload)
    except (OSError, KeyError, zipfile.BadZipFile, ElementTree.ParseError) as exc:
        raise _reject_source(
            path,
            "OPENXML_INVALID",
            f"{path.name} is not a valid {capability_id.upper()} package: {exc}.",
            capability_id=capability_id,
            action=f"Open the artifact in its authoring application and save a valid {path.suffix.upper()} copy.",
        ) from exc


def _docx2python_render_section(section: list) -> str | None:
    """Render one docx2python body section as markdown-table text if it's a
    real table (more than one row, or more than one cell in its one row),
    or as plain paragraph text otherwise. docx2python wraps every paragraph
    in the same row/cell/paragraph-list shape as a 1x1 "table" -- this is
    what tells a real w:tbl apart from ordinary body text without a second
    XML pass."""
    rows = []
    for row in section:
        cells = ["\n".join(p for p in cell if p.strip()).strip() for cell in row]
        if any(cells):
            rows.append(cells)
    if not rows:
        return None
    if len(rows) == 1 and len(rows[0]) == 1:
        return rows[0][0]

    ncols = max(len(row) for row in rows)
    lines = []
    for row_index, row in enumerate(rows):
        padded = row + [""] * (ncols - len(row))
        safe = [cell.replace("|", "\\|").replace("\n", " ") for cell in padded]
        lines.append("| " + " | ".join(safe) + " |")
        if row_index == 0:
            lines.append("| " + " | ".join(["---"] * ncols) + " |")
    return "\n".join(lines)


def _parse_docx_openxml_text_only(path: Path, max_words: int, src: dict) -> list[Chunk]:
    """Paragraph-text-only fallback: no table structure, no images, no
    comments, no tracked-changes revision status. This is the whole DOCX
    pipeline when docx2python isn't installed -- a real, optional
    dependency (see PAN-103), not assumed always present."""
    root = _read_openxml_root(path, "word/document.xml", "docx")
    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    paragraphs: list[str] = []
    for paragraph in root.iter(namespace + "p"):
        text = "".join(node.text or "" for node in paragraph.iter(namespace + "t")).strip()
        if text:
            paragraphs.append(text)
    return _chunks_from_numbered_lines(paragraphs, path, max_words, "docx", src, "paragraphs")


def parse_docx(path: Path, max_words: int = CHUNK_WORDS,
               source_record: dict | None = None) -> list[Chunk]:
    """Extract DOCX paragraphs and real table structure via docx2python.

    Falls back to plain paragraph text (no tables) when docx2python isn't
    installed -- a real, optional dependency. Verified against a real
    41-table Keystone document: the previous paragraph-only reader
    flattened every table into run-on text with no row/column boundary at
    all (e.g. a 6-row, 2-column "what changed" table became one
    undifferentiated paragraph); this preserves the grid.

    Tracked-changes (w:ins/w:del) status and comments are not yet
    surfaced here -- PAN-103's docx-revisions integration is separate,
    real remaining work, not silently assumed done by this function.
    """
    src = source_record or _source_record(path)
    _read_openxml_root(path, "word/document.xml", "docx")  # validity check, existing error contract

    try:
        from docx2python import docx2python
    except ImportError:
        return _parse_docx_openxml_text_only(path, max_words, src)

    with docx2python(str(path)) as doc:
        blocks = [
            rendered
            for section in doc.body
            if (rendered := _docx2python_render_section(section)) is not None
        ]
    return _chunks_from_numbered_lines(blocks, path, max_words, "docx", src, "blocks")


def _pptx_table_text(table: Any) -> str:
    """Render a native PPTX table as a structured grid, not flattened text."""
    lines = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
    return "Table:\n" + "\n".join(lines)


def _pptx_chart_text(chart: Any) -> str:
    """Pull real category/series/value data out of a native OOXML chart.

    PE decks' charts are almost always live chart parts (c:chart), not
    flat pictures -- the cached numCache/strCache in the chart XML holds
    exact numbers. python-pptx exposes this directly, so a bar chart or
    revenue bridge becomes checkable data instead of an opaque image with
    nothing to verify a claim against.
    """
    try:
        plot = chart.plots[0]
    except (IndexError, ValueError):
        return ""
    categories = [str(c) for c in plot.categories]
    lines = [f"Chart ({chart.chart_type}):"]
    for series in plot.series:
        values = list(series.values)
        if categories and len(categories) == len(values):
            pairs = ", ".join(f"{cat}={val}" for cat, val in zip(categories, values))
        else:
            pairs = ", ".join(str(v) for v in values)
        lines.append(f"  {series.name}: {pairs}")
    return "\n".join(lines)


def _pptx_shape_text(shape: Any) -> list[str]:
    """Text, table grid, or chart data for one shape, prefixed by its own
    name so a claim can trace back to the exact shape on the slide, not
    just the slide as a whole."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    parts: list[str] = []
    if getattr(shape, "has_chart", False):
        chart_text = _pptx_chart_text(shape.chart)
        if chart_text:
            parts.append(f"[{shape.name}] {chart_text}")
    elif getattr(shape, "has_table", False):
        parts.append(f"[{shape.name}] {_pptx_table_text(shape.table)}")
    elif getattr(shape, "has_text_frame", False):
        text = (shape.text_frame.text or "").strip()
        if text:
            parts.append(text)
    elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
        # A chart pasted as a flat image (no native chart XML) has no
        # series data to read here. This is a declared coverage limit, not
        # a silent drop: a human or a later vision-capable pass can see
        # exactly which shape on which slide still needs review, instead
        # of the picture simply never having existed as far as any chunk
        # is concerned.
        parts.append(
            f"[{shape.name}] IMAGE_NOT_EXTRACTED: a picture is present on this "
            "slide with no native chart data behind it; its content (chart, "
            "diagram, or photo) was not read."
        )
    elif getattr(shape, "shape_type", None) is None and getattr(shape, "has_chart", None) is not None:
        # A GraphicFrame whose graphicData python-pptx doesn't recognize as
        # chart/table/OLE -- shape_type's own docstring names this exact
        # case as SmartArt, and it is also what a modern "chart-ex" chart
        # (waterfall, funnel, histogram, box-and-whisker -- the family that
        # covers PowerPoint's native EBITDA-bridge chart type) produces,
        # since chart-ex uses a different graphicData namespace that
        # python-pptx's has_chart never matches. Both fall through every
        # branch above with no text, no marker, and no error today. Naming
        # the gap explicitly is cheap even before either is truly readable.
        kind = "shape"
        graphic_data_uri = getattr(getattr(shape, "_element", None), "graphicData_uri", None)
        if graphic_data_uri and "chartex" in graphic_data_uri.lower():
            kind = "modern chart (waterfall/funnel/histogram/box-whisker-family)"
        parts.append(
            f"[{shape.name}] UNSUPPORTED_GRAPHIC_FRAME: a {kind} is present on "
            "this slide that this parser does not yet read (SmartArt or a "
            "chart-ex chart type); its content was not extracted."
        )
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            parts.extend(_pptx_shape_text(sub_shape))
    return parts


def _pptx_reading_order(shapes: Any, slide_height: Any) -> list[Any]:
    """Order a slide's shapes as a person reads them, not as PowerPoint stored them.

    `slide.shapes` yields z-order -- the sequence shapes were added to the
    XML -- which has nothing to do with where they sit. A value box drawn
    before its label emits before that label, so the text is all present and
    the relationship between the pieces is gone: content correct, position
    lost.

    Shapes are banded into rows first, then read left to right within a row.
    Sorting on `top` alone would interleave two columns of a two-column slide,
    which is worse than z-order rather than better. The band is a twentieth of
    the slide, coarse enough that a caption sitting a few points above its
    neighbour still counts as the same row.

    A shape with no explicit position (a placeholder inheriting from the
    layout) keeps its original relative order at the end rather than being
    assigned a position it does not have.
    """
    band = max(1, int(slide_height / 20)) if slide_height else 1

    def key(item: tuple[int, Any]) -> tuple:
        index, shape = item
        top, left = getattr(shape, "top", None), getattr(shape, "left", None)
        if top is None or left is None:
            return (1, 0, 0, index)
        return (0, int(top) // band, int(left), index)

    return [shape for _, shape in sorted(enumerate(shapes), key=key)]


def parse_pptx(path: Path, max_words: int = CHUNK_WORDS,
               source_record: dict | None = None,
               pdf_engine: Any = None) -> list[Chunk]:
    """Extract shape text, native table grids, native chart data, and
    speaker notes per slide. A chart pasted as a flat picture (no native
    chart XML) has no series data to read here and is not invented; a
    declared IMAGE_NOT_EXTRACTED marker names the shape so the gap is
    visible instead of the picture simply never having existed.

    Native chart/table data (COLUMN_CLUSTERED (51): Q1=10.0, Q2=15.0,
    Q3=22.0 -- verified against the real revenue-review fixture) is exact:
    it comes from the chart's own XML, not a model reading rendered pixels.
    It stays the primary and only path for any slide that has it.

    Only a slide carrying a declared IMAGE_NOT_EXTRACTED or
    UNSUPPORTED_GRAPHIC_FRAME marker -- meaning python-pptx found nothing
    to read there at all -- additionally gets a PDF-rendered second look:
    the whole deck exported once via LibreOffice headless and read through
    the same convert_page pipeline a PDF page uses. This recovers content
    for slides with SmartArt, chart-ex charts, or a chart pasted as a flat
    image; it never overrides a slide that already has native data, so it
    cannot reintroduce the PAN-100 "confidently wrong chart pixels" failure
    on a slide that did not need it."""
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as exc:
        raise _reject_source(
            path,
            "READER_UNAVAILABLE",
            "PPTX parsing is unavailable because python-pptx is not installed.",
            capability_id="pptx",
        ) from exc
    try:
        presentation = Presentation(str(path))
    except (PackageNotFoundError, OSError, KeyError, ValueError) as exc:
        raise _reject_source(
            path,
            "OPENXML_INVALID",
            f"{path.name} is not a valid PPTX package: {exc}.",
            capability_id="pptx",
            action="Open the deck in PowerPoint and save a valid .pptx copy.",
        ) from exc

    src = source_record or _source_record(path)
    chunks: list[Chunk] = []
    rendered_slides: dict[int, str] | None = None      # computed lazily, once
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in _pptx_reading_order(slide.shapes, presentation.slide_height):
            parts.extend(_pptx_shape_text(shape))
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                parts.append(f"Speaker notes: {notes_text}")

        if _slide_needs_pdf_fallback(parts):
            if rendered_slides is None:
                rendered_slides = _render_pptx_as_pdf(path, pdf_engine) or {}
            fallback = rendered_slides.get(slide_number, "").strip()
            if fallback:
                parts.append(
                    f"[pdf-fallback] slide rendered as an image and read because "
                    f"native extraction found an unreadable shape here:\n{fallback}"
                )

        body = "\n".join(part for part in parts if part)
        if not body:
            continue
        chunks.extend(_split_words(
            body, max_words, f"{path.name}::slide:{slide_number}", str(path), "pptx", src,
            page_or_slide_number=slide_number,
        ))
    return chunks


def _slide_needs_pdf_fallback(parts: list[str]) -> bool:
    """Whether a slide carries a declared native-reading gap worth a second,
    PDF-rendered look. The two markers this checks are the only ones
    _pptx_shape_text emits for "found the shape, could not read it" --
    IMAGE_NOT_EXTRACTED (a flat picture/chart with no native XML) and
    UNSUPPORTED_GRAPHIC_FRAME (SmartArt or a chart-ex chart type)."""
    return any("IMAGE_NOT_EXTRACTED" in part or "UNSUPPORTED_GRAPHIC_FRAME" in part
              for part in parts)


def _render_pptx_as_pdf(path: Path, pdf_engine: Any) -> dict[int, str] | None:
    """Export a whole deck to PDF via LibreOffice headless, then read every
    page with the normal PDF page-conversion path. Returns None (never
    raises) if soffice is missing or the export fails -- a fallback that
    can itself fail hard would be worse than no fallback.

    `pdf_engine`, if supplied, is a FACTORY: given the exported PDF's path,
    it returns a bound convert_page(image, page_num) callable -- not a
    pre-bound callable itself. parse_pdf's convert_page can be bound ahead
    of time because the caller already knows which file it is converting;
    here the PDF this function converts does not exist until soffice creates
    it, so whatever engine reads it has to learn that path too. The eval
    harness's _paddle_convert is already exactly this shape (Path -> bound
    callable), so it plugs in directly.
    """
    import shutil
    import subprocess
    import tempfile
    import uuid

    if shutil.which("soffice") is None:
        return None

    with tempfile.TemporaryDirectory() as outdir:
        profile = Path(tempfile.gettempdir()) / f"lo-profile-{uuid.uuid4().hex}"
        try:
            subprocess.run(
                ["soffice", "--headless", "--norestore",
                 f"-env:UserInstallation=file://{profile}",
                 "--convert-to", "pdf", "--outdir", outdir, str(path)],
                capture_output=True, timeout=120, check=True,
            )
        except Exception as exc:
            print(f"  [PPTX] soffice export failed for {path.name}: {exc}", file=sys.stderr)
            return None
        finally:
            shutil.rmtree(profile, ignore_errors=True)

        exported = Path(outdir) / f"{path.stem}.pdf"
        if not exported.exists():
            return None

        convert_page = pdf_engine(exported) if pdf_engine is not None else _default_convert_page()
        if convert_page is None:
            return None

        import pdfplumber
        pages: dict[int, str] = {}
        try:
            with pdfplumber.open(exported) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    try:
                        image = page.to_image(resolution=200).original
                        markdown, _ = convert_page(image, page_num)
                        if not _is_degenerate_repetition(markdown):
                            pages[page_num] = markdown
                    except Exception as exc:
                        print(f"  [PPTX] page {page_num} render failed: {exc}", file=sys.stderr)
        except Exception as exc:
            print(f"  [PPTX] could not open exported PDF: {exc}", file=sys.stderr)
            return None
        return pages


def parse_transcript(path: Path, max_words: int = CHUNK_WORDS,
                     source_record: dict | None = None) -> list[Chunk]:
    src = source_record or _source_record(path)
    if path.suffix.lower() == ".txt":
        return parse_plain_text(path, max_words, src, "transcript")
    try:
        text = path.read_text(encoding="utf-8-sig").replace("\r\n", "\n")
    except UnicodeDecodeError as exc:
        raise _reject_source(
            path, "TEXT_ENCODING_UNSUPPORTED", f"Transcript is not valid UTF-8: {exc}.",
            capability_id="transcript_export",
            action="Export the transcript as UTF-8 WebVTT, SRT, or plain text.",
        ) from exc
    blocks = [block.strip() for block in re.split(r"\n\s*\n", text) if block.strip()]
    chunks: list[Chunk] = []
    cue_number = 0
    for block in blocks:
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines or lines[0].upper() == "WEBVTT":
            continue
        timing_index = next((i for i, line in enumerate(lines) if "-->" in line), None)
        if timing_index is None:
            continue
        cue_number += 1
        timecode = lines[timing_index]
        body = "\n".join(lines[timing_index + 1:]).strip()
        if not body:
            continue
        locator = f"{path.name}::cue:{cue_number}:{timecode.replace(' ', '')}"
        chunks.extend(_split_words(body, max_words, locator, str(path), "transcript", src))
    return chunks


def _email_body(message: Any) -> str:
    plain: list[str] = []
    html_parts: list[str] = []
    parts = message.walk() if message.is_multipart() else [message]
    for part in parts:
        if part.is_multipart() or part.get_content_disposition() == "attachment":
            continue
        content_type = part.get_content_type()
        if content_type not in {"text/plain", "text/html"}:
            continue
        try:
            content = part.get_content()
        except (LookupError, UnicodeDecodeError):
            payload = part.get_payload(decode=True) or b""
            content = payload.decode("utf-8", errors="replace")
        if content_type == "text/plain":
            plain.append(str(content))
        else:
            reader = _VisibleHTML()
            reader.feed(str(content))
            html_parts.append(reader.text())
    return "\n".join(plain or html_parts).strip()


def _strip_quoted_reply_history(body: str) -> tuple[str, bool]:
    """Drop quoted prior-message text and signatures from a reply body.

    Without this, a thread's later .mbox message re-surfaces every earlier
    message's full text inline -- the same statement gets chunked once per
    reply, each time attributed to the later message's own known_at instead
    of the message that actually said it. email_reply_parser is a small,
    zero-dependency, regex/heuristic library built exactly for this (Gmail
    "On DATE, X wrote:" headers, ">" quote markers, "---- Original
    Message ----" blocks, signature blocks) -- deliberately not `talon`
    (PAN-104's other candidate), which pulls in scikit-learn/scipy/numpy for
    a signature-detection feature this ticket doesn't need.
    """
    from email_reply_parser import EmailReplyParser

    try:
        reply_only = EmailReplyParser.parse_reply(body).strip()
    except Exception:
        return body, False
    # A pure forward with no original commentary parses down to just the
    # "---- Forwarded message ----"-style boundary line, not to nothing --
    # never let stripping silently erase an entire message's only content.
    if not reply_only or reply_only == body.strip() or _is_boilerplate_only(reply_only):
        return body, False
    return reply_only, True


_QUOTE_BOUNDARY_LINE = re.compile(
    r"^[-_=]{3,}|^(from|date|subject|to|cc|sent):", re.IGNORECASE
)


def _is_boilerplate_only(text: str) -> bool:
    """True if every non-blank line is a separator/header-style boundary
    marker (e.g. "---- Forwarded message ----", "From: ...") rather than
    actual message content.
    """
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return bool(lines) and all(_QUOTE_BOUNDARY_LINE.match(line) for line in lines)


def _email_date(message: Any) -> str | None:
    raw = message.get("Date")
    if not raw:
        return None
    try:
        return parsedate_to_datetime(raw).isoformat()
    except (TypeError, ValueError, OverflowError):
        return None


def _email_headers(message: Any) -> dict[str, str]:
    """The header fields a claim about "who sent this" needs to point at.

    Subject previously only reached section_heading, and From/To/Message-Id
    were not surfaced anywhere on the chunk -- so a caller had the body but
    no way to answer "who is this from" without re-parsing the raw file
    itself. These are message metadata, not attachment content, so exposing
    them does not touch invariant 6 (artifacts are never copied into the
    vault); a header is not the artifact.
    """
    headers = {}
    for name in ("Subject", "From", "To", "Cc", "Message-Id"):
        value = message.get(name)
        if value:
            headers[name.lower().replace("-", "_")] = str(value)
    return headers


def _email_attachments(message: Any) -> list[str]:
    """Attachment filenames only -- never their content (invariant 6)."""
    names = []
    for part in message.walk():
        if part.get_content_disposition() == "attachment":
            name = part.get_filename()
            if name:
                names.append(name)
    return names


def parse_email(path: Path, max_words: int = CHUNK_WORDS,
                source_record: dict | None = None) -> list[Chunk]:
    src = source_record or _source_record(path)
    if path.suffix.lower() == ".eml":
        try:
            messages = [BytesParser(policy=email_policy.default).parsebytes(path.read_bytes())]
        except (OSError, ValueError) as exc:
            raise _reject_source(
                path, "EMAIL_INVALID", f"Cannot parse RFC 822 email: {exc}.",
                capability_id="email_export",
                action="Export the message as RFC 822 .eml with complete headers.",
            ) from exc
    else:
        # Default mbox factory yields compat32 Message objects with no
        # get_content() -- _email_body would crash on every real .mbox file.
        # Use the same modern EmailPolicy the .eml branch above already does.
        box = mailbox.mbox(
            path, factory=lambda f: BytesParser(policy=email_policy.default).parse(f), create=False
        )
        try:
            messages = [message for message in box]
        finally:
            box.close()
    chunks: list[Chunk] = []
    for message_number, message in enumerate(messages, 1):
        body = _email_body(message)
        if not body:
            continue
        body, quoted_history_stripped = _strip_quoted_reply_history(body)
        if not body:
            continue
        attachment_names = _email_attachments(message)
        headers = _email_headers(message)
        message_chunks = _split_words(
            body, max_words, f"{path.name}::message:{message_number}:body",
            str(path), "email", src,
            section_heading=str(message.get("Subject") or "(no subject)"),
        )
        document_date = _email_date(message)
        for chunk in message_chunks:
            chunk.provenance = {
                "headers": headers,
                "excluded_attachments": len(attachment_names),
                "attachment_filenames": attachment_names,
                "attachment_policy": "SEPARATE_SOURCE_ENVELOPE_REQUIRED",
                "quoted_reply_history_stripped": quoted_history_stripped,
            }
            chunk.period_context = {
                "document_date": document_date,
                "document_date_source": "email-header:Date" if document_date else None,
            }
        chunks.extend(message_chunks)
    return chunks


IMAGE_SUFFIXES = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff")


def parse_image(path: Path, max_words: int = CHUNK_WORDS,
                source_record: dict | None = None, convert_page: Any = None,
                vision_fallback: Any = None) -> list[Chunk]:
    """Read a standalone image (a screenshot, a chart exported as PNG) with
    the same model pipeline a PDF page gets, instead of leaving image
    formats entirely unsupported.

    A raster image has no PDF structure, so it is wrapped as a one-page PDF
    (a real, if degenerate, PDF -- pdfplumber opens it, page.to_image()
    just re-renders the embedded raster) and run through the identical
    per-page conversion parse_pdf uses: same convert_page injection point,
    same degenerate-repetition check, same declared-gap philosophy for
    anything the model cannot read.

    `vision_fallback`, if supplied, is called with the ORIGINAL image path
    when the local model path produced nothing usable -- never by default.
    This function makes no network call itself and never will: any fallback
    that leaves this machine is the caller's explicit choice to inject, the
    same way convert_page already is. Wiring one in as a production default
    is a policy-table decision (invariant 7), not something to default here.
    """
    import tempfile

    from PIL import Image

    convert_page = convert_page or _default_convert_page()
    src = source_record or _source_record(path)

    markdown, picture_classes = "", []
    if convert_page is not None:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as handle:
            Image.open(path).convert("RGB").save(handle.name, "PDF")
            wrapped = Path(handle.name)
        try:
            import pdfplumber
            with pdfplumber.open(wrapped) as pdf:
                image = pdf.pages[0].to_image(resolution=200).original
                markdown, picture_classes = convert_page(image, 1)
                if _is_degenerate_repetition(markdown):
                    markdown = ""
        except Exception as exc:
            print(f"  [IMAGE] model failed on {path.name} ({exc})", file=sys.stderr)
            markdown = ""
        finally:
            wrapped.unlink(missing_ok=True)

    if not markdown.strip() and vision_fallback is not None:
        try:
            transcription = vision_fallback(path)
        except Exception as exc:
            transcription = f"[vision-fallback failed: {exc}]"
        if transcription and transcription.strip():
            markdown = (f"[vision-fallback, MODEL-DERIVED transcription, not read text]\n"
                       f"{transcription.strip()}")

    if not markdown.strip():
        raise _reject_source(
            path, "IMAGE_UNREADABLE",
            "No local model or vision fallback produced content for this image.",
            capability_id="native_image",
            action="Install the PDF model stack (see deploy/README.md) or supply a vision_fallback.",
        )

    body_parts = [markdown.strip()]
    for picture_class in picture_classes:
        body_parts.append(
            f"[picture] IMAGE_NOT_EXTRACTED: a {picture_class} is present in this "
            "image; its content was not reliably extracted."
        )
    return _chunk_markdown_blocks(
        "\n\n".join(body_parts), max_words, "img", str(path), src, 1,
    )


def _pdfplumber_text_only_pdf(path: Path, max_words: int, src: dict) -> list[Chunk]:
    """Text-only fallback: page.extract_text(), no tables, no scanned pages.

    This is the whole PDF pipeline when the Granite-Docling stack (torch,
    transformers, docling_core) isn't installed -- a real, heavy, optional
    dependency (see PAN-99/PAN-100). Kept intact rather than replaced so a
    deployment without it still gets real, if degraded, PDF support instead
    of losing the format outright.
    """
    import pdfplumber

    chunks: list[Chunk] = []
    try:
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if not text.strip():
                    continue
                words = text.split()
                if len(words) <= max_words:
                    body = text.strip()
                    chunks.append(Chunk(
                        chunk_id=_chunk_hash(body),
                        locator=f"p{page_num}",
                        body=body,
                        source_path=str(path),
                        source_type="pdf",
                        source_record=src,
                        word_count=len(words),
                        page_or_slide_number=page_num,
                    ))
                else:
                    sub = _split_words(text, max_words, f"p{page_num}",
                                       str(path), "pdf", src,
                                       page_or_slide_number=page_num)
                    chunks.extend(sub)
    except UnsupportedSourceError:
        raise
    except Exception as exc:
        raise _reject_source(
            path,
            "PDF_INVALID",
            f"The native PDF reader could not open {path.name}: {exc}.",
            capability_id="native_pdf",
            action="Repair the PDF or export it as searchable PDF or UTF-8 text.",
        ) from exc
    if not chunks:
        raise _reject_source(
            path,
            "OCR_REQUIRED",
            f"{path.name} contains no extractable text; PANTA will not pretend that image pixels were parsed.",
            capability_id="scanned_pdf_ocr",
        )
    return chunks


_GRANITE_DOCLING_MODEL_ID = "ibm-granite/granite-docling-258M"
_granite_docling_pipeline_cache: dict[str, Any] = {}
_granite_docling_availability_cache: bool | None = None


def _granite_docling_available() -> bool:
    """Whether torch/transformers/docling_core are importable, checked once.

    Re-running `import torch` on every parse_pdf call (rather than once) hit
    a real, reproducible torch-internal crash on the second call in the same
    process: `RuntimeError: function '_has_torch_function' already has a
    docstring` inside torch/overrides.py. Likely a torch re-init quirk
    triggered by module-cache interactions in the test suite, not something
    to guess a workaround for -- caching the check once, like the model
    pipeline itself is already cached, avoids the repeat import entirely.
    """
    global _granite_docling_availability_cache
    if _granite_docling_availability_cache is None:
        try:
            import docling_core  # noqa: F401
            import torch  # noqa: F401
            import transformers  # noqa: F401
            _granite_docling_availability_cache = True
        except Exception as exc:  # noqa: BLE001
            # Not just ImportError: importing transformers reads package
            # metadata off disk, and with this venv inside an iCloud-synced
            # folder that raised OSError(89, "Operation canceled") on an
            # evicted file -- which escaped an ImportError-only guard and
            # crashed extraction instead of taking the plain-text fallback
            # this function exists to select. Any failure to load the model
            # stack means the model is unavailable; that is the whole
            # question being asked here, so any exception answers it.
            if not isinstance(exc, ImportError):
                print(f"  [PDF] model stack unavailable ({type(exc).__name__}: {exc}); "
                      f"falling back to plain-text PDF extraction", file=sys.stderr)
            _granite_docling_availability_cache = False
    return _granite_docling_availability_cache


_TEXT_RESIDUAL_ENABLED = os.environ.get("PE_OS_TEXT_RESIDUAL", "1") != "0"
# '$' and ',' are dropped so "$50.4" and "50.4" test as the same token:
# a chart model writes the bare number where the page carries the currency.
_RESIDUAL_STRIP = re.compile(r"[^0-9a-z.%/+-]")


def _residual_norm(token: str) -> str:
    """Normalise a token for presence testing, not for display."""
    # Trailing punctuation is stripped after the class filter: '.' has to stay
    # in the keep-set for "35.3", but leaving it on the end made the emitted
    # "CAGR." fail to match the text layer's "CAGR" and report as missing.
    return _RESIDUAL_STRIP.sub("", token.lower()).strip(".")


_VALUE_STRIP = re.compile(r"[^0-9.%-]")


def _value_norm(token: str) -> str:
    """Reduce a token to its numeric identity: '$35.3' and '35.3' must match.

    Currency marks, commas and stray punctuation differ between what a chart
    model writes and how the same number sits in the text layer, and treating
    those as different values would fail every corroboration for cosmetic
    reasons.
    """
    return _VALUE_STRIP.sub("", token).strip(".-")


def _page_text_tokens(page: Any) -> list[str]:
    """The page's PDF text-layer tokens, or [] if it cannot be read."""
    try:
        return [(w.get("text") or "").strip() for w in page.extract_words()]
    except Exception:
        return []                      # never fail a page over a coverage check


def _chart_structure_warnings(payload: str) -> list[str]:
    """Signs the chart table is garbled even though its numbers are real.

    Corroborating values says nothing about their arrangement. On Goldman's
    page 16 a stacked bar chart came back with its stack components promoted
    to column headers -- "Year | Total | $1.2 | $1.3" -- and a body that was
    mostly empty. Every number was genuinely on the page, so the value check
    passed and printed a word that reads as general reassurance.

    These two signals are deliberately crude. They are meant to catch a table
    that has visibly collapsed, not to grade layout quality.
    """
    rows = []
    for line in payload.splitlines():
        if "|" not in line:
            continue
        raw = line.strip()
        cells = [c.strip() for c in raw.split("|")]
        # Strip ONLY the phantom cells created by a markdown row's edge pipes.
        # Dropping every trailing blank instead would erase the empty cells
        # that are the evidence of a collapsed table -- the check would then
        # under-report exactly the rows it exists to catch.
        if raw.startswith("|"):
            cells = cells[1:]
        if raw.endswith("|") and cells:
            cells = cells[:-1]
        if not cells or all(set(c) <= set("-: ") for c in cells if c):
            continue                   # markdown separator row, not data
        rows.append(cells)
    if len(rows) < 2:
        return []

    warnings: list[str] = []
    header, body = rows[0], rows[1:]
    numeric_header = [c for c in header if c and any(ch.isdigit() for ch in c)]
    if header and len(numeric_header) >= max(2, len(header) // 3):
        warnings.append(
            f"{len(numeric_header)} of {len(header)} header cells hold numeric "
            "values rather than labels")
    flat = [c for row in body for c in row]
    empty = [c for c in flat if not c]
    if flat and len(empty) * 2 >= len(flat):
        warnings.append(f"{len(empty)} of {len(flat)} body cells are empty")
    return warnings


def _chart_corroboration(page: Any, markdown: str) -> list[str]:
    """Check chart-recognition numbers against the page's own text layer.

    PAN-100 found chart recognition produces "structurally confident but
    factually wrong" values, which is why it is opt-in. The finding stands,
    but it is now testable: in a vector PDF the chart's data labels are real
    text objects, so any number the model reports should already be present
    as read text. A value that is NOT in the text layer was not read off the
    page -- it was invented, and that is exactly PAN-100's failure caught
    mechanically instead of trusted.

    Corroboration raises confidence in the VALUES only. Which value belongs
    to which bar stays model-asserted: two numbers can both be real and still
    be mapped to the wrong bars. So this never promotes the output past a
    proposal for human confirmation, and under invariant 3 it cannot be
    `derived` -- "a model looked at pixels" is not an inspectable derivation.
    """
    # startswith, not "in": the IMAGE_NOT_EXTRACTED marker mentions
    # "[chart-recognition output follows]" and is not itself chart output.
    blocks = [b for b in markdown.split("\n\n")
              if b.lstrip().startswith("[chart-recognition,")]
    if not blocks:
        return []
    known = {_value_norm(t) for t in _page_text_tokens(page)}
    known.discard("")
    if not known:
        return []                      # no text layer: nothing to check against

    notes: list[str] = []
    for block in blocks:
        # Drop the header line: it carries bbox=[174, 339, ...], whose pixel
        # coordinates are metadata about where the chart is, not values read
        # off it. Validating them against the text layer flags every chart.
        payload = "\n".join(block.splitlines()[1:])
        claimed, seen = [], set()
        for token in re.findall(r"[^\s|]+", payload):
            key = _value_norm(token)
            if not key or not any(c.isdigit() for c in key) or key in seen:
                continue
            seen.add(key)
            claimed.append((token, key))
        if not claimed:
            continue
        missing = [tok for tok, key in claimed if key not in known]
        if missing:
            notes.append(
                f"[validation] UNCORROBORATED: {len(missing)} of {len(claimed)} value(s) "
                f"in the chart-recognition output above do NOT appear in this page's PDF "
                f"text layer: {' \u00b7 '.join(missing)}. Nothing on the page reads that "
                "way, so these were not read -- they were inferred. This is the PAN-100 "
                "failure mode; treat the whole chart block as unreliable."
            )
        else:
            defects = _chart_structure_warnings(payload)
            if defects:
                notes.append(
                    f"[validation] VALUES CORROBORATED, STRUCTURE SUSPECT: all "
                    f"{len(claimed)} value(s) appear in this page's PDF text layer, so "
                    f"none were invented -- but the table above does not read like a "
                    f"chart: {'; '.join(defects)}. The numbers are real; their "
                    "arrangement is not. Do not take any pairing from this block "
                    "without reading the chart region itself."
                )
            else:
                notes.append(
                    f"[validation] CORROBORATED: all {len(claimed)} value(s) in the "
                    "chart-recognition output above also appear in this page's PDF text "
                    "layer as read text, so none were invented. The MAPPING of value to "
                    "chart element remains model-asserted and still needs human "
                    "confirmation."
                )
    return notes


def _text_layer_residual(page: Any, emitted: str) -> list[str]:
    """Text objects in the page's PDF text layer that the model did not emit.

    The layout models classify a chart region as a picture and drop the text
    inside it -- but in a vector PDF those axis labels and data values are
    real text objects, not pixels. On a Goldman investor-day page the entire
    loss was $50.4, $35.3, 2017-2019, 2020-2022: numbers sitting in the file,
    discarded because they fell inside a box labelled "chart".

    Recovering them by bbox was the obvious approach and it does not work.
    Page rotation (270 deg on that deck) leaves pdfplumber reporting word
    boxes in unrotated space, degenerate -- five tokens all claiming the same
    0.8pt-tall line -- so no scale maps the layout model's box onto them. A
    set difference needs no geometry at all, and it covers Granite and
    classic Docling too, neither of which reports a box.

    This is READ TEXT, epistemically `observed`. What it deliberately does
    NOT do is say which bar a value belongs to: that association is the part
    a model would have to infer, and inferring it here is exactly the
    confidently-wrong failure PAN-100 recorded.
    """
    if not _TEXT_RESIDUAL_ENABLED:
        return []
    seen = {_residual_norm(t) for t in re.findall(r"\S+", emitted)}
    seen.discard("")
    residual, reported = [], set()
    for text in _page_text_tokens(page):
        key = _residual_norm(text)
        if not key or key in seen or key in reported:
            continue
        reported.add(key)
        residual.append(text)
    return residual


def _granite_device() -> str:
    """Pick the device for Granite-Docling: CUDA if present, else CPU.

    MPS is deliberately never returned. PAN-99 found generate() hangs
    indefinitely on Apple GPU for this model, and a hang is worse than a
    slow CPU run because it never surfaces as an error. Set
    PE_OS_GRANITE_DEVICE to force a choice (including "mps", if someone
    later wants to retest that).
    """
    forced = os.environ.get("PE_OS_GRANITE_DEVICE", "").strip()
    if forced:
        return forced
    try:
        import torch
        if torch.cuda.is_available():
            return "cuda"
    except Exception:
        pass
    return "cpu"


def _granite_docling_pipeline() -> tuple[Any, Any, Any]:
    """Lazily load and cache the Granite-Docling processor/model/torch module.

    Loaded once per process, not once per page or per document -- load time
    is a few seconds and is wasted if repeated. `use_cache=True` on generate()
    is part of the confirmed-working configuration, not a guess.

    Device is resolved by _granite_device(): CUDA when present, CPU
    otherwise, and never MPS -- Apple GPU was found to hang indefinitely on
    generate() for this model (PAN-99 research). CPU+float32 was the local
    Mac answer to that hang; it is not a property of the model, so pinning
    it here would make a GPU deployment buy nothing for this engine.
    """
    if "model" in _granite_docling_pipeline_cache:
        cached = _granite_docling_pipeline_cache
        return cached["torch"], cached["processor"], cached["model"]
    import torch
    from transformers import AutoModelForImageTextToText, AutoProcessor

    device = _granite_device()
    # bfloat16 on CUDA: this is a 258M model, so the win is bandwidth, not
    # capacity. float32 stays the CPU default because CPU bf16 is slower,
    # not faster. Override either with PE_OS_GRANITE_DTYPE.
    dtype_name = os.environ.get(
        "PE_OS_GRANITE_DTYPE", "bfloat16" if device == "cuda" else "float32")
    dtype = getattr(torch, dtype_name)

    processor = AutoProcessor.from_pretrained(_GRANITE_DOCLING_MODEL_ID)
    model = AutoModelForImageTextToText.from_pretrained(
        _GRANITE_DOCLING_MODEL_ID, dtype=dtype, _attn_implementation="sdpa",
    ).to(device)
    model.eval()
    print(f"  [PDF] Granite-Docling on {device} ({dtype_name})", file=sys.stderr)
    _granite_docling_pipeline_cache.update(
        torch=torch, processor=processor, model=model, device=device)
    return torch, processor, model


def _granite_docling_convert_page(image: Any, max_new_tokens: int = 2048) -> tuple[str, list[str]]:
    """Convert one rendered page image to markdown plus a list of picture
    classifications (e.g. "bar_chart") found on the page.

    Uses docling_core's own DoclingDocument.pictures[].meta.classification
    rather than regexing the raw doctags -- a real, structured field, not a
    fragile string match.
    """
    from docling_core.types.doc import DoclingDocument
    from docling_core.types.doc.document import DocTagsDocument

    torch, processor, model = _granite_docling_pipeline()
    messages = [{"role": "user", "content": [{"type": "image"}, {"type": "text", "text": "Convert this page to docling."}]}]
    prompt = processor.apply_chat_template(messages, add_generation_prompt=True)
    device = _granite_docling_pipeline_cache.get("device", "cpu")
    inputs = processor(text=prompt, images=[image], return_tensors="pt").to(device)
    with torch.no_grad():
        generated_ids = model.generate(**inputs, max_new_tokens=max_new_tokens, do_sample=False, use_cache=True)
    trimmed_ids = generated_ids[:, inputs.input_ids.shape[1]:]
    doctags = processor.batch_decode(trimmed_ids, skip_special_tokens=False)[0].lstrip()

    doctags_doc = DocTagsDocument.from_doctags_and_image_pairs([doctags], [image])
    doc = DoclingDocument.load_from_doctags(doctags_doc, document_name="page")
    markdown = doc.export_to_markdown()

    picture_classes: list[str] = []
    for picture in doc.pictures:
        classification = getattr(picture.meta, "classification", None)
        predictions = getattr(classification, "predictions", None) if classification else None
        if predictions:
            picture_classes.append(predictions[0].class_name)
        else:
            picture_classes.append("unclassified")
    return markdown, picture_classes


_DEGENERATE_REPETITION_THRESHOLD = 6


def _is_degenerate_repetition(markdown: str) -> bool:
    """True if the same non-blank line repeats often enough to be a
    generation repetition-loop rather than real page content.

    Confirmed via direct test on a heavily degraded (skewed, ~70dpi
    effective, noisy) scan: Granite-Docling ran to completion with no
    exception and produced a short phrase ~80 times, recovering none of
    the page's real content. Genuine document text essentially never
    repeats an identical line this many times within one page.
    """
    lines = [line.strip() for line in markdown.splitlines() if line.strip()]
    if not lines:
        return False
    counts: dict[str, int] = {}
    for line in lines:
        counts[line] = counts.get(line, 0) + 1
    return max(counts.values()) >= _DEGENERATE_REPETITION_THRESHOLD


def _chunk_markdown_blocks(markdown: str, max_words: int, locator_prefix: str,
                            source_path: str, source_record: dict,
                            page_or_slide_number: int) -> list[Chunk]:
    """Group markdown into word-bounded chunks without ever splitting inside
    a block (a table, a paragraph, a heading) -- a naive word-count
    rejoin-and-resplit (as _split_words does) would flatten a markdown
    table's newlines into a single run-on line, destroying the row/column
    structure this whole path exists to produce (the same real failure
    mode found for PPTX shape text -- see PAN-101)."""
    blocks = [b for b in re.split(r"\n\s*\n", markdown) if b.strip()]
    chunks: list[Chunk] = []
    pending: list[str] = []
    pending_words = 0
    start_block = 0

    def flush(end_block: int) -> None:
        nonlocal pending, pending_words
        if not pending:
            return
        body = "\n\n".join(pending)
        locator = f"{locator_prefix}:b{start_block}-{end_block}"
        chunks.append(Chunk(
            chunk_id=_chunk_hash(body),
            locator=locator,
            body=body,
            source_path=source_path,
            source_type="pdf",
            source_record=source_record,
            word_count=len(body.split()),
            page_or_slide_number=page_or_slide_number,
        ))
        pending, pending_words = [], 0

    for i, block in enumerate(blocks):
        block_words = len(block.split())
        if pending and pending_words + block_words > max_words:
            flush(i)
            start_block = i
        pending.append(block)
        pending_words += block_words
    flush(len(blocks))

    if len(chunks) == 1:
        chunks[0].locator = locator_prefix
    return chunks


def _default_convert_page():
    """The page-image converter parse_pdf and parse_image fall back to when
    no caller supplies one: Granite-Docling if the model stack is installed,
    else None (meaning "no local model available")."""
    if not _granite_docling_available():
        return None
    return lambda image, page_num: _granite_docling_convert_page(image)


def parse_pdf(path: Path, max_words: int = CHUNK_WORDS,
              source_record: dict | None = None,
              convert_page: Any = None) -> list[Chunk]:
    """Convert each page via Granite-Docling-258M (real table structure,
    reads pixels directly so scanned pages need no separate OCR pass --
    verified against a real 9-page financial-narrative PDF: dense tables
    with distinct economic-basis EBITDA figures extracted with zero errors,
    section headings and callout boxes correctly separated from body text).
    Falls back to plain pdfplumber text extraction -- no tables, requires a
    real text layer -- when the model stack (torch/transformers/docling_core)
    isn't installed; this is a real, heavy, optional dependency, not assumed
    always present. A picture/chart on a page becomes a declared
    IMAGE_NOT_EXTRACTED-style marker naming its detected type, never
    invented data: a real waterfall/EBITDA-bridge chart was tested directly
    and produced a structurally confident but factually wrong table (real
    deltas swapped between categories) -- exactly the "confident wrong
    answer" this codebase's coverage-limit philosophy exists to prevent, so
    chart pixels are named as present, not blindly converted to numbers.
    """
    try:
        import pdfplumber
    except ImportError as exc:
        raise _reject_source(
            path,
            "READER_UNAVAILABLE",
            "Native PDF parsing is unavailable because the approved local reader is not installed.",
            capability_id="native_pdf",
        ) from exc
    src = source_record or _source_record(path)

    # `convert_page` swaps ONLY the page-image -> (markdown, picture classes)
    # step, so an alternative model is compared against Granite through the
    # identical chunking, locators, degenerate-output check and per-page
    # fallback. Anything else would compare two pipelines, not two models.
    # Default None keeps the production path exactly as it was.
    if convert_page is None:
        convert_page = _default_convert_page()
        if convert_page is None:
            return _pdfplumber_text_only_pdf(path, max_words, src)

    chunks: list[Chunk] = []
    try:
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    image = page.to_image(resolution=200).original
                    markdown, picture_classes = convert_page(image, page_num)
                    if _is_degenerate_repetition(markdown):
                        # Not an exception -- the model ran to completion and
                        # produced well-formed-looking markdown, but a badly
                        # degraded scan can send generation into a repetition
                        # loop (confirmed via direct test: a heavily skewed,
                        # low-DPI, noisy page produced the same short phrase
                        # ~80 times, recovering none of the page's real
                        # content). Nothing here raises, so without this
                        # check the garbage would flow through as a normal
                        # chunk with no coverage-limit signal at all.
                        raise RuntimeError("degenerate repetition loop in generated output")
                except Exception as exc:
                    # A single page's model failure must not lose the rest
                    # of the document -- fall back to plain text for just
                    # this page rather than aborting the whole PDF.
                    print(f"  [PDF] page {page_num}: model failed ({exc}), falling back to text", file=sys.stderr)
                    text = page.extract_text() or ""
                    if not text.strip():
                        continue
                    words = text.split()
                    if len(words) <= max_words:
                        body = text.strip()
                        chunks.append(Chunk(
                            chunk_id=_chunk_hash(body), locator=f"p{page_num}", body=body,
                            source_path=str(path), source_type="pdf", source_record=src,
                            word_count=len(words), page_or_slide_number=page_num,
                        ))
                    else:
                        chunks.extend(_split_words(
                            text, max_words, f"p{page_num}", str(path), "pdf", src,
                            page_or_slide_number=page_num,
                        ))
                    continue

                body_parts = [markdown.strip()] if markdown.strip() else []
                for picture_class in picture_classes:
                    body_parts.append(
                        f"[picture] IMAGE_NOT_EXTRACTED: a {picture_class} is present on this "
                        "page; its content was not reliably extracted. Chart pixel data is "
                        "not converted to numbers here -- a real waterfall/bridge chart test "
                        "showed this can produce confidently wrong values (PAN-100)."
                    )
                body_parts.extend(_chart_corroboration(page, "\n\n".join(body_parts)))
                residual = _text_layer_residual(page, "\n\n".join(body_parts))
                if residual:
                    shown = " \u00b7 ".join(residual[:40])
                    more = f" (+{len(residual) - 40} more)" if len(residual) > 40 else ""
                    numeric = [t for t in residual if any(c.isdigit() for c in t)]
                    # A page number is a stray token; nine numbers are a region
                    # the layout model never proposed, so nothing ever read it.
                    weight = (f" {len(numeric)} of them carry digits, which usually means a "
                              "data region (a chart, an axis, a callout) was not detected "
                              "at all rather than merely unread."
                              if len(numeric) >= 3 else "")
                    body_parts.append(
                        f"[coverage] TEXT_LAYER_RESIDUAL: {len(residual)} text object(s) "
                        f"in this page's PDF text layer do not appear in the extraction "
                        f"above: {shown}{more}. These are READ TEXT, not inferred.{weight} "
                        "Which chart element each belongs to is NOT resolved here."
                    )
                body = "\n\n".join(body_parts)
                if not body.strip():
                    continue
                chunks.extend(_chunk_markdown_blocks(
                    body, max_words, f"p{page_num}", str(path), src, page_num,
                ))
    except UnsupportedSourceError:
        raise
    except Exception as exc:
        raise _reject_source(
            path,
            "PDF_INVALID",
            f"The native PDF reader could not open {path.name}: {exc}.",
            capability_id="native_pdf",
            action="Repair the PDF or export it as searchable PDF or UTF-8 text.",
        ) from exc
    if not chunks:
        raise _reject_source(
            path,
            "OCR_REQUIRED",
            f"{path.name} contains no extractable content on any page.",
            capability_id="scanned_pdf_ocr",
        )
    return chunks


def _xlsx_merged_ranges(path: Path, sheet_name: str) -> list[tuple[int, int, int, int]]:
    """Read one worksheet's merge ranges directly from its XML part.

    openpyxl's read_only worksheets -- required here for real models with
    thousands of cells -- do not expose ``merged_cells`` at all, and a
    second non-read-only load of the same workbook just to read this would
    force a full in-memory parse of a file already loaded twice. Merge
    ranges are a handful of small, well-defined XML elements; read them
    directly instead. Best-effort: any failure returns no ranges rather
    than breaking the rest of the parse.
    """
    ns = {
        "m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    try:
        from openpyxl.utils import range_boundaries
        with zipfile.ZipFile(path) as archive:
            workbook_xml = ElementTree.fromstring(archive.read("xl/workbook.xml"))
            sheet_rid = next(
                (
                    sheet_el.get(f"{{{ns['r']}}}id")
                    for sheet_el in workbook_xml.findall("m:sheets/m:sheet", ns)
                    if sheet_el.get("name") == sheet_name
                ),
                None,
            )
            if sheet_rid is None:
                return []
            rels_xml = ElementTree.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
            target = next(
                (
                    rel.get("Target")
                    for rel in rels_xml
                    if rel.get("Id") == sheet_rid
                ),
                None,
            )
            if not target:
                return []
            # A relationship Target is either package-root-relative
            # ("/xl/worksheets/sheet2.xml", the common case for Excel-authored
            # files) or relative to the xl/ part itself ("worksheets/sheet2.xml").
            target = target.lstrip("/")
            sheet_path = target if target.startswith("xl/") else f"xl/{target}"
            sheet_xml = ElementTree.fromstring(archive.read(sheet_path))
            ranges = []
            for merge_cell in sheet_xml.findall("m:mergeCells/m:mergeCell", ns):
                ref = merge_cell.get("ref") or ""
                if ":" not in ref:
                    continue
                min_col, min_row, max_col, max_row = range_boundaries(ref)
                ranges.append((min_row, min_col, max_row, max_col))
            return ranges
    except (OSError, KeyError, zipfile.BadZipFile, ElementTree.ParseError, ValueError, ImportError):
        return []


def parse_xlsx(path: Path, max_words: int = CHUNK_WORDS,
               source_record: dict | None = None) -> list[Chunk]:
    """Create reproducible, cell-addressable chunks from an Excel workbook.

    A workbook is not prose: formulas and their cached outputs carry different
    meanings.  The chunk body keeps both, while the locator names the exact
    sheet and cell range so a reviewer can verify any extracted claim.
    """
    try:
        import openpyxl
        from openpyxl.utils import get_column_letter
    except ImportError as exc:
        raise _reject_source(
            path,
            "READER_UNAVAILABLE",
            "Open XML workbook parsing is unavailable because openpyxl is not installed.",
            capability_id="openxml_workbook",
        ) from exc
    try:
        formulas = openpyxl.load_workbook(path, data_only=False, read_only=True)
        values = openpyxl.load_workbook(path, data_only=True, read_only=True)
    except Exception as exc:
        raise _reject_source(
            path,
            "WORKBOOK_INVALID",
            f"Cannot read workbook {path.name}: {exc}.",
            capability_id="openxml_workbook",
            action="Open the workbook in Excel and save a valid .xlsx/.xlsm copy without flattening formulas.",
        ) from exc
    src = source_record or _source_record(path)
    chunks: list[Chunk] = []
    workbook_has_formulas = any(
        isinstance(cell.value, str) and cell.value.startswith("=")
        for ws in formulas.worksheets for row in ws.iter_rows() for cell in row
    )
    effective_max_words = max(max_words, 1200) if workbook_has_formulas else max_words
    for sheet_name in formulas.sheetnames:
        ws, value_ws = formulas[sheet_name], values[sheet_name]
        # Models can contain thousands of raw input/output rows.  Their
        # complete deterministic cell graph is captured separately; L2
        # should see only the formula-bearing rows needed to name the
        # economics, not every cell -- but only on sheets that actually
        # compute (a dedicated Inputs sheet, by definition, has no formulas
        # of its own; gating this per-workbook instead of per-sheet dropped
        # every raw-value sheet whenever any other sheet in the file had a
        # formula).
        sheet_has_formulas = any(
            isinstance(cell.value, str) and cell.value.startswith("=")
            for row in ws.iter_rows() for cell in row
        )
        # Only the top-left cell of a merged range carries a value in
        # openpyxl; a title or section header merged across a row would
        # otherwise vanish everywhere except its first column. The
        # read_only worksheet needed for real models has no random cell
        # access and no merged_cells at all, so ranges come from
        # _xlsx_merged_ranges (direct XML) and each top-left value is
        # captured in sequential order as that row is actually reached.
        # Formula strings are deliberately excluded from propagation -- a
        # merged formula cell's cached value belongs to its own
        # coordinate, not to whichever coordinate propagation reaches, so
        # propagating it would print a wrong "cached=" pairing. Titles and
        # section labels are the real, common case and are always plain
        # text.
        merges_by_start_row: dict[int, list[tuple[int, int, int, int]]] = {}
        for (min_row, min_col, max_row, max_col) in _xlsx_merged_ranges(path, sheet_name):
            merges_by_start_row.setdefault(min_row, []).append((min_row, min_col, max_row, max_col))
        merged_value_by_coord: dict[tuple[int, int], Any] = {}
        # Running per-column header state (PAN-102): a row that is mostly
        # or entirely text with no formula of its own is treated as a
        # header/title row, and its text becomes the header context
        # attached to data cells below it, in that column, until a later
        # header-like row replaces it.  Without this, a bare coordinate
        # like "C15=42.3" carries no period or line-item context once it
        # leaves its own chunk.
        column_header: dict[int, str] = {}
        pending: list[str] = []
        start_row = end_row = None
        def flush() -> None:
            nonlocal pending, start_row, end_row
            if not pending or start_row is None or end_row is None:
                return
            body = f"Workbook: {path.name}\nSheet: {sheet_name}\n" + "\n".join(pending)
            chunks.append(Chunk(
                chunk_id=_chunk_hash(body),
                locator=f"{path.name}::{sheet_name}!{start_row}:{end_row}",
                body=body, source_path=str(path), source_type="xlsx",
                source_record=src, word_count=len(body.split()),
                section_heading=sheet_name,
            ))
            pending, start_row, end_row = [], None, None
        for row_number, (row, cached_row) in enumerate(zip(ws.iter_rows(), value_ws.iter_rows()), start=1):
            # Read-only cells with no value are lightweight EmptyCell
            # placeholders that carry no .row/.column/.coordinate at all --
            # position within the row tuple (1-indexed) is used instead of
            # trusting cell attributes, since both cell types support it.
            for (min_row, min_col, max_row, max_col) in merges_by_start_row.get(row_number, []):
                top_left_value = row[min_col - 1].value if min_col - 1 < len(row) else None
                if not isinstance(top_left_value, str) or top_left_value.startswith("="):
                    continue
                for r in range(min_row, max_row + 1):
                    for c in range(min_col, max_col + 1):
                        merged_value_by_coord[(r, c)] = top_left_value

            row_entries: list[tuple[Any, Any, Any, int]] = []
            for col_index, (cell, cached_cell) in enumerate(zip(row, cached_row), start=1):
                raw = cell.value
                if raw is None:
                    merged = merged_value_by_coord.get((row_number, col_index))
                    if merged is not None:
                        raw = merged
                if raw is None:
                    continue
                row_entries.append((cell, cached_cell, raw, col_index))
            if not row_entries:
                continue

            has_formula_in_row = any(
                isinstance(raw, str) and raw.startswith("=") for _, _, raw, _ in row_entries
            )
            text_entries = [
                (col_index, raw) for _, _, raw, col_index in row_entries
                if isinstance(raw, str) and not raw.startswith("=")
            ]
            # A period-header row ("Line Item | Unit | Q2'26 | Q3'26 | ...")
            # mixes text labels with date values across most of its width --
            # it must count as header-like too, not just a uniform-text row,
            # or it fails the "mostly text" threshold on its own date cells
            # and gets read as a data row instead.
            date_entries = [
                (col_index, raw) for _, _, raw, col_index in row_entries
                if isinstance(raw, (datetime, date)) and not isinstance(raw, bool)
            ]
            numeric_entries = [
                (col_index, raw) for _, _, raw, col_index in row_entries
                if isinstance(raw, (int, float)) and not isinstance(raw, bool)
            ]
            header_like_entries = text_entries + date_entries

            # A bare numeric row with no text/date label of its own and no
            # formula, on a sheet that otherwise computes, is exactly the
            # "thousands of raw rows" case the original filter targeted --
            # still skipped. But a row carrying ANY text or date (a title, a
            # section header, a period-header row, or a label next to a
            # raw historical figure) is never noise and must never be
            # silently dropped just because it isn't itself a formula.
            if sheet_has_formulas and not has_formula_in_row and not header_like_entries:
                continue

            is_header_like = (
                not has_formula_in_row
                and not numeric_entries
                and len(header_like_entries) >= 2
                and len(header_like_entries) >= len(row_entries) - 1
            )
            if is_header_like:
                for col_index, raw in header_like_entries:
                    column_header[col_index] = raw
            row_label = text_entries[0][1] if text_entries else None

            cells: list[str] = []
            for cell, cached_cell, raw, col_index in row_entries:
                coordinate = f"{get_column_letter(col_index)}{row_number}"
                is_formula = isinstance(raw, str) and raw.startswith("=")
                if is_formula:
                    piece = f"{coordinate}=FORMULA({raw}); cached={cached_cell.value!r}"
                else:
                    piece = f"{coordinate}={raw}"
                if not is_header_like:
                    context_bits = []
                    if row_label is not None and row_label != raw:
                        context_bits.append(str(row_label))
                    header = column_header.get(col_index)
                    if header is not None and header != raw:
                        context_bits.append(str(header))
                    role = _xlsx_cell_semantic_role(cell, is_formula)
                    if role:
                        context_bits.append(f"role={role}")
                    if context_bits:
                        piece += f" [{'; '.join(context_bits)}]"
                cells.append(piece)

            line = " | ".join(cells)
            projected = len((" ".join(pending + [line])).split())
            if pending and projected > effective_max_words:
                flush()
            if start_row is None:
                start_row = row_number
            end_row = row_number
            pending.append(line)
        flush()
    return chunks


def _xlsx_cell_semantic_role(cell: Any, is_formula: bool) -> str | None:
    """Classify a cell by the standard IB/LBO-model font-color convention.

    Blue = hardcoded input, black/default = an ordinary same-sheet formula,
    green = a formula that links to another sheet.  Confirmed against a
    real Keystone LBO model's actual color usage (not assumed from the
    convention alone) before trusting it: this file uses exactly
    0000FF/000000/008000 for input/formula/cross-sheet-link respectively.
    Any other color is left unclassified rather than guessed -- a wrong
    label is worse than no label.
    """
    color = cell.font.color if cell.font else None
    rgb = getattr(color, "rgb", None)
    if not isinstance(rgb, str) or len(rgb) < 6:
        return None
    hex6 = rgb[-6:].upper()
    if is_formula:
        return "cross_sheet_link" if hex6 == "008000" else None
    return "input" if hex6 == "0000FF" else None


def parse_source(path: Path, max_words: int = CHUNK_WORDS,
                 source_record: dict | None = None,
                 convert_page: Any = None,
                 vision_fallback: Any = None,
                 pdf_engine: Any = None) -> list[Chunk]:
    src = source_record or _source_record(path)
    envelope = src.get("source_envelope") or {}
    capability = resolve_source_capability(path, {
        "document_type": envelope.get("document_type") or src.get("doc_type"),
        "parser_route": envelope.get("parser_route"),
    })
    capability_id = str(capability["capability_id"])
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        chunks = parse_pdf(path, max_words, src, convert_page=convert_page)
    elif suffix in IMAGE_SUFFIXES:
        chunks = parse_image(path, max_words, src, convert_page=convert_page,
                             vision_fallback=vision_fallback)
    elif suffix in (".md", ".markdown"):
        chunks = parse_markdown(path, max_words, src)
    elif suffix == ".txt" and capability_id == "transcript_export":
        chunks = parse_transcript(path, max_words, src)
    elif suffix == ".txt":
        chunks = parse_plain_text(path, max_words, src)
    elif suffix in (".html", ".htm"):
        chunks = parse_html(path, max_words, src)
    elif suffix == ".csv":
        chunks = parse_csv(path, max_words, src)
    elif suffix == ".docx":
        chunks = parse_docx(path, max_words, src)
    elif suffix == ".pptx":
        chunks = parse_pptx(path, max_words, src, pdf_engine=pdf_engine)
    elif suffix in (".srt", ".vtt"):
        chunks = parse_transcript(path, max_words, src)
    elif suffix in (".eml", ".mbox"):
        chunks = parse_email(path, max_words, src)
    elif suffix in (".xlsx", ".xlsm"):
        chunks = parse_xlsx(path, max_words, src)
    elif suffix == ".xls":
        raise _reject_source(
            path,
            "LEGACY_EXCEL_UNSUPPORTED",
            "Legacy .xls is not supported; convert the workbook to .xlsx before extraction.",
            capability_id="legacy_excel",
        )
    elif suffix == ".msg":
        raise _reject_source(
            path,
            "OUTLOOK_MSG_UNSUPPORTED",
            "Binary Outlook .msg is not parsed because doing so would require an unapproved decoder.",
            capability_id="outlook_msg",
        )
    else:
        raise _reject_source(
            path,
            "UNSUPPORTED_SOURCE",
            f"Unsupported source type: {suffix or '<no extension>'}.",
        )
    if not chunks:
        raise _reject_source(
            path,
            "NO_EXTRACTABLE_TEXT",
            f"{path.name} contains no safely extractable text.",
            capability_id=capability_id,
            action="Verify the export contains visible text and upload it again; attachments require their own SourceEnvelope.",
        )
    return _decorate_chunks(chunks, path, src, capability_id)


def load_manifest(
    manifest: str,
    deal: str,
    source_dir: Path = VAULT_INBOX,
) -> list[Path]:
    """Return ordered list of source files for the given manifest."""
    keys = MANIFEST_SOURCES.get(manifest, [])
    paths = []
    for key in keys:
        # Try stem match in vault/inbox
        for ext in (
            ".md", ".markdown", ".txt", ".html", ".htm", ".csv", ".pdf",
            ".docx", ".pptx", ".xlsx", ".xlsm", ".srt", ".vtt", ".eml", ".mbox",
        ):
            candidate = source_dir / f"{key}{ext}"
            if candidate.exists():
                paths.append(candidate)
                break
    return paths


# ─────────────────────────────────────────────────────────────────────────────
# L2 — LLM annotator (schema-constrained via tool_use)
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class RawClaim:
    metric: str
    value: Any
    unit: str | None
    period: str | None
    perimeter: str | None
    epistemic_class: str
    direction: str
    topic: str
    definition_id: str | None
    statement: str
    locator: str
    source_id: str
    source_path: str
    known_at: str
    derivation: str | None = None
    author: str | None = None
    # Identity dimensions. Defaulted so a cached raw_claims run predating these
    # fields still loads — it simply scores as unresolvable until re-extracted.
    entity: str = "unspecified"
    period_canonical: str = "none"
    scope: str = "unspecified"
    # Which slice of the quantity. "total" means the whole; blank collides a
    # component with its own total, which reads as a contradiction.
    measurement: str = "total"
    claim_kind: str = "QUANTITATIVE"
    bound: str = "EXACT"
    basis: str = "unspecified"
    scenario: str = "unspecified"
    source_version_id: str | None = None
    # The real concept name when `metric` is "Other" — the enum's escape hatch.
    # Defaulted so a cached raw_claims run predating the field still loads.
    metric_label: str | None = None
    # Which annotate_chunk() call this claim came from, and its position in
    # that call's own claims array -- lets a derived claim reference another
    # claim from the SAME tool call before either has a real claim_id (that
    # only exists after validate() runs). See derive_relations().
    chunk_id: str = ""
    batch_index: int = -1
    derivation_operand_batch_indices: list[int] = field(default_factory=list)
    alternative_to_batch_index: int | None = None
    # Position on this archetype's evidence ladder. See archetype_pack.
    evidence_state: str | None = None


def _is_fatal_provider_error(exc: Exception) -> bool:
    """Return True when retrying more chunks with the same key cannot succeed."""
    status_code = getattr(exc, "status_code", None)
    message = str(exc).lower()
    # OpenRouter uses HTTP 402 both for an exhausted account and for its
    # temporary in-flight spending cap.  The latter explicitly asks callers to
    # wait for outstanding requests to settle, so treating every 402 as fatal
    # throws away an otherwise resumable extraction batch.
    if status_code == 402 and "in_flight_budget_exhausted" in message:
        return False
    if status_code in {401, 402, 403}:
        return True
    return any(
        marker in message
        for marker in (
            "billing_error",
            "permission_error",
            "key limit exceeded",
            "invalid api key",
        )
    )


def _provider_retry_delay(exc: Exception, attempt: int) -> float | None:
    """Return a bounded retry delay for transient provider failures."""
    status_code = getattr(exc, "status_code", None)
    message = str(exc)
    transient_in_flight_cap = (
        status_code == 402
        and "in_flight_budget_exhausted" in message.lower()
    )
    if status_code not in {408, 409, 429, 500, 502, 503, 504, 529} and not transient_in_flight_cap:
        return None

    response = getattr(exc, "response", None)
    headers = getattr(response, "headers", {}) if response is not None else {}
    raw_delay = headers.get("Retry-After") if headers else None
    if raw_delay is None:
        match = re.search(
            r"retry_after_seconds(?:_raw)?['\"\s:]+([0-9]+(?:\.[0-9]+)?)",
            message,
            re.IGNORECASE,
        )
        if match is None:
            match = re.search(
                r"retry-after['\"\s:]+['\"]?([0-9]+(?:\.[0-9]+)?)",
                message,
                re.IGNORECASE,
            )
        raw_delay = match.group(1) if match else None
    try:
        delay = float(raw_delay) if raw_delay is not None else 2 ** attempt
    except (TypeError, ValueError):
        delay = 2 ** attempt
    return max(1.0, min(delay, 30.0))


def annotate_chunk(
    chunk: Chunk,
    client,
    deal: str,
    rate_limit_delay: float = 0.25,
    *,
    raise_errors: bool = False,
    archetype: str = DEFAULT_ARCHETYPE,
) -> list[RawClaim]:
    src = chunk.source_record
    prompt = (
        f"DEAL: {deal}\n"
        f"SOURCE: {src['source_id']} ({src['name']}) — {src['doc_type']}\n"
        f"SOURCE PARTY: {src.get('party') or 'not available'}\n"
        f"EFFECTIVE DATE: {src.get('effective_date') or 'not available'}\n"
        f"KNOWN AT: {src['known_at']}\n"
        f"FRAGMENT LOCATOR: {chunk.locator}\n\n"
        f"SECTION HEADING: {chunk.section_heading or 'not available'}\n"
        f"PAGE OR SLIDE: {chunk.page_or_slide_number or 'not available'}\n\n"
        f"{chunk.body}"
    )
    try:
        request = {
            "model": MODEL,
            # Match the larger schema capacity: Excel chunks commonly contain
            # more than four claim-bearing rows and need room for tool JSON.
            "max_tokens": MAX_TOKENS,
            "system": SYSTEM_PROMPT,
            "tools": [claim_tool_for(archetype)],
            "messages": [{"role": "user", "content": prompt}],
        }
        # Sampling temperature. SDK 1.0.0 dropped `temperature` from the typed
        # signature of messages.create, but the API still honours it, and the SDK
        # forwards extra_body verbatim — so this is the only route to it here.
        #
        # It is not a nicety. Measured on one 92-word Keystone chunk at the API
        # default of 1.0: the claim count was stable at 11 across four runs, yet
        # all four identity fingerprints differed and not one of the 33
        # identities recurred in every run — the same EBITDA was FirmView twice
        # and something else the other two times.
        #
        # Extraction varying in *what* it finds would be visible. Varying in how
        # it *classifies* what it found is invisible and worse: the same quantity
        # lands under different identities, so grouping and contradiction
        # detection quietly stop working.
        extra_body = {"temperature": 0}
        provider_extra = openrouter_extra_body()
        if provider_extra:
            extra_body.update(provider_extra)
        request["extra_body"] = extra_body
        # A reasoning model will otherwise spend the entire output budget
        # thinking and never emit the tool call -- see thinking_parameter().
        thinking = thinking_parameter()
        if thinking:
            request["thinking"] = thinking
        # Pinning the tool makes GLM generate a tool call that never
        # terminates -- see forces_tool_choice().
        if forces_tool_choice():
            request["tool_choice"] = {"type": "tool", "name": "emit_claims"}
        resp = client.messages.create(**request)
        time.sleep(rate_limit_delay)
    except Exception as e:
        if raise_errors:
            raise
        print(f"  [L2 ERROR] {chunk.chunk_id}: {e}", file=sys.stderr)
        return []

    raw_claims: list[RawClaim] = []
    for block in resp.content:
        if block.type == "tool_use" and block.name == "emit_claims":
            batch = block.input.get("claims", [])
            batch_size = len(batch)
            for batch_index, c in enumerate(batch):
                # Indices the model gave must point at another claim IN this
                # same batch -- not itself, not out of range. A dangling or
                # self-referencing index is dropped, not fabricated into an
                # edge to something arbitrary.
                operand_indices = [
                    idx for idx in (c.get("derivation_operand_indices") or [])
                    if isinstance(idx, int) and 0 <= idx < batch_size and idx != batch_index
                ]
                alt_idx = c.get("alternative_to_index")
                if not (isinstance(alt_idx, int) and 0 <= alt_idx < batch_size and alt_idx != batch_index):
                    alt_idx = None
                period = _non_empty_l2_text(c.get("period"))
                if period is None:
                    effective_date = _non_empty_l2_text(src.get("effective_date"))
                    period = f"as of {effective_date}" if effective_date else "unknown"

                # There is no reliable economic-scope inference when the model
                # omitted the perimeter. Preserve that absence explicitly rather
                # than inventing an entity, consolidation basis, or adjustment view.
                perimeter = _non_empty_l2_text(c.get("perimeter")) or "unknown"

                hint = _non_empty_l2_text(c.get("locator_hint"))
                if hint is None:
                    hint = _non_empty_l2_text(chunk.section_heading)
                if hint is None and chunk.page_or_slide_number is not None:
                    hint = f"page or slide {chunk.page_or_slide_number}"
                if hint is None:
                    # The fragment locator is deterministic L1 metadata. Reusing it
                    # does not add a fabricated location and the branch below avoids
                    # appending the locator to itself.
                    hint = chunk.locator
                # Resolve to the section this claim actually came from when
                # the hint names one; otherwise keep the chunk locator, with
                # the hint appended only when it adds something (the model
                # sometimes echoes the locator it was given, which produced
                # "file.md::## Heading:file.md::## Heading").
                locator = _claim_locator(chunk, hint)
                raw_claims.append(RawClaim(
                    metric=c["metric"],
                    metric_label=_non_empty_l2_text(c.get("metric_label")),
                    value=c.get("value"),
                    unit=c.get("unit"),
                    period=period,
                    perimeter=perimeter,
                    epistemic_class=c.get("epistemic_class", "asserted"),
                    direction=c.get("direction", "context"),
                    topic=c.get("topic", "OTHER"),
                    definition_id=c.get("definition_id"),
                    statement=c.get("statement", ""),
                    locator=locator,
                    source_id=src["source_id"],
                    source_path=chunk.source_path,
                    known_at=src["known_at"],
                    source_version_id=(
                        (src.get("source_envelope") or {}).get("source_version_id")
                    ),
                    derivation=c.get("derivation"),
                    author=c.get("author"),
                    entity=c.get("entity") or "unspecified",
                    period_canonical=c.get("period_canonical") or "none",
                    scope=c.get("scope") or "unspecified",
                    measurement=c.get("measurement") or "total",
                    claim_kind=c.get("claim_kind") or "QUANTITATIVE",
                    bound=c.get("bound") or "EXACT",
                    basis=c.get("basis") or "unspecified",
                    scenario=c.get("scenario") or "unspecified",
                    chunk_id=chunk.chunk_id,
                    batch_index=batch_index,
                    derivation_operand_batch_indices=operand_indices,
                    alternative_to_batch_index=alt_idx,
                    evidence_state=_non_empty_l2_text(c.get("evidence_state")),
                ))
    return raw_claims


# ─────────────────────────────────────────────────────────────────────────────
# L3 — Validator / normalizer (deterministic)
# ─────────────────────────────────────────────────────────────────────────────

def _parse_float(v: Any) -> float | None:
    if v is None or v == "":
        return None
    try:
        value = float(str(v).replace(",", "").replace("%", "").strip())
        return value if math.isfinite(value) else None
    except (TypeError, ValueError):
        return None


def _claim_locator(chunk: Chunk, hint: str | None) -> str:
    """Where THIS claim came from, not where its chunk started and ended.

    parse_markdown merges adjacent short sections, so one chunk can span
    "## Reported accounts + ## Seller CIM + ## Independent QoE". Stamping
    that whole span on every claim drawn from it loses the one thing a
    locator exists to give a reviewer: the place to look. When the model's
    own locator_hint names one of the sections the chunk actually covers,
    resolve to that section alone.

    Falls back to the previous behaviour (chunk locator, hint appended when
    it adds something) whenever the hint names nothing recognisable -- a
    hint like "table row 3" is real information and still worth keeping,
    it just cannot narrow the section by itself.
    """
    prefix = chunk.locator.split("::", 1)[0] if "::" in chunk.locator else chunk.locator
    normalize = lambda s: re.sub(r"[^a-z0-9]", "", (s or "").lower())  # noqa: E731
    if hint:
        needle = normalize(hint)
        for heading in chunk.section_headings:
            trimmed = normalize(heading)
            if trimmed and (trimmed in needle or needle in trimmed):
                return f"{prefix}::{heading}"
    if hint and hint not in chunk.locator and chunk.locator not in hint:
        return f"{chunk.locator}:{hint}"
    return chunk.locator


def _non_empty_l2_text(value: Any) -> str | None:
    """Return trimmed non-empty L2 text without coercing non-string values."""
    if not isinstance(value, str):
        return None
    value = value.strip()
    return value or None


def _normalize_period(raw: str | None) -> str:
    if not raw:
        return "unknown"
    if raw.upper().strip() in PERIOD_MAP:
        return PERIOD_MAP[raw.upper().strip()]
    if re.match(r"^\d{4}-\d{2}-\d{2}$", raw):
        return raw
    return f"RAW:{raw}"


@dataclass
class CanonicalClaim:
    # --- CAP-003 required fields (aligned to benchmark v1.1) ---
    claim_id: str           # ks-sha256 stable content-addressed ID
    statement: str
    source_id: str          # SRC-CIM, SRC-QOE, etc.
    locator: str
    epistemic_class: str    # asserted / observed / derived / attested
    value: float | None
    value_raw: Any
    unit: str | None
    definition_id: str | None
    period: str             # raw period label from source
    period_iso: str         # normalized ISO date or "LTM"
    perimeter: str
    ground_truth_flag: bool  # True only for benchmark validation claims
    validation_only: bool    # True only for quarantined benchmark claims
    notes: str | None
    # --- compiler metadata (NOT in frozen CIC schema) ---
    metric: str             # from METRIC_ENUM — for dedup and reporting
    source_path: str
    known_at: str
    direction: str
    topic: str
    derivation: str | None
    author: str | None
    # Identity dimensions — see tools/object_identity.py for how they key a claim.
    entity: str = "unspecified"
    period_canonical: str = "none"
    scope: str = "unspecified"
    measurement: str = "total"
    claim_kind: str = "QUANTITATIVE"
    bound: str = "EXACT"
    basis: str = "unspecified"
    scenario: str = "unspecified"
    source_version_id: str | None = None
    # Real concept name when metric == "Other". Kept out of the identity on
    # purpose: normalize_metric() maps "Other" to "", so the claim is
    # unresolvable by design — visible, never silently matched.
    metric_label: str | None = None
    validation_errors: list[str] = field(default_factory=list)
    # A deterministic review finding can be material without being grounds to
    # discard evidence. Keep it visible in validation_errors while preserving
    # the separate admission rule used by the graph assembler.
    nonblocking_validation_errors: list[str] = field(default_factory=list)
    # Carried through from RawClaim so derive_relations() can resolve
    # same-batch operand references after claim_id exists. See RawClaim.
    chunk_id: str = ""
    batch_index: int = -1
    derivation_operand_batch_indices: list[int] = field(default_factory=list)
    alternative_to_batch_index: int | None = None
    # Position on this archetype's evidence ladder. See archetype_pack.
    evidence_state: str | None = None


def validate(raw: RawClaim, archetype: str = DEFAULT_ARCHETYPE) -> CanonicalClaim:
    errors: list[str] = []
    nonblocking_errors: list[str] = []
    if raw.metric not in metric_vocabulary(archetype):
        errors.append(f"unknown metric: '{raw.metric}'")
    metric_label = (raw.metric_label or "").strip() or None
    if raw.metric == "Other" and not metric_label:
        # "Other" without a name is the failure it exists to prevent: an
        # unlabelled hole tells a reader nothing, where "Other / Pre-money
        # valuation" tells them exactly which concept the ontology is missing.
        errors.append("metric 'Other' requires metric_label naming the real concept")
    if raw.metric != "Other" and metric_label:
        errors.append(
            f"metric_label '{metric_label}' set on non-Other metric '{raw.metric}'"
        )
    value = _parse_float(raw.value)
    try:
        if not math.isfinite(float(str(raw.value).replace(",", "").replace("%", "").strip())):
            errors.append("non-finite numeric value")
    except (TypeError, ValueError):
        pass  # Textual quantities and ranges are valid; keep their raw value.
    if raw.claim_kind not in CLAIM_KIND_ENUM:
        errors.append(f"invalid claim_kind: '{raw.claim_kind}'")
    if raw.bound not in BOUND_ENUM:
        errors.append(f"invalid bound: '{raw.bound}'")
    period_iso = _normalize_period(raw.period)
    # RAW: prefix means period didn't match PERIOD_MAP — store as-is, no longer reject.
    # period_iso carries the raw label for non-standard periods.
    perimeter = raw.perimeter or "unknown"
    ec = raw.epistemic_class if raw.epistemic_class in EPISTEMIC_CLASS_ENUM else "asserted"
    if raw.epistemic_class not in EPISTEMIC_CLASS_ENUM:
        errors.append(f"invalid epistemic_class: '{raw.epistemic_class}'")
    if ec == "derived" and not (raw.derivation or "").strip():
        errors.append("derived claim missing derivation field")
    elif ec == "derived":
        result = verify_derivation(raw.derivation, raw.value)
        if result["status"] in {
            "computed_disagrees_with_both",
            "value_disagrees_with_text",
        }:
            error = (
                "derivation arithmetic disagrees with stored value "
                f"(status={result['status']}, computed={result.get('computed')}, "
                f"claimed={raw.value})"
            )
            errors.append(error)
            nonblocking_errors.append(error)
    # Completeness against the archetype pack's own declaration (dictionary
    # section 4): a concept states which identity fields a claim of it must
    # carry, so "basis: unspecified" on a concept whose required_identity
    # names basis is incomplete by the pack's rule rather than by a rule
    # anyone wrote into this file. Non-blocking on purpose -- an incomplete
    # claim is still evidence, and 4.2 is explicit that a claim is never
    # forced onto a neighbouring concept to make it look complete.
    incomplete = missing_required_identity(raw)
    if incomplete:
        error = (
            "incomplete identity for this concept: missing "
            + ", ".join(incomplete)
        )
        errors.append(error)
        nonblocking_errors.append(error)
    # The schema tells the model not to emit a CHARACTERISATION, and the model
    # labels one correctly and emits it anyway — observed on "low capital
    # expenditure", classified CHARACTERISATION and returned regardless.
    #
    # So the rule is enforced here instead of asked for. A deterministic filter
    # is the right home for it in any case: it cannot drift between runs, and the
    # claim lands in rejected_claims.json with a reason rather than silently
    # never existing. A seller's adjective kept out of the case is a decision
    # somebody can review; one that was never extracted is invisible.
    if str(raw.claim_kind or "").upper() == "CHARACTERISATION":
        errors.append(
            "characterisation without checkable content — a descriptor, not evidence"
        )
    # Richer than a plain (metric, value, period, perimeter) hash: entity,
    # scope, basis, measurement, scenario, source and epistemic_class all
    # feed the identity too, from tools/object_identity.py's canonical
    # scheme (PAN-63). Two claims that share a metric/value/period/perimeter
    # but differ on any of those dimensions are different claims and must
    # not collide onto the same claim_id.
    claim_id = canonical_claim_id({
        "entity": raw.entity,
        "metric": raw.metric,
        "period": raw.period,
        "period_canonical": raw.period_canonical,
        "scope": raw.scope,
        "basis": raw.basis,
        "measurement": raw.measurement,
        "scenario": raw.scenario,
        "unit": raw.unit,
        "source_id": raw.source_id,
        "source_version_id": raw.source_version_id,
        "locator": raw.locator,
        "epistemic_class": ec,
        # Hash the value exactly as E3 will persist it. Non-numeric bounded or
        # textual quantities (for example "30-60 days") keep ``value_raw``;
        # hashing None here made their emitted IDs impossible to reproduce from
        # e3_claims.json.
        "value": value if value is not None else raw.value,
        "perimeter": perimeter,
    })
    return CanonicalClaim(
        claim_id=claim_id,
        statement=raw.statement,
        source_id=raw.source_id,
        locator=raw.locator,
        epistemic_class=ec,
        value=value,
        value_raw=raw.value,
        unit=raw.unit,
        definition_id=raw.definition_id,
        period=raw.period or "",
        period_iso=period_iso,
        perimeter=perimeter,
        ground_truth_flag=False,
        validation_only=False,
        notes=None,
        metric=raw.metric,
        source_path=raw.source_path,
        known_at=raw.known_at,
        direction=raw.direction,
        topic=raw.topic,
        derivation=raw.derivation,
        author=raw.author,
        source_version_id=raw.source_version_id,
        entity=raw.entity or "unspecified",
        period_canonical=raw.period_canonical or "none",
        scope=raw.scope or "unspecified",
        measurement=raw.measurement or "total",
        claim_kind=raw.claim_kind or "QUANTITATIVE",
        bound=raw.bound or "EXACT",
        basis=raw.basis or "unspecified",
        scenario=raw.scenario or "unspecified",
        metric_label=metric_label,
        validation_errors=errors,
        nonblocking_validation_errors=nonblocking_errors,
        chunk_id=raw.chunk_id,
        batch_index=raw.batch_index,
        derivation_operand_batch_indices=raw.derivation_operand_batch_indices,
        alternative_to_batch_index=raw.alternative_to_batch_index,
        evidence_state=raw.evidence_state,
    )


# ─────────────────────────────────────────────────────────────────────────────
# L4 — Graph assembler (deterministic)
# ─────────────────────────────────────────────────────────────────────────────

@dataclass
class SubGraph:
    claims: list[CanonicalClaim]
    conflicts: list[dict]
    rejected: list[dict]
    admitted_count: int
    rejected_count: int
    conflict_count: int


def assemble(claims: list[CanonicalClaim]) -> SubGraph:
    admitted: dict[str, CanonicalClaim] = {}
    conflicts: list[dict] = []
    rejected: list[dict] = []
    for c in claims:
        blocking_errors = [
            error for error in c.validation_errors
            if error not in c.nonblocking_validation_errors
        ]
        if blocking_errors:
            rejected.append({
                "claim_id": c.claim_id,
                "metric": c.metric,
                "errors": c.validation_errors,
                "statement": c.statement[:80],
            })
            continue
        sid = c.claim_id
        if sid in admitted:
            existing = admitted[sid]
            if existing.value != c.value:
                conflicts.append({
                    "claim_id": sid,
                    "metric": c.metric,
                    "source_a": existing.source_id,
                    "value_a": existing.value,
                    "source_b": c.source_id,
                    "value_b": c.value,
                    "note": "Same stable_id, different values — check perimeter or period.",
                })
        else:
            admitted[sid] = c
    return SubGraph(
        claims=list(admitted.values()),
        conflicts=conflicts,
        rejected=rejected,
        admitted_count=len(admitted),
        rejected_count=len(rejected),
        conflict_count=len(conflicts),
    )


def resolve_operand_claim_ids(claims: list[CanonicalClaim]) -> dict[str, list[str]]:
    """claim_id -> the real claim_ids of its same-batch derivation operands.

    The dictionary's rule for a derived output (section 11.3): it "porta
    input IDs/version/hash" -- a derived value has to name what it was
    computed FROM, or nobody downstream can re-check it. The model can only
    reference operands positionally while it writes them (claim_id does not
    exist yet), so this resolves those positions once, and both the
    DERIVED_FROM edges and each claim's own derivation.operand_claim_ids are
    built from this same answer rather than computing it twice.
    """
    by_batch_position = {
        (c.chunk_id, c.batch_index): c.claim_id
        for c in claims
        if c.chunk_id and c.batch_index >= 0
    }
    resolved: dict[str, list[str]] = {}
    for c in claims:
        ids: list[str] = []
        for operand_index in c.derivation_operand_batch_indices:
            target_id = by_batch_position.get((c.chunk_id, operand_index))
            if target_id and target_id != c.claim_id and target_id not in ids:
                ids.append(target_id)
        if ids:
            resolved[c.claim_id] = ids
    return resolved


def derive_relations(claims: list[CanonicalClaim]) -> list[dict[str, str]]:
    """DERIVED_FROM and ALTERNATIVE_TO edges from same-batch references (see
    RawClaim / CLAIM_TOOL's derivation_operand_indices and
    alternative_to_index).

    Deterministic, not model-trusted: an index has to resolve to a REAL
    admitted claim from the SAME chunk_id, or it is dropped rather than
    turned into an edge pointing at something arbitrary or already rejected.
    Only covers references the model saw in its own tool call -- a
    derivation or scenario link whose counterpart lives in a different
    chunk produces no edge here, same limitation as the cross-chunk-context
    gap parse_markdown's merge already narrows but does not eliminate.
    """
    operands = resolve_operand_claim_ids(claims)
    by_batch_position = {
        (c.chunk_id, c.batch_index): c.claim_id
        for c in claims
        if c.chunk_id and c.batch_index >= 0
    }
    relations: list[dict[str, str]] = []
    for c in claims:
        for target_id in operands.get(c.claim_id, ()):
            relations.append({
                "source_claim_id": c.claim_id,
                "relation": "DERIVED_FROM",
                "target_type": "claim",
                "target_id": target_id,
            })
        if c.alternative_to_batch_index is not None:
            target_id = by_batch_position.get((c.chunk_id, c.alternative_to_batch_index))
            if target_id is not None and target_id != c.claim_id:
                relations.append({
                    "source_claim_id": c.claim_id,
                    "relation": "ALTERNATIVE_TO",
                    "target_type": "claim",
                    "target_id": target_id,
                })
    return relations


# ─────────────────────────────────────────────────────────────────────────────
# Output — E3 manifest format + graph.json for comparison
# ─────────────────────────────────────────────────────────────────────────────

def _to_e3_manifest(graph: SubGraph, deal: str, manifest: str,
                    sources_used: list[dict],
                    workbook_graph_summary: list[dict[str, Any]] | None = None) -> dict:
    """
    Produce an E3 extraction manifest.
    Claims are in CAP-003 format: only the frozen fields, no compiler metadata.
    Compiler metadata lives in the adjacent extraction_metadata section.
    """
    claims_output = []
    for c in graph.claims:
        claims_output.append({
            "claim_id": c.claim_id,
            "statement": c.statement,
            "source_id": c.source_id,
            "source_version_id": c.source_version_id,
            "locator": c.locator,
            "epistemic_class": c.epistemic_class,
            "value": str(c.value) if c.value is not None else c.value_raw,
            "unit": c.unit,
            "definition_id": c.definition_id,
            "period": c.period,
            "perimeter": c.perimeter,
            "ground_truth_flag": c.ground_truth_flag,
            "validation_only": c.validation_only,
            "notes": c.notes,
        })
    return {
        "schema_version": "e3-1.0",
        "manifest_id": manifest,
        "deal": deal,
        "extractor": "extract_v2",
        "sources_included": [s["source_id"] for s in sources_used],
        "claims": claims_output,
        "conflicts": graph.conflicts,
        "extraction_metadata": {
            "admitted_count": graph.admitted_count,
            "rejected_count": graph.rejected_count,
            "conflict_count": graph.conflict_count,
            "rejected": graph.rejected,
            "compiler_fields_per_claim": [
                {
                    "claim_id": c.claim_id,
                    "metric": c.metric,
                    "metric_label": c.metric_label,
                    "value_raw": c.value_raw,
                    "known_at": c.known_at,
                    "direction": c.direction,
                    "topic": c.topic,
                    "derivation": c.derivation,
                    "author": c.author,
                    "entity": c.entity,
                    "period_canonical": c.period_canonical,
                    "scope": c.scope,
                    "measurement": c.measurement,
                    "claim_kind": c.claim_kind,
                    "bound": c.bound,
                    "basis": c.basis,
                    "scenario": c.scenario,
                    # Computed by validate(), consumed by assemble() to decide
                    # admission — but never written to e3_claims.json before
                    # this line, so a flagged-but-admitted claim's flag was
                    # silently lost the moment it left the Python process. A
                    # review flag nobody downstream can see is not a flag.
                    "nonblocking_validation_errors": c.nonblocking_validation_errors,
                }
                for c in graph.claims
            ],
            # Formula cells are captured deterministically at L1 and written
            # beside this manifest.  Keep the E3 lightweight, but make the
            # immutable workbook graph discoverable by the product/compiler.
            "workbook_formula_graphs": workbook_graph_summary or [],
        },
    }


def _capture_workbook_graphs(source_paths: list[Path],
                             source_records: dict[Path, dict] | None = None) -> dict[str, Any]:
    """Capture lossless formula graphs for every Open XML workbook in an intake.

    V2 used to feed formula text to the LLM but then discarded the deterministic
    cell/dependency graph.  That made a formula impossible to inspect or bind
    after an upload.  The sidecar is deliberately separate from E3 claims:
    formulas remain source evidence, never LLM-invented claims.
    """
    from tools.source_graph import capture

    workbooks: list[dict[str, Any]] = []
    for source_path in source_paths:
        if source_path.suffix.lower() not in {".xlsx", ".xlsm"}:
            continue
        graph = capture(source_path).to_json()
        stats = graph.get("stats", {})
        workbooks.append({
            "source_id": (source_records or {}).get(source_path, _source_record(source_path))["source_id"],
            "source_path": str(source_path),
            "source_filename": source_path.name,
            "graph": graph,
            "summary": {
                "formula_count": int(stats.get("by_kind", {}).get("formula", 0)),
                "precedent_edge_count": int(stats.get("precedent_edges", 0)),
                "defined_name_count": int(stats.get("defined_names", 0)),
                "cached_formula_value_count": sum(
                    1 for cell in graph.get("cells", {}).values()
                    if cell.get("kind") == "formula" and cell.get("cached_value") is not None
                ),
                "evaluated_formula_count": sum(
                    1 for cell in graph.get("cells", {}).values()
                    if cell.get("kind") == "formula"
                    and cell.get("evaluation_status") == "CALCULATED_ACYCLIC"
                ),
                "cyclic_formula_count": sum(
                    1 for cell in graph.get("cells", {}).values()
                    if cell.get("kind") == "formula"
                    and cell.get("evaluation_status") == "CYCLIC_COMPONENT"
                ),
            },
        })
    return {"schema": "workbook-formula-graphs-1.0", "workbooks": workbooks}


def _compare_graphs(new_claims: list[dict], old_path: Path) -> dict:
    with open(old_path) as f:
        old_graph = json.load(f)
    old_nodes = {n["id"]: n for n in old_graph.get("nodes", [])
                 if n.get("type") == "claim"}
    new_by_metric_period = {}
    for c in new_claims:
        key = (c.get("metric", c.get("claim_id", "")), c.get("period", ""), c.get("perimeter", ""))
        new_by_metric_period[key] = c
    old_by_key = {}
    for n in old_nodes.values():
        key = (n.get("metric", ""), n.get("as_of", ""), n.get("perimeter", ""))
        old_by_key[key] = n
    only_new = [c for k, c in new_by_metric_period.items() if k not in old_by_key]
    only_old = [n for k, n in old_by_key.items() if k not in new_by_metric_period]
    changes = []
    for k in new_by_metric_period:
        if k in old_by_key:
            nv = str(new_by_metric_period[k].get("value", ""))
            ov = str(old_by_key[k].get("value", ""))
            if nv != ov:
                changes.append({"metric": k[0], "period": k[1],
                                 "old_value": ov, "new_value": nv})
    return {
        "new_total": len(new_claims),
        "old_total": len(old_nodes),
        "only_in_new": len(only_new),
        "only_in_old": len(only_old),
        "value_changes": len(changes),
        "samples_only_new": [
            {"metric": c.get("metric", c.get("claim_id", "")),
             "value": c.get("value"), "statement": c.get("statement", "")[:80]}
            for c in only_new[:5]
        ],
        "samples_only_old": [
            {"metric": n.get("metric", ""), "value": n.get("value"),
             "statement": n.get("statement", "")[:80]}
            for n in only_old[:5]
        ],
        "value_change_samples": changes[:5],
    }


def _sha256_file(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _w(path: Path, obj: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(obj, indent=2, default=str), encoding="utf-8")
    size = path.stat().st_size
    sha = _sha256_file(path)
    print(f"  {path.name:<45} {size//1024:4d}KB  sha256:{sha[:16]}...")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main() -> int:
    ap = argparse.ArgumentParser(description="extract_v2 / E3 pipeline")
    src_group = ap.add_mutually_exclusive_group(required=True)
    src_group.add_argument("--manifest", choices=["K-PRE", "K-IC", "K-LIVE", "ALL"],
                           help="Run over all sources in manifest (no temporal leakage)")
    src_group.add_argument(
        "--source",
        help=(
            "Single source file (.pdf, .docx, .pptx, .xlsx/.xlsm, .csv, "
            ".md/.txt/.html, .srt/.vtt, .eml/.mbox)"
        ),
    )
    ap.add_argument(
        "--input-dir",
        type=Path,
        help=(
            "Directory containing manifest source files. Defaults to vault/inbox; "
            "use this for external or sensitive corpora without copying them into the repo."
        ),
    )
    ap.add_argument("--deal", required=True, help="Deal slug (e.g. keystone)")
    ap.add_argument(
        "--source-envelope", type=Path,
        help="Path to a panta.source-envelope/1.0 JSON record for a --source intake.",
    )
    ap.add_argument(
        "--source-catalog", type=Path,
        help="Per-deal source catalog in CSV or JSON-list format.",
    )
    ap.add_argument("--output", default="pipeline_out/e3",
                    help="Output directory (default: pipeline_out/e3)")
    ap.add_argument("--dry-run", action="store_true",
                    help="Show chunks without calling LLM")
    ap.add_argument("--compare", default=None,
                    help="Path to existing graph.json to compare against")
    ap.add_argument("--workers", type=int, default=3,
                    help="Parallel LLM workers (default: 3)")
    ap.add_argument("--chunk-words", type=int, default=CHUNK_WORDS,
                    help=f"Words per chunk (default: {CHUNK_WORDS})")
    ap.add_argument(
        "--resume-partial-cache",
        action="store_true",
        help=(
            "Migrate a legacy/incomplete raw cache by treating chunks with "
            "existing claims as complete and retrying the rest"
        ),
    )
    ap.add_argument(
        "--via-cli",
        action="store_true",
        help=(
            "Call the model through the Claude Code CLI instead of the API "
            "(subscription-backed). The schema is instructed, not enforced — do "
            "not compare a score from this path against one from the API path"
        ),
    )
    ap.add_argument("--cli-model", default="haiku",
                    help="Model alias for --via-cli (default: haiku)")
    args = ap.parse_args()

    source_catalog = load_source_catalog(args.source_catalog) if args.source_catalog else None
    if args.source_catalog and not source_catalog:
        print(f"WARNING: source catalog is empty or unreadable: {args.source_catalog}",
              file=sys.stderr)
    if source_catalog:
        gaps = sorted(unmapped_source_types(source_catalog))
        if gaps:
            print(f"WARNING: unmapped source_type values: {', '.join(gaps)}", file=sys.stderr)

    # ── Collect source paths ──────────────────────────────────────────────
    if args.manifest:
        source_dir = args.input_dir or VAULT_INBOX
        if not source_dir.is_absolute():
            source_dir = ROOT / source_dir
        source_paths = load_manifest(args.manifest, args.deal, source_dir)
        manifest_label = args.manifest
        if not source_paths:
            print(f"ERROR: No sources found for manifest {args.manifest} in {source_dir}",
                  file=sys.stderr)
            print(f"  Expected stems: {MANIFEST_SOURCES[args.manifest]}")
            return 1
        print(f"\n[Manifest {args.manifest}] {len(source_paths)} sources:")
        for p in source_paths:
            rec = _source_record(p, source_catalog)
            print(f"  {rec['source_id']:<16} {p.name}  (known_at: {rec['known_at']})")
    else:
        source_path = Path(args.source)
        if not source_path.exists():
            source_path = ROOT / args.source
        if not source_path.exists():
            print(f"ERROR: source not found: {args.source}", file=sys.stderr)
            return 1
        source_paths = [source_path]
        manifest_label = "SINGLE"

    source_records: dict[Path, dict] = {}
    if source_catalog:
        for source_path in source_paths:
            record = _source_record(source_path, source_catalog)
            if source_record_from_catalog(source_path, source_catalog) is not None:
                source_records[source_path] = record
    if args.source_envelope:
        if args.manifest:
            print("ERROR: --source-envelope can only be used with --source", file=sys.stderr)
            return 2
        try:
            envelope = json.loads(args.source_envelope.read_text(encoding="utf-8"))
            if envelope.get("schema") != "panta.source-envelope/1.0":
                raise ValueError("unsupported source envelope schema")
            if str(envelope.get("case_id")) != str(args.deal):
                raise ValueError("source envelope case_id does not match --deal")
            source_records[source_paths[0]] = extractor_source_record(envelope)
        except (OSError, ValueError, json.JSONDecodeError) as exc:
            print(f"ERROR: invalid --source-envelope: {exc}", file=sys.stderr)
            return 2

    out_dir = ROOT / args.output / manifest_label
    out_dir.mkdir(parents=True, exist_ok=True)

    # ── L1: Parse all sources ─────────────────────────────────────────────
    print(f"\n[L1] Parsing {len(source_paths)} source(s)...")
    all_chunks: list[Chunk] = []
    for sp in source_paths:
        try:
            chunks = parse_source(sp, args.chunk_words, source_records.get(sp))
        except UnsupportedSourceError as exc:
            print(f"ERROR: {json.dumps(exc.to_dict(), sort_keys=True)}", file=sys.stderr)
            return 2
        print(f"  {sp.name:<50} {len(chunks):3d} chunks")
        all_chunks.extend(chunks)
    print(f"  Total: {len(all_chunks)} chunks  "
          f"(avg {sum(c.word_count for c in all_chunks)//max(len(all_chunks),1)} w/chunk)")

    # Persist the lossless workbook structure before L2.  This is essential
    # even for a dry run or a later LLM failure: the formula text, precedents
    # and cached values are facts contained in the uploaded source.
    workbook_graphs = _capture_workbook_graphs(source_paths, source_records)
    workbook_graph_path = out_dir / "workbook_formula_graphs.json"
    _w(workbook_graph_path, workbook_graphs)
    for workbook in workbook_graphs["workbooks"]:
        summary = workbook["summary"]
        print(
            f"  {workbook['source_filename']:<50} "
            f"{summary['formula_count']:5d} formulas  "
            f"{summary['precedent_edge_count']:5d} precedent edges"
        )

    chunks_debug = [
        {"chunk_id": c.chunk_id, "locator": c.locator,
         "source_id": c.source_record["source_id"],
         "provenance": c.provenance,
         "period_context": c.period_context,
         "section_heading": c.section_heading,
         "page_or_slide_number": c.page_or_slide_number,
         "word_count": c.word_count, "preview": c.body[:100].replace("\n", " ") + "..."}
        for c in all_chunks
    ]
    _w(out_dir / "chunks_debug.json", chunks_debug)

    if args.dry_run:
        print(f"\n[DRY-RUN] Chunks written to {out_dir}/chunks_debug.json")
        _print_chunk_summary(all_chunks)
        return 0

    # ── L2: Annotate ──────────────────────────────────────────────────────
    # Raw claims cache: completed modern runs are reused. Interrupted/failed modern
    # runs retain their successful chunks and retry only the outstanding chunk IDs.
    # A cache without the companion status file is treated as a legacy complete cache.
    raw_cache_path = out_dir / "raw_claims_cache.json"
    chunk_status_path = out_dir / "l2_chunk_status.json"
    failed_chunks_path = out_dir / "failed_chunks.json"
    cached_raw: list[RawClaim] = []
    completed_chunk_ids: set[str] = set()
    completed_chunk_models: dict[str, str] = {}
    pending_chunks = list(all_chunks)
    modern_cache = raw_cache_path.exists() and chunk_status_path.exists()

    if raw_cache_path.exists() and not modern_cache and not args.resume_partial_cache:
        print(f"\n[L2] Cache found — loading raw claims from {raw_cache_path.name} (skipping API calls)")
        cached = json.loads(raw_cache_path.read_text())
        all_raw = [RawClaim(**c) for c in cached]
        print(f"  Loaded {len(all_raw)} raw claims from cache")
    else:
        if raw_cache_path.exists() and args.resume_partial_cache and not modern_cache:
            cached = json.loads(raw_cache_path.read_text())
            cached_raw = [RawClaim(**c) for c in cached]
            cached_locators = {claim.locator for claim in cached_raw}
            completed_chunk_ids = {
                chunk.chunk_id
                for chunk in all_chunks
                if any(
                    locator == chunk.locator
                    or locator.startswith(chunk.locator + ":")
                    for locator in cached_locators
                )
            }
            pending_chunks = [
                chunk for chunk in all_chunks
                if chunk.chunk_id not in completed_chunk_ids
            ]
            print(
                f"\n[L2] Migrating partial legacy cache — "
                f"{len(completed_chunk_ids)}/{len(all_chunks)} chunks have claims, "
                f"{len(pending_chunks)} to retry"
            )
        elif modern_cache:
            cached = json.loads(raw_cache_path.read_text())
            cached_raw = [RawClaim(**c) for c in cached]
            status = json.loads(chunk_status_path.read_text())
            completed_chunk_ids = set(status.get("completed_chunk_ids", []))
            completed_chunk_models = {
                str(chunk_id): str(model)
                for chunk_id, model in status.get("completed_chunk_models", {}).items()
                if chunk_id in completed_chunk_ids
            }
            if len(completed_chunk_models) < len(completed_chunk_ids):
                existing_model = os.environ.get(
                    "PEOS_EXISTING_CACHE_MODEL",
                    "legacy-cache",
                ).strip() or "legacy-cache"
                for chunk_id in completed_chunk_ids:
                    completed_chunk_models.setdefault(chunk_id, existing_model)
            pending_chunks = [
                chunk for chunk in all_chunks
                if chunk.chunk_id not in completed_chunk_ids
            ]
            if not pending_chunks and status.get("complete"):
                all_raw = cached_raw
                print(
                    f"\n[L2] Complete chunk cache found — loaded "
                    f"{len(all_raw)} raw claims (skipping API calls)"
                )
            else:
                print(
                    f"\n[L2] Resuming partial cache — "
                    f"{len(completed_chunk_ids)}/{len(all_chunks)} chunks complete, "
                    f"{len(pending_chunks)} to retry"
                )

        if modern_cache and not pending_chunks and status.get("complete"):
            pass
        else:
            if args.via_cli:
                # Subscription-backed path for re-runs. The schema is instructed
                # rather than enforced here (the CLI has no tool_choice), so a
                # number produced this way must not be compared against one
                # produced via the API — see tools/llm_cli_provider.py.
                from tools.llm_cli_provider import CliClient
                client = CliClient(model=args.cli_model)
                print(f"\n[L2] via Claude CLI (model={args.cli_model}) — "
                      f"schema instructed, not enforced")
            else:
                api_key = configured_api_key()
                if not api_key:
                    print(f"ERROR: {missing_key_message()}.", file=sys.stderr)
                    print("  Or use --dry-run to inspect chunks.")
                    return 1
                try:
                    import anthropic
                    client = anthropic.Anthropic(**anthropic_client_kwargs(api_key))
                except ImportError:
                    print("ERROR: pip install anthropic", file=sys.stderr)
                    return 1
            print(f"\n[L2] Annotating {len(pending_chunks)} chunk(s) "
                  f"(workers={args.workers}, model={MODEL})...")
            all_raw = list(cached_raw)
            processed = 0
            batch_errors: list[dict[str, str]] = []

            def _checkpoint(*, complete: bool = False) -> None:
                from dataclasses import asdict

                raw_cache_path.write_text(
                    json.dumps([asdict(r) for r in all_raw], indent=2, default=str),
                    encoding="utf-8",
                )
                chunk_status_path.write_text(
                    json.dumps(
                        {
                            "schema_version": "l2-chunk-status-1.0",
                            "total_chunks": len(all_chunks),
                            "completed_chunk_ids": sorted(completed_chunk_ids),
                            "completed_chunk_models": dict(
                                sorted(completed_chunk_models.items())
                            ),
                            "models_used": sorted(set(completed_chunk_models.values())),
                            "failed_chunks": batch_errors,
                            "complete": complete,
                        },
                        indent=2,
                    ),
                    encoding="utf-8",
                )
                failed_chunks_path.write_text(
                    json.dumps(batch_errors, indent=2),
                    encoding="utf-8",
                )

            def _process(chunk: Chunk) -> tuple[Chunk, list[RawClaim]]:
                for attempt in range(MAX_PROVIDER_RETRIES + 1):
                    try:
                        return chunk, annotate_chunk(
                            chunk,
                            client,
                            args.deal,
                            raise_errors=True,
                        )
                    except Exception as exc:
                        delay = _provider_retry_delay(exc, attempt)
                        if delay is None or attempt >= MAX_PROVIDER_RETRIES:
                            raise
                        print(
                            f"  [L2 RETRY] {chunk.chunk_id}: "
                            f"{type(exc).__name__}; retry "
                            f"{attempt + 1}/{MAX_PROVIDER_RETRIES} in {delay:g}s",
                            file=sys.stderr,
                        )
                        time.sleep(delay)
                raise AssertionError("unreachable provider retry state")

            with ThreadPoolExecutor(max_workers=args.workers) as pool:
                futures = {pool.submit(_process, c): c for c in pending_chunks}
                for fut in as_completed(futures):
                    chunk = futures[fut]
                    processed += 1
                    try:
                        _, raw_claims = fut.result()
                    except Exception as exc:
                        message = str(exc).splitlines()[0][:1000]
                        batch_errors.append(
                            {
                                "chunk_id": chunk.chunk_id,
                                "locator": chunk.locator,
                                "error_type": type(exc).__name__,
                                "message": message,
                            }
                        )
                        print(
                            f"  [L2 ERROR] {chunk.chunk_id}: "
                            f"{type(exc).__name__}: {message}",
                            file=sys.stderr,
                        )
                        if _is_fatal_provider_error(exc):
                            for queued in futures:
                                if queued is not fut:
                                    queued.cancel()
                            print(
                                "  [L2] Fatal provider/key error; remaining "
                                "queued chunks were cancelled.",
                                file=sys.stderr,
                            )
                            break
                    else:
                        all_raw.extend(raw_claims)
                        completed_chunk_ids.add(chunk.chunk_id)
                        completed_chunk_models[chunk.chunk_id] = MODEL
                        if raw_claims:
                            print(f"  [{processed:03d}/{len(pending_chunks):03d}] "
                                  f"{chunk.locator[:55]:<55} → {len(raw_claims)} claim(s)")
                    if processed % 10 == 0:
                        _checkpoint()

            is_complete = (
                not batch_errors
                and len(completed_chunk_ids) == len(all_chunks)
            )
            _checkpoint(complete=is_complete)
            print(f"  Raw claims checkpointed → {raw_cache_path.name}")
            if not is_complete:
                outstanding = len(all_chunks) - len(completed_chunk_ids)
                print(
                    f"ERROR: L2 incomplete — {outstanding} chunk(s) outstanding; "
                    "rerun the same command to retry only outstanding chunks.",
                    file=sys.stderr,
                )
                return 3

    print(f"  Total raw: {len(all_raw)}")

    # ── L3: Validate ──────────────────────────────────────────────────────
    print("\n[L3] Validating...")
    canonicals = [validate(r) for r in all_raw]
    invalid = [
        c for c in canonicals
        if any(error not in c.nonblocking_validation_errors
               for error in c.validation_errors)
    ]
    if invalid:
        print(f"  Rejected: {len(invalid)}")
        for c in invalid[:5]:
            print(f"    x {c.metric}: {c.validation_errors}")

    # ── L4: Assemble ──────────────────────────────────────────────────────
    print("\n[L4] Assembling...")
    graph = assemble(canonicals)
    print(f"  Admitted: {graph.admitted_count}  "
          f"Rejected: {graph.rejected_count}  "
          f"Conflicts: {graph.conflict_count}")
    if graph.conflicts:
        for conf in graph.conflicts:
            print(f"    conflict: {conf['metric']} {conf['value_a']} vs {conf['value_b']}")

    # ── Output ────────────────────────────────────────────────────────────
    print(f"\n[Output] Writing to {out_dir}/...")
    sources_used = [source_records.get(p, _source_record(p)) for p in source_paths]
    workbook_graph_summary = [
        {
            "source_id": item["source_id"],
            "source_filename": item["source_filename"],
            **item["summary"],
            "artifact": "workbook_formula_graphs.json",
        }
        for item in workbook_graphs["workbooks"]
    ]
    e3 = _to_e3_manifest(
        graph,
        args.deal,
        manifest_label,
        sources_used,
        workbook_graph_summary,
    )
    if chunk_status_path.exists():
        l2_status = json.loads(chunk_status_path.read_text())
        e3["extraction_metadata"]["l2_complete"] = bool(
            l2_status.get("complete")
        )
        e3["extraction_metadata"]["llm_models"] = list(
            l2_status.get("models_used", [])
        )
    _w(out_dir / "e3_claims.json", e3)
    _w(out_dir / "rejected_claims.json", graph.rejected)
    _w(out_dir / "conflict_report.json", graph.conflicts)

    # Metric distribution
    metric_counts: dict[str, int] = {}
    for c in graph.claims:
        metric_counts[c.metric] = metric_counts.get(c.metric, 0) + 1

    report_lines = [
        f"E3 Extraction Report — {manifest_label}",
        "=" * 50,
        f"  Deal        : {args.deal}",
        f"  Manifest    : {manifest_label}",
        f"  Sources     : {len(source_paths)}",
        f"  Chunks      : {len(all_chunks)}",
        f"  Raw claims  : {len(all_raw)}",
        f"  Admitted    : {graph.admitted_count}",
        f"  Rejected    : {graph.rejected_count}",
        f"  Conflicts   : {graph.conflict_count}",
        "",
        "Admitted by metric:",
    ]
    for metric, count in sorted(metric_counts.items(), key=lambda x: -x[1]):
        report_lines.append(f"  {count:3d}  {metric}")
    if graph.conflicts:
        report_lines += ["", "Conflicts:"]
        for conf in graph.conflicts:
            report_lines.append(
                f"  {conf['metric']}: {conf['value_a']} ({conf['source_a']}) "
                f"vs {conf['value_b']} ({conf['source_b']})"
            )
    report_text = "\n".join(report_lines) + "\n"
    (out_dir / "extraction_report.txt").write_text(report_text, encoding="utf-8")

    # Hash manifest for version pinning
    e3_path = out_dir / "e3_claims.json"
    manifest_record = {
        "manifest_id": manifest_label,
        "deal": args.deal,
        "sources": [s["source_id"] for s in sources_used],
        "admitted_count": graph.admitted_count,
        "e3_claims_sha256": _sha256_file(e3_path),
    }
    _w(out_dir / "manifest_hash.json", manifest_record)

    print()
    print(report_text)

    # ── Compare ───────────────────────────────────────────────────────────
    if args.compare:
        compare_path = Path(args.compare)
        if not compare_path.exists():
            compare_path = ROOT / args.compare
        if compare_path.exists():
            print(f"[Compare] vs {compare_path.name}...")
            diff = _compare_graphs(e3["claims"], compare_path)
            _w(out_dir / "comparison.json", diff)
            print(f"  New: {diff['new_total']}  Old: {diff['old_total']}  "
                  f"Only-new: {diff['only_in_new']}  Only-old: {diff['only_in_old']}  "
                  f"Changed: {diff['value_changes']}")

    return 0


def _print_chunk_summary(chunks: list[Chunk]) -> None:
    print(f"\n  Chunk summary ({len(chunks)} total):")
    for c in chunks[:12]:
        preview = c.body[:55].replace("\n", " ")
        print(f"    {c.source_record['source_id']:<14} {c.locator:<50} "
              f"{c.word_count:3d}w  \"{preview}...\"")
    if len(chunks) > 12:
        print(f"    ... and {len(chunks) - 12} more")


if __name__ == "__main__":
    sys.exit(main())
