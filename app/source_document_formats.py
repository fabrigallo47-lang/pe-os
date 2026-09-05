"""Focused readers for native fixture-supported slides, mail, Word sections and images."""
import base64
import html
import io
import re
from email import policy
from email.parser import BytesParser

from fastapi import HTTPException


def rectangle(locator, width, height):
    match = re.search(r':rect:([\d.,-]+)$', locator)
    if not match:
        if ':rect:' in locator:
            raise HTTPException(422, 'The cited image region is malformed.')
        return None
    try:
        box = [float(n) for n in match[1].split(',')]
        valid = len(box) == 4 and 0 <= box[0] < box[2] <= width and 0 <= box[1] < box[3] <= height
    except ValueError:
        valid = False
    if not valid:
        raise HTTPException(422, 'The cited image region is outside this source version.')
    return box


def word_sections(document):
    from docx import Document
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    word = Document(io.BytesIO(document.data))
    blocks = []
    for child in word.element.body:
        if child.tag.endswith('}p'):
            paragraph = Paragraph(child, word)
            pictures = []
            for blip in child.xpath('.//a:blip'):
                relation = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                part = word.part.related_parts.get(relation)
                if part is not None and part.content_type in {'image/png', 'image/jpeg', 'image/gif'}:
                    pictures.append((part.content_type, part.blob))
            blocks.append({'text': paragraph.text, 'heading': paragraph.style.name.startswith('Heading'), 'pictures': pictures})
        elif child.tag.endswith('}tbl'):
            table = Table(child, word)
            blocks.append({'text': '\n'.join(' | '.join(cell.text for cell in row.cells) for row in table.rows), 'heading': False, 'pictures': []})
    return blocks


def extra_position(document, locator):
    address = locator.split('::', 1)[-1]
    if document.suffix == '.pptx':
        from pptx import Presentation
        deck = Presentation(io.BytesIO(document.data))
        match = re.fullmatch(r'slide:([1-9]\d*)(?::w\d+-\d+)?', address)
        slide = int(match[1]) if match else 1
        if slide > len(deck.slides):
            raise HTTPException(422, 'The cited slide does not exist in this source version.')
        return {'kind': 'text', 'native': 'pptx', 'status': 'LOCATED' if match else 'UNRESOLVED', 'slide': slide,
                'label': f'Slide {slide} · native text, tables and chart data' if match else 'Exact slide not supplied'}
    if document.suffix == '.eml':
        message = BytesParser(policy=policy.default).parsebytes(document.data)
        valid = address in {'message:1:body', 'message:1:headers'} or address in {
            f'message-id:{message.get("Message-ID")}:body', f'message-id:{message.get("Message-ID")}:headers'}
        return {'kind': 'text', 'native': 'email', 'status': 'LOCATED' if valid else 'UNRESOLVED',
                'part': 'headers' if address.endswith(':headers') else 'body',
                'label': 'Email ' + ('headers' if address.endswith(':headers') else 'body') if valid else 'Exact message part not supplied'}
    if document.suffix == '.docx':
        blocks = word_sections(document)
        title = address.removeprefix('section:')
        matches = [i for i, block in enumerate(blocks) if block['heading'] and block['text'] == title]
        start = matches[0] if len(matches) == 1 else None
        end = next((i for i in range(start + 1, len(blocks)) if blocks[i]['heading']), len(blocks)) if start is not None else None
        return {'kind': 'text', 'native': 'word-section', 'status': 'LOCATED' if start is not None else 'UNRESOLVED',
                'start': start, 'end': end, 'label': f'Section: {title}' if start is not None else 'Exact Word section could not be located'}
    from PIL import Image
    with Image.open(io.BytesIO(document.data)) as picture:
        width, height = picture.size
    box = rectangle(address, width, height)
    if address.split(':rect:', 1)[0] != 'image:1':
        box = None
    return {'kind': 'image', 'native': 'image', 'status': 'LOCATED' if box or address == 'image:1' else 'UNRESOLVED',
            'box': box, 'width': width, 'height': height, 'label': 'Cited image region' if box else 'Original image'}


def extra_body(document, position):
    native = position['native']
    if native == 'pptx':
        from pptx import Presentation
        from tools.extract_v2_physical import _pptx_shape_text, _pptx_reading_order
        deck = Presentation(io.BytesIO(document.data))
        slide = deck.slides[position['slide'] - 1]
        parts = [text for shape in _pptx_reading_order(slide.shapes, deck.slide_height) for text in _pptx_shape_text(shape)]
        return '<p>Original slide content. Text, tables and chart values retain their native source; this view does not reproduce the slide layout.</p><div id="selection">' + ''.join('<pre>' + html.escape(part) + '</pre>' for part in parts) + '</div>'
    if native == 'email':
        from tools.extract_v2_physical import _email_body, _email_attachments
        message = BytesParser(policy=policy.default).parsebytes(document.data)
        headers = '\n'.join(f'{key}: {message.get(key, "")}' for key in ('Subject', 'From', 'To', 'Date', 'Message-ID'))
        body = _email_body(message)
        return ('<section' + (' id="selection"' if position['part'] == 'headers' else '') + '><pre>' + html.escape(headers)
                + '</pre></section><section' + (' id="selection"' if position['part'] == 'body' else '')
                + '><h2>Message body</h2><pre>' + html.escape(body) + '</pre></section><h2>Attachments in the original email</h2><ul>'
                + ''.join('<li>' + html.escape(name) + '</li>' for name in _email_attachments(message)) + '</ul>')
    if native == 'word-section':
        blocks = word_sections(document)
        rendered = []
        for i, block in enumerate(blocks):
            selected = position['start'] is not None and position['start'] <= i < position['end']
            images = ''.join(f'<img alt="Original embedded image" class="pdf" src="data:{mime};base64,{base64.b64encode(blob).decode()}">' for mime, blob in block['pictures'])
            rendered.append('<div' + (' class="selected"' if selected else '') + (' id="selection"' if i == position['start'] else '')
                            + '><pre>' + html.escape(block['text']) + '</pre>' + images + '</div>')
        return ''.join(rendered)
    encoded = base64.b64encode(document.data).decode()
    box = position['box']
    highlight = ''
    if box:
        left, top, right, bottom = box
        w, h = position['width'], position['height']
        highlight = f'<span id="selection" style="position:absolute;left:{left/w*100}%;top:{top/h*100}%;width:{(right-left)/w*100}%;height:{(bottom-top)/h*100}%;outline:3px solid #d7b83d;pointer-events:none"></span>'
    return f'<div style="position:relative">{highlight}<img class="pdf" alt="Original image" src="data:{document.media_type};base64,{encoded}"></div>'
